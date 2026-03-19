#!/usr/bin/env python3
"""Design Thinking Workshop App — Backend Server

Serves the DTW collaborative workshop app with:
  GET    /                              — serve index.html
  GET    /style.css, /app.js            — static assets
  GET    /exercises.json                — exercise catalog
  POST   /api/sessions                  — create session
  GET    /api/sessions                  — list all sessions
  GET    /api/sessions/<id>             — get session state
  POST   /api/sessions/<id>/advance     — next exercise
  POST   /api/sessions/<id>/back        — previous exercise
  POST   /api/sessions/<id>/summary     — generate AI summary
  GET    /api/sessions/<id>/export      — download session JSON
  GET    /api/sessions/<id>/export-chat — download full chat log as .docx

SocketIO events:
  Client→Server:   join_session, send_message, facilitator_join
  Server→Client:   session_state, new_message, exercise_changed,
                   summary_ready, participant_joined, error
"""
import io
import json
import logging
import os
import ssl
import subprocess
import sys
import tempfile
import urllib.request
import uuid
from datetime import datetime, timezone
from pathlib import Path

# ── Load .env (no third-party library needed) ──────────────────────────────────
_env_path = os.path.join(os.path.dirname(__file__), ".env")
if os.path.exists(_env_path):
    with open(_env_path) as _f:
        for _line in _f:
            _line = _line.strip()
            if _line and not _line.startswith("#") and "=" in _line:
                _k, _v = _line.split("=", 1)
                os.environ.setdefault(_k.strip(), _v.strip())

import eventlet
import eventlet.tpool
eventlet.monkey_patch()

from flask import Flask, jsonify, request, send_from_directory, Response
from flask_socketio import SocketIO, emit, join_room

# ── Config ─────────────────────────────────────────────────────────────────────
PORT       = int(os.environ.get("PORT", 5001))
BASE_DIR   = Path(__file__).parent
SESSIONS_DIR = BASE_DIR / "sessions"
SESSIONS_DIR.mkdir(exist_ok=True)
ARTIFACTS_DIR = BASE_DIR / "artifacts"
ARTIFACTS_DIR.mkdir(exist_ok=True)

# ── S3 storage (optional — falls back to local filesystem when not configured) ──
S3_BUCKET  = os.environ.get("S3_BUCKET", "")           # e.g. "lilly-dtw-sessions"
S3_REGION  = os.environ.get("S3_REGION", "us-east-1")
S3_PREFIX  = os.environ.get("S3_PREFIX", "dtw")        # key prefix inside bucket
# Credentials are picked up automatically from:
#   AWS_ACCESS_KEY_ID + AWS_SECRET_ACCESS_KEY env vars, ~/.aws/credentials, or IAM role
_S3_CLIENT = None  # initialised lazily on first use


def _s3():
    """Return a boto3 S3 client, creating it on first call."""
    global _S3_CLIENT
    if _S3_CLIENT is None:
        import boto3
        _S3_CLIENT = boto3.client("s3", region_name=S3_REGION)
    return _S3_CLIENT


def _s3_enabled() -> bool:
    return bool(S3_BUCKET)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    datefmt="%Y-%m-%dT%H:%M:%SZ",
)
log = logging.getLogger("dtw")

app = Flask(__name__, static_folder=str(BASE_DIR))
app.config["SECRET_KEY"] = os.environ.get("FLASK_SECRET_KEY", os.urandom(24).hex())
socketio = SocketIO(app, cors_allowed_origins="*", async_mode="eventlet")

# In-memory map: room_code → session_id  (rebuilt from disk on startup)
_room_code_map: dict[str, str] = {}


# ── Utilities ──────────────────────────────────────────────────────────────────

def utcnow() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def room_code_from_id(session_id: str) -> str:
    return session_id.replace("-", "")[:6].upper()


def session_path(session_id: str) -> Path:
    return SESSIONS_DIR / f"session_{session_id}.json"


def _s3_session_key(session_id: str) -> str:
    return f"{S3_PREFIX}/sessions/session_{session_id}.json"


def load_session(session_id: str) -> dict | None:
    if _s3_enabled():
        try:
            obj = _s3().get_object(Bucket=S3_BUCKET, Key=_s3_session_key(session_id))
            return json.loads(obj["Body"].read().decode("utf-8"))
        except Exception:
            return None
    p = session_path(session_id)
    if not p.exists():
        return None
    with open(p, encoding="utf-8") as f:
        return json.load(f)


def save_session(session: dict) -> None:
    body = json.dumps(session, indent=2, ensure_ascii=False).encode("utf-8")
    if _s3_enabled():
        _s3().put_object(
            Bucket=S3_BUCKET,
            Key=_s3_session_key(session["session_id"]),
            Body=body,
            ContentType="application/json",
        )
        return
    p = session_path(session["session_id"])
    tmp = p.with_suffix(".tmp")
    tmp.write_bytes(body)
    tmp.replace(p)


def _rebuild_room_code_map() -> None:
    """Rebuild in-memory room_code→session_id map from storage at startup."""
    if _s3_enabled():
        try:
            paginator = _s3().get_paginator("list_objects_v2")
            prefix = f"{S3_PREFIX}/sessions/"
            for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=prefix):
                for obj in page.get("Contents", []):
                    try:
                        raw = _s3().get_object(Bucket=S3_BUCKET, Key=obj["Key"])
                        s = json.loads(raw["Body"].read().decode("utf-8"))
                        rc = s.get("room_code") or room_code_from_id(s["session_id"])
                        _room_code_map[rc] = s["session_id"]
                    except Exception:
                        pass
        except Exception as exc:
            log.warning("Could not rebuild room code map from S3: %s", exc)
        return
    for p in SESSIONS_DIR.glob("session_*.json"):
        try:
            with open(p, encoding="utf-8") as f:
                s = json.load(f)
            rc = s.get("room_code") or room_code_from_id(s["session_id"])
            _room_code_map[rc] = s["session_id"]
        except Exception:
            pass


_rebuild_room_code_map()


def load_exercises() -> list[dict]:
    p = BASE_DIR / "exercises.json"
    with open(p, encoding="utf-8") as f:
        return json.load(f)


def get_exercise_meta(exercise_id: str) -> dict | None:
    for ex in load_exercises():
        if ex["id"] == exercise_id:
            return ex
    return None


# ── AI Summary (Lilly LLM Gateway) ────────────────────────────────────────────

LLM_GATEWAY_URL = os.environ.get(
    "LLM_GATEWAY_URL",
    "https://lilly-code-server.api.gateway.llm.lilly.com"
).rstrip("/")
LLM_MODEL = os.environ.get("LLM_MODEL", "claude-sonnet-4-6")

# ── Facilitator access control ────────────────────────────────────────────────

# Names exactly as they appear in the JWT 'name' claim (case-insensitive match)
FACILITATORS: list[str] = [
    "Matt Wambsganss",
]


_keychain: dict = {}
_keychain_lock = eventlet.semaphore.Semaphore(1)


def _read_keychain_token() -> str:
    """Read the lilly-code JWT from the macOS keychain (populated by Lilly Code VS Code extension).
    The keychain value is a JSON blob: {"access_token": "eyJ...", "expires_at": "..."}.
    Uses eventlet.tpool so the blocking subprocess runs in a real OS thread (not a green thread),
    and a semaphore to prevent concurrent callers from each triggering a separate keychain prompt.
    After reading, re-applies the partition-list ACL so future reads don't prompt again.
    """
    # Fast path — return cached value if available
    if _keychain.get("token"):
        return _keychain["token"]

    with _keychain_lock:
        # Re-check inside lock in case another caller just populated it
        if _keychain.get("token"):
            return _keychain["token"]

        def _run_security():
            return subprocess.run(
                ["security", "find-generic-password", "-s", "lilly-code", "-w"],
                capture_output=True, text=True,
            )

        try:
            result = eventlet.tpool.execute(_run_security)
            if result.returncode != 0:
                raise subprocess.CalledProcessError(result.returncode, "security")
            raw = result.stdout.strip()
            try:
                data = json.loads(raw)
                token = data.get("access_token", "")
            except json.JSONDecodeError:
                token = raw  # fallback: treat raw value as the token
            _keychain["token"] = token
            # Grant python3/security permanent ACL access so future reads don't prompt.
            # Runs silently — Lilly Code extension resets the ACL on each token refresh,
            # so we re-apply it here each time we read.
            eventlet.tpool.execute(lambda: subprocess.run(
                ["security", "set-generic-password-partition-list",
                 "-s", "lilly-code", "-a", "sso-token",
                 "-S", "apple-tool:,apple:", "-k", ""],
                capture_output=True, text=True,
            ))
            return token
        except subprocess.CalledProcessError:
            log.warning("lilly-code keychain entry not found — AI summaries will fail until you sign in via VS Code")
            return ""


def _decode_jwt_payload(token: str) -> dict:
    """Decode the payload of a JWT without signature verification."""
    import base64
    try:
        parts = token.split(".")
        if len(parts) < 2:
            return {}
        payload = parts[1]
        # Add padding
        payload += "=" * (-len(payload) % 4)
        return json.loads(base64.urlsafe_b64decode(payload).decode("utf-8"))
    except Exception:
        return {}


def _identity_from_keychain() -> dict:
    """Return {name, email, is_facilitator} from the keychain JWT, or empty strings."""
    token = _keychain.get("token", "") or _read_keychain_token()
    if not token:
        return {"name": "", "email": "", "is_facilitator": False}
    claims = _decode_jwt_payload(token)
    name  = claims.get("name") or claims.get("given_name", "")
    email = claims.get("upn") or claims.get("email") or claims.get("preferred_username", "")
    is_fac = name.lower() in [f.lower() for f in FACILITATORS]
    return {"name": name, "email": email, "is_facilitator": is_fac}


def _get_llm_token() -> str:
    """Return a valid gateway token. Use the cached value to avoid repeated keychain prompts;
    the cache is cleared on 401/403 so a VS Code re-login is still picked up automatically."""
    cached = _keychain.get("token", "")
    if cached:
        return cached
    return _read_keychain_token()

def _ssl_ctx() -> ssl.SSLContext:
    ctx = ssl.create_default_context()
    ctx.check_hostname = False
    ctx.verify_mode = ssl.CERT_NONE
    return ctx


def generate_exercise_summary(exercise_meta: dict, messages: list[dict]) -> str:
    token = _get_llm_token()
    if not token:
        raise ValueError(
            "No LLM gateway token found. Sign in via the Lilly Code VS Code extension to populate the keychain."
        )

    transcript_lines = []
    for m in messages:
        ts = m.get("timestamp", "")[:16].replace("T", " ")
        transcript_lines.append(f"[{ts}] {m.get('sender','?')}: {m.get('text','')}")
    transcript = "\n".join(transcript_lines) if transcript_lines else "(no messages)"

    participant_names = list({m.get("sender", "?") for m in messages})
    participant_count = len(participant_names)

    system_prompt = (
        "You are an expert Design Thinking Workshop facilitator trained in IBM Enterprise Design Thinking methodology. "
        "Your role is to synthesize live workshop discussions into structured, executive-ready outputs. "
        "Outputs should be actionable, specific, and grounded entirely in what participants actually said — never invent content. "
        "Format responses in clean Markdown. Use ## for section headers and - for bullets. "
        "Each section header becomes a slide in the deliverable deck, so keep headers short and descriptive. "
        "Write as if preparing a professional workshop output that will be shared with senior leadership."
    )

    # Exercise-specific output structure mapped to IBM field guide methodology
    output_formats = {
        "intros_outcomes": """## Output Structure
## Participants in the Room
-- (Name and role/team for each participant who introduced themselves)

## Desired Outcomes
-- (Summarized outcomes — synthesize participant input into clear, concise statements; do not quote verbatim)

## Common Themes
-- (2–3 themes that emerged across multiple desired outcomes)

## Shared Definition of Success
-- (1–2 sentences capturing what the group collectively wants to achieve today)""",

        "hopes_fears": """## Output Structure
## Hopes
- (Each hope, close to participant's own words — organized by theme)

## Fears
- (Each fear or concern, close to participant's own words — organized by theme)

## Key Tensions
- (Hopes and fears that are in direct conflict with each other)

## Implications for the Workshop
- (How should the facilitator use these to shape the day)""",

        "persona_cards": """## Output Structure
For each distinct persona mentioned, create a section:

## Persona: [Role Title]
- **Goals:** what they are trying to accomplish
- **Pain Points:** top 2–3 frustrations
- **Representative Quote:** one quote that captures their experience

Repeat for each persona. End with:

## Common Threads Across Personas
- (Shared goals or frustrations that span multiple roles)""",

        "stakeholder_map": """## Output Structure
## Core Process Owners
- (Roles who own steps in the process — executors)

## Decision Makers
- (Roles who approve, gate, or direct the work)

## Impacted Parties
- (Roles who feel the downstream effects but don't control the process)

## Enablers & Support Roles
- (Roles that provide tools, data, or support but aren't in the critical path)

## Key Friction Points
- (Where stakeholder relationships are unclear or cause conflict)""",

        "empathy_maps": """## Output Structure
Create ONE section per persona. Use EXACTLY this format for each:

## [Persona Name / Role]
- **CONTEXT:** [1–2 sentences: who they are, their role in this process, and the single defining challenge they face]
- **SAYS:** "[A direct quote or paraphrase of something this persona says aloud]"
- **SAYS:** "[Another representative quote]"
- **DOES:** [A concrete action or behavior they perform in the process]
- **DOES:** [Another action or behavior]
- **THINKS:** [An internal assumption, mental model, or belief they hold]
- **THINKS:** [Another thought or concern they rarely voice]
- **FEELS:** [An emotional state — frustration, anxiety, pride, confusion]
- **FEELS:** [Another emotion or tension they experience]
- **PAINS:** [Top frustration or pain point — specific and grounded in the transcript]
- **PAINS:** [Second key pain point]
- **GAINS:** [What success looks like for this persona — concrete and measurable]
- **GAINS:** [Another gain or goal they are working toward]

Repeat the ## [Persona Name / Role] block for every persona identified in the conversation.
End with:

## Common Threads Across Personas
- (Shared pain points, emotions, or beliefs that span multiple personas)""",

        "asis_process_map": """## Output Structure
## Process Overview
- (Brief description of the process being mapped and key roles involved)

## Process Steps
- **Step 1 — [Step Name]:** Owner: [Role]. What happens. Pain points if any.
- (Continue for each step in sequence)

## Key Handoff Points
- (Where work moves between teams — flag delays or confusion here)

## Pain Points by Stage
- (Grouped frustrations: where does the most friction occur and why)

## Root Cause Patterns
- (Structural issues: parallel execution, ownership fragmentation, communication gaps)""",

        "needs_statements": """## Output Structure
## Formal Needs Statements
- [Persona] needs a way to [action] so that [benefit].
- (One statement per line — restate participant input in this exact format)

## Needs Grouped by Theme
### [Theme Name]
- (Related needs that address the same underlying problem)

## Highest-Priority Needs
- (Top 3 needs most cited or most impactful — these should anchor the ideation phase)""",

        "assumptions_grid": """## Output Structure
## High-Risk / Uncertain — Must Validate First
- (Assumptions or questions that are both high-stakes and unproven)

## High-Risk / Certain — Monitor Closely
- (Known facts that are high-stakes — ensure team stays aligned)

## Low-Risk / Uncertain — Worth Exploring
- (Open questions that are good to answer but not blocking)

## Top 3 Most Pressing Unknowns
1. [Unknown] — Why it matters / what to do to validate it
2. [Unknown] — Why it matters / what to do to validate it
3. [Unknown] — Why it matters / what to do to validate it""",

        "big_idea_vignettes": """## Output Structure
Group ideas into 3–5 thematic clusters:

## [Theme: e.g., Automation & AI]
- **[Idea Headline]:** one-sentence description
- (All ideas in this cluster)

## [Theme: e.g., Process Clarity & Ownership]
- **[Idea Headline]:** one-sentence description

(Repeat for each theme)

## Standout Ideas
- (Ideas that appeared more than once or generated notable enthusiasm)

## Ideas for Fast Follow-up
- (Ideas simple enough to act on quickly — quick wins)""",

        "prioritization_grid": """## Output Structure
## No-Brainers — High Importance, High Feasibility
1. [Idea] — why it's both important and achievable
2. (List all no-brainer items — these are immediate priorities)

## Big Bets — High Importance, Lower Feasibility
1. [Idea] — what makes it hard, why it's still worth pursuing
2. (List big bets — these need planning and investment)

## Utilities — Lower Importance, High Feasibility
- [Idea] — worth doing eventually but not the top priority

## Top 5 Prioritized Opportunities
1. [Opportunity name and one-sentence rationale]
2.
3.
4.
5.""",

        "tobe_process_map": """## Output Structure
## Improved Process Overview
- (What fundamentally changes from the As-Is state)

## Future-State Process Steps
- **Step 1 — [Step Name]:** What changes. Which pain points resolved. New capability required.
- (Continue for each step)

## Pain Points Resolved
- (Explicit mapping: old pain point → how the new process addresses it)

## New Capabilities Required
- (What tools, systems, or behaviors must exist for this future state to work)

## Quick Wins vs. Longer-Term Changes
- **Quick Wins:** changes that can be made immediately
- **Longer-term:** changes requiring system or org investment""",

        "experience_roadmap": """## Output Structure
## Near-Term (1–3 months) — Our User Can…
- (Capability statements starting with "Our user can…")

## Mid-Term (4–6 months) — Our User Can…
- (Capability statements starting with "Our user can…")

## Long-Term (7–12 months) — Our User Can…
- (Capability statements starting with "Our user can…")

## What We Will Learn at Each Stage
- Near-term learning: [what validated assumptions will tell us]
- Mid-term learning: [what to measure and refine]""",

        "hills_objectives": """## Output Structure
For each Hill/Objective:

## Objective [N]: [Short Name]
- **Who:** [the specific user this is for]
- **What:** [what they will be able to do that they can't today]
- **Wow:** [measurable outcome that proves success]
- **Pain Addressed:** [which current pain points this resolves]
- **Deliverables:** [concrete outputs required to achieve this Hill]
- **Quantified Value:** [estimated time saved, errors reduced, or experience improved]

(Repeat for each Hill — aim for 2–3 total)""",

        "gantt_roadmap": """## Output Structure
## Short-Term (Months 1–3)
- [Initiative name] — owner, effort level (Low/Med/High), dependencies

## Mid-Term (Months 4–6)
- [Initiative name] — owner, effort level, dependencies

## Long-Term (Months 7–12)
- [Initiative name] — owner, effort level, dependencies

## Dependencies & Sequencing
- (What must happen before what — flag blockers)

## High-Priority Items
- (Initiatives that are blockers for the rest of the roadmap)""",

        "resource_plan": """## Output Structure
## Resource Summary by Initiative
| Initiative | Owner | Skills Required | Effort | Dependencies | Gap? |
|------------|-------|-----------------|--------|--------------|------|
(Fill in each initiative)

## Ownership Gaps — Needs Assignment
- (Initiatives with no clear owner — flag as risks)

## Cross-Team Dependencies
- (Where multiple teams must coordinate — identify the connection point)

## Recommended First Hire or Partnership
- (If resource gaps exist, what role or team should be engaged first)""",

        "feedback_grid": """## Output Structure
## Things That Worked
- (What participants said was effective — keep doing these)

## Things to Change
- (What participants want done differently next time)

## Questions We Still Have
- (Open questions not resolved in the workshop)

## New Ideas to Try
- (Ideas sparked by the workshop that weren't in scope today)

## Top 2–3 Actions from This Feedback
- (Concrete changes to make based on the most-cited feedback items)""",

        "executive_summary": """## Output Structure
## Problem Statement
- (The core challenge this workshop addressed, in 2–3 sentences)

## Key Users
- (Who is affected — roles, teams, and their primary pain)

## Top 3 Insights
1. [Insight — what the team learned that they didn't know before]
2.
3.

## Top Pain Points
- (Most impactful friction points identified across all exercises)

## Top 5 Ideas
1. [Idea and brief rationale]
2.
3.
4.
5.

## Agreed Hills / Objectives
- (The 2–3 Hill statements the team committed to)

## Immediate Next Steps
- [Action] — Owner: [Name], Target: [Date]
- (2–4 concrete next steps with owners)""",

        "playback_deck": """## Output Structure
## Workshop Context & Objectives
- (Why we ran this workshop, what problem we set out to solve)

## Key Personas & Challenges
- (Who is affected and what their core experience looks like today)

## Key Insights from Empathy Work
- (The most important things we learned about how people experience the current state)

## Hills / Success Outcomes
- (The 2–3 Hills the team aligned on — WHO / WHAT / WOW)

## Roadmap Highlights
- (Near-term, mid-term, and long-term milestones)

## Next Steps & Asks
- (What leadership needs to support, decide, or resource)""",
    }

    ex_id = exercise_meta.get("id", "")
    output_format = output_formats.get(ex_id, """## Output Structure
## Key Themes
- (2–4 dominant themes from the discussion)

## Exercise Findings
- (Key findings organized by the exercise's specific framework)

## Notable Quotes
- "[quote]" — SenderName

## Key Takeaways
1. (Actionable takeaway)
2. (Actionable takeaway)
3. (Actionable takeaway)""")

    user_prompt = f"""## Workshop Exercise: {exercise_meta['name']}
**Phase {exercise_meta['phase']} — {exercise_meta['phase_name']}**
**Purpose:** {exercise_meta['description']}
**Desired Outcome:** {exercise_meta['desired_outcome']}

## Participant Discussion Transcript
({len(messages)} messages from {participant_count} participant{'s' if participant_count != 1 else ''}: {', '.join(participant_names)})

{transcript}

## Your Synthesis Task
{exercise_meta['prompt_hint']}

{output_format}

Ground every bullet in the transcript. Be specific — use participant names where it adds value. Write executive-ready prose: concise, direct, no filler. This output becomes the PPTX and DOCX deliverable for this exercise."""

    # Combine system + user into a single user message (gateway pattern)
    combined_prompt = f"{system_prompt}\n\n{user_prompt}"

    payload = json.dumps({
        "model":      LLM_MODEL,
        "max_tokens": 2500,
        "messages":   [{"role": "user", "content": combined_prompt}],
    }).encode()

    req = urllib.request.Request(
        LLM_GATEWAY_URL + "/v1/messages",
        data=payload,
        headers={
            "Authorization":      f"Bearer {token}",
            "Content-Type":       "application/json",
            "anthropic-version":  "2023-06-01",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, context=_ssl_ctx()) as resp:
            data = json.loads(resp.read())
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        # If 401/403, the token may have expired — clear it so next call re-reads from keychain
        if e.code in (401, 403):
            _keychain.pop("token", None)
        log.error("LLM Gateway HTTP %s: %s", e.code, body[:500])
        raise RuntimeError(f"LLM Gateway returned {e.code}: {body[:200]}") from e

    return data.get("content", [{}])[0].get("text", "")


# ── HTTP Routes ────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return send_from_directory(str(BASE_DIR), "index.html")


@app.route("/<path:filename>")
def static_files(filename):
    return send_from_directory(str(BASE_DIR), filename)


@app.route("/api/sessions", methods=["POST"])
def create_session():
    body = request.json or {}
    session_name = (body.get("session_name") or "").strip()
    exercise_order = body.get("exercise_order") or []

    if not session_name:
        return jsonify({"error": "session_name is required"}), 400
    if not exercise_order:
        return jsonify({"error": "exercise_order must not be empty"}), 400

    session_id = str(uuid.uuid4())
    room_code  = room_code_from_id(session_id)

    # Ensure room code uniqueness (collision is astronomically unlikely but handle it)
    attempts = 0
    while room_code in _room_code_map and attempts < 10:
        session_id = str(uuid.uuid4())
        room_code  = room_code_from_id(session_id)
        attempts  += 1

    # Build exercise state entries
    exercises_state: dict = {}
    for ex_id in exercise_order:
        exercises_state[ex_id] = {
            "messages": [],
            "summary": None,
            "summary_generated_at": None,
            "status": "pending",
        }

    session = {
        "session_id": session_id,
        "session_name": session_name,
        "room_code": room_code,
        "created_at": utcnow(),
        "exercise_order": exercise_order,
        "current_exercise_index": 0,
        "status": "active",
        "participants": [],
        "exercises": exercises_state,
    }

    save_session(session)
    _room_code_map[room_code] = session_id

    log.info("Created session %s (room %s) — %d exercises", session_id, room_code, len(exercise_order))
    return jsonify({"session_id": session_id, "room_code": room_code}), 201


@app.route("/api/sessions/test", methods=["POST"])
def create_test_session():
    """Create a session pre-filled with synthetic participant chats for every exercise."""
    import random

    FAKE_PARTICIPANTS = [
        {"name": "Sarah Chen",       "role": "Digital Marketing Manager"},
        {"name": "Marcus Rodriguez", "role": "Campaign Operations Lead"},
        {"name": "Priya Nair",       "role": "Marketing Analytics Analyst"},
        {"name": "James Okafor",     "role": "Brand Strategy Director"},
        {"name": "Lily Thompson",    "role": "Content & Creative Specialist"},
    ]

    SYNTHETIC_CHATS: dict[str, list[str]] = {
        "intros_outcomes": [
            "Hi everyone — Sarah Chen, Digital Marketing Manager. I run campaign delivery across 6 brands and I'm here because our trafficking process is costing us weeks of lead time that we don't have.",
            "Marcus Rodriguez, Campaign Operations. I own the trafficking workflows end-to-end and I want to get clear on which specific handoff points are creating the most delay so we can fix them.",
            "Priya Nair, Analytics. I'm joining because I'm often the last person in the chain — and by the time I get campaign parameters to set up tracking, we've already missed the launch window.",
            "James Okafor, Brand Strategy Director. My goal today is to walk out with a prioritized, actionable list of improvements we can actually resource and execute — not just a list of complaints.",
            "Lily Thompson, Content and Creative. I want the rest of this group to understand what it's like on the creative side — waiting 4 days to hear if an asset is approved, not knowing which version is final.",
            "Sarah: My main outcome for today — identify the top 2 or 3 root causes of trafficking delay and agree on who owns fixing them.",
            "Marcus: I want us to map the as-is process completely honestly, including the workarounds we've all built because the official process doesn't work.",
            "Priya: I want to leave with a clear definition of what 'done' looks like for each handoff — so analytics isn't always scrambling at launch.",
            "James: I need us to separate the problems we can solve this quarter from the platform changes that need a 12-month roadmap. Both matter, but they need different owners.",
            "Lily: Success for me is if Marcus and Sarah leave today with a real understanding of how much revision-chasing affects our output quality — and how fixable it actually is.",
            "Sarah: One more thing — I want to be honest that some of this is a brief quality problem, not just a trafficking problem. We're sending incomplete information and then wondering why things break downstream.",
            "Marcus: Agreed. And the feedback loop is broken too — when trafficking flags an issue, it often goes into a Slack thread that nobody owns, and the fix never makes it back to the original process.",
            "Priya: From analytics' view, the biggest desired outcome is same-day tracking setup at launch. Right now we average 2.3 days post-launch before tracking is live. That's data we can never recover.",
            "James: So let me restate what I'm hearing: we want a process where briefs arrive complete, traffickers can act immediately, and analytics is configured at launch. Is that a fair summary of what success looks like?",
            "Lily: Yes — and I'd add: creative knows which feedback is final before starting revisions. Right now I often do two full rounds before I find out the first round wasn't actually approved.",
        ],
        "hopes_fears": [
            "Hope: We leave today with a brief standard that every campaign manager agrees to use — not a suggestion, a requirement. Fear: We design a perfect process on paper but nothing changes because nobody owns enforcement.",
            "Hope: We identify the 2–3 changes that have the biggest impact on time-to-launch. Fear: We spend all day on root causes and run out of time to define solutions.",
            "Hope: Data quality in the first 72 hours post-launch improves to the point where we can actually do mid-flight optimization. Fear: Any solution requires platform changes that won't be funded until next fiscal year.",
            "Hope: Leadership backs the changes we recommend today with real budget and protected time. Fear: We define great Hills but they become 'good ideas' that get deprioritized the minute Q3 planning starts.",
            "Hope: Creative and Ops teams build genuine empathy for each other's constraints today. Fear: We go deep on creative-ops friction but ignore the Campaign Manager's role in upstream brief quality.",
            "Sarah: I hope we stop treating trafficking as a black box. Every campaign manager should understand what a trafficker actually does with a brief — that alone would improve brief quality by 50%.",
            "Marcus: My fear is that we conflate tool problems with process problems. A new platform won't fix unclear ownership. We need to agree on process first.",
            "Priya: I hope we get specific about what 'complete' means for analytics handoff. Right now it means different things to different people — which is how we end up with late tracking setup.",
            "James: Fear: we scope this too broadly and try to boil the ocean. I'd rather leave with 3 things we're actually going to do than 15 things nobody will resource.",
            "Lily: Hope: this workshop surfaces the invisible cost of revision cycles — how many hours creative spends re-doing work that was signed off, then un-signed off. That cost is real and nobody tracks it.",
            "Sarah: Fear that I have — we're going to identify root causes and point at Marketing Ops as the problem, when actually the problem starts with how briefs are written. Including by me.",
            "Marcus: Hope: by the end of today, everyone in this room understands the full end-to-end process — not just their slice. That shared understanding is worth more than any tool change.",
            "Priya: Fear: we build a beautiful new process flow and then discover that our trafficking platform literally cannot support it. We need a technical feasibility check built into whatever we decide.",
            "James: Hope: the Hills we write today are specific enough that we can actually measure success. 'Faster trafficking' is not a Hill. 'Zero campaigns with incomplete briefs by Q3' is.",
            "Lily: Fear: we run out of time to get to the roadmap and ideation. The root cause mapping is important but I want us to actually design the future state, not just document the broken present.",
        ],
        "persona_cards": [
            "The Campaign Manager juggles 15+ active campaigns simultaneously. There is no single place to check status — she cross-references Slack threads, email chains, a shared spreadsheet that's always out of date, and direct messages to traffickers.",
            "The Campaign Manager's biggest pain: she finds out about trafficking delays when something breaks, never proactively. She described her job as 'reactive by design' — the system forces her into firefighting mode.",
            "The Trafficking Specialist receives briefs with missing or conflicting fields on roughly 40% of submissions. He has built his own personal intake checklist to catch errors — but that checklist is in his head, not documented anywhere.",
            "The Trafficking Specialist spends an estimated 35–40% of his week on rework — re-entering corrected data, following up on missing specs, and managing platform errors caused by upstream gaps.",
            "The Analytics Lead is effectively locked out of mid-flight optimization because reporting data isn't trustworthy until 48–72 hours post-launch. She described the current state as 'flying blind for the first three days of every campaign.'",
            "The Analytics Lead has built shadow tracking using manual UTM spreadsheets as a workaround. She knows this is fragile and unsustainable but has no better option within the current workflow.",
            "The Creative Producer operates in a constant state of unclear prioritization. The word 'urgent' is used on 80%+ of requests, which means it signals nothing. She described it as 'everything is on fire, so nothing is.'",
            "The Creative Producer's biggest complaint isn't volume — it's revision ambiguity. She frequently completes a second or third round of revisions only to discover the first round was never formally approved.",
            "The Brand Director wants a real-time campaign health dashboard and has been asking for one for 18 months. What he receives instead is a static PowerPoint deck, updated manually before Monday reviews.",
            "The Brand Director spends 2–3 hours every Sunday manually reconciling campaign numbers from 3 different sources before his Monday morning leadership review. He considers this an unacceptable use of his time.",
            "Sarah: The Campaign Manager's representative quote — 'I have 18 campaigns running right now and I genuinely don't know the status of 6 of them without making 6 phone calls.'",
            "Marcus: The Trafficking Specialist's representative quote — 'I've gotten to the point where I don't trust any brief I receive on the first submission. I always assume I'll need to follow up at least once.'",
            "Priya: The Analytics Lead's representative quote — 'By the time I can trust the data, the window to optimize has already closed. We're doing post-mortems, not real-time management.'",
            "Lily: The Creative Producer's representative quote — 'I get told something is urgent, I drop everything to deliver it, and then it sits in someone's inbox for four days waiting for approval. That's demoralizing.'",
            "James: The Brand Director's representative quote — 'I ask for a dashboard and I get a PowerPoint. Every week. That's not visibility — that's theater.'",
        ],
        "stakeholder_map": [
            "Core process owners: Campaign Managers initiate the process with a brief and own campaign performance. Trafficking Specialists are the primary executors — every asset flows through them before going live.",
            "Decision makers: Brand Directors hold final approval authority on creative and budget reallocation. They are also the primary recipients of campaign performance reporting.",
            "Impacted parties — Paid Media: Media buyers lose optimization windows when trafficking is delayed. A 3-day delay at launch can mean 15–20% of the campaign budget is spent without optimization active.",
            "Impacted parties — Legal and Compliance: Legal reviews assets before trafficking, but the review SLA is frequently missed because briefs arrive late. When briefs are delayed, Legal is blamed for slowing down launch.",
            "Impacted parties — External Agencies: When Lilly runs agency-supported campaigns, agencies are waiting on the same incomplete briefs and experiencing the same delays — but they have contractual SLAs we're violating.",
            "Enablers — Marketing Operations: MarOps owns the trafficking platform (currently Workfront) but is not looped into process design decisions. They find out about workflow failures after they've already caused delays.",
            "Enablers — IT/Digital: IT owns platform integrations and data pipelines. Analytics' mid-flight reporting gap requires an IT-owned integration between the trafficking platform and the analytics stack.",
            "Key friction point 1: Brand Directors and Campaign Managers rarely speak directly during execution. Decisions made by Brand Directors (like creative changes mid-flight) don't reach Campaign Managers until work has already been done.",
            "Key friction point 2: Legal and Trafficking operate on different SLA clocks and don't coordinate. Legal approval is a prerequisite for trafficking, but neither team has visibility into the other's queue.",
            "Key friction point 3: Marketing Ops is accountable for the platform but has no ownership of the process. When the platform surfaces a problem, there's no clear owner to resolve it.",
            "Sarah: There's a missing role in this map — someone who owns cross-functional process health. Right now that accountability doesn't formally exist. James would be the logical owner, but it's not in anyone's job description.",
            "Marcus: The agency relationship is under-represented in how we think about this process. We treat agencies as external, but they're fully embedded in our workflow. Their delays become our delays.",
            "Priya: I'd add that IT is a bottleneck we never name. Any real fix to the analytics reporting gap requires IT prioritization, and we have zero leverage over that queue right now.",
            "James: This map is showing me that we have too many stakeholders with impact and not enough with accountability. We need to rationalize ownership before we redesign the process.",
            "Lily: Creative has no upstream visibility and no downstream feedback. We produce into a black hole — we don't know if our assets actually performed, which means we can't improve them.",
        ],
        "empathy_maps": [
            "Campaign Manager Sarah SAYS: 'I have no idea where things stand until something breaks.' THINKS: There's no system giving me a live view — I have to chase information manually every single day. FEELS: Chronically anxious and reactive. She described herself as 'always behind, even when I'm on top of things.' DOES: Sends 3–5 status check Slack messages per campaign per day. Maintains a personal tracking spreadsheet in parallel with the official system. PAINS: No proactive visibility — she learns about problems after they've already impacted the campaign. GAINS: Real-time status view that eliminates daily status chasing.",
            "Trafficking Specialist Marcus SAYS: 'I get briefs with missing fields constantly — probably 4 out of 10 have at least one gap I have to chase.' THINKS: The upstream process is broken and I've compensated by building my own safeguards. If I didn't, nothing would work. FEELS: Burdened and quietly resigned. He's normalized a level of friction that shouldn't exist. DOES: Has built a personal 22-point intake checklist that he applies to every brief before accepting it. Sends a follow-up email within 2 hours of receiving any brief. PAINS: Briefs arrive incomplete, forcing him to interrupt campaign managers mid-day for missing information. GAINS: A validated, complete brief on first submission — zero rework before he can begin trafficking.",
            "Analytics Lead Priya SAYS: 'Clean data at launch would change everything — right now we're doing post-mortems, not optimization.' THINKS: The integration problem is solvable but nobody has prioritized it. I could fix this in a sprint if I had IT resources. FEELS: Professionally frustrated — she has the analytical capability but not the data infrastructure to use it. DOES: Builds manual UTM spreadsheets as a workaround. Sets calendar reminders to check campaign launch notifications from Slack instead of receiving automated alerts. PAINS: Analytics setup is delayed 2–3 days post-launch because she isn't automatically notified when a campaign goes live. GAINS: Same-day analytics configuration at launch, enabling mid-flight optimization from day one.",
            "Creative Producer Lily SAYS: 'Everything is urgent — which means nothing is. And I never know which round of feedback is actually final.' THINKS: The feedback process is broken at the approval layer, not the creative layer. I deliver what's asked, then the ask changes. FEELS: Demoralized and undervalued. Her best work goes unseen because the review process obscures quality. DOES: Sends assets to Slack AND email AND the project management platform simultaneously — she doesn't trust any single channel to reach the approver. Checks for feedback notifications every 45 minutes when an asset is in review. PAINS: Unclear prioritization means all requests are treated as equally urgent, creating constant context-switching. GAINS: Single consolidated feedback round with clear approval status before any revision begins.",
            "Brand Director James SAYS: 'I ask for a dashboard and I get a PowerPoint. Every single week.' THINKS: My team is spending hours on manual reporting that should be automated. This is a resource waste I can see but can't seem to fix. FEELS: Frustrated by the gap between what the business should be capable of and what it actually delivers. DOES: Spends 2–3 hours every Sunday manually reconciling numbers from three different sources before Monday morning reviews. Has requested a live dashboard 4 times in 18 months. PAINS: Must manually build Monday review decks from static data exports because no live reporting exists. GAINS: Real-time campaign health dashboard that auto-updates, eliminating manual prep for leadership reviews.",
        ],
        "asis_process_map": [
            "Phase 1 — Brief Creation: Campaign Manager creates a brief in a Word document template. There is no required field validation — briefs can be submitted with empty sections. Average brief has 3–4 incomplete fields.",
            "Phase 1 continued: Brief is emailed simultaneously to Creative and Trafficking. There is no version control — if the brief changes after submission, teams receive a new email. The old version remains in inboxes.",
            "Phase 2 — Legal Review: Legal must approve all campaign creative before trafficking can begin. Legal is copied on the brief email but there is no formal handoff — they learn about the timeline from context clues, not a system.",
            "Phase 2 continued: Legal SLA is 3 business days. In practice, legal review takes 4.7 days on average because briefs arrive with missing legal disclaimers, triggering a back-and-forth loop before review can start.",
            "Phase 3 — Creative Development: Creative begins work upon receiving the brief email. Average revision cycles: 2.1 rounds. Primary cause of revision: brief specifications changed after creative began, or feedback was not consolidated before delivery.",
            "Phase 3 continued: Creative delivers assets via email attachment AND uploads to Workfront AND posts to a Slack channel — because different stakeholders check different places. Asset management is fragmented across 3 systems.",
            "Phase 4 — Trafficking: Specialist manually enters campaign specifications into the trafficking platform. This step takes 45–90 minutes per campaign. 30% of campaigns require re-entry due to spec errors traced back to the original brief.",
            "Phase 4 continued: Trafficking has no automated connection to Legal approval status. Traffickers must ask Campaign Managers if Legal has approved before proceeding — adding a step that should be a system trigger.",
            "Phase 5 — QA: QA team reviews the live trafficking setup against the original brief email. If the brief has been revised since submission, QA may be checking against an outdated version without knowing it.",
            "Phase 5 continued: QA errors traced to brief version mismatch account for 22% of all QA failures. These are preventable — they exist because there is no single source of truth for the current brief version.",
            "Phase 6 — Launch: Campaign goes live when Trafficking marks it ready in the platform. Analytics is not automatically notified. Priya finds out via a Slack message from Marcus, usually 4–6 hours after launch.",
            "Phase 6 continued: Analytics tracking setup takes 2–3 days post-launch due to manual parameter entry. Mid-flight data is unavailable during this window. Optimization opportunities in the first 72 hours are lost.",
            "Key handoff pain: Creative → Trafficking is not a formal handoff. Assets are uploaded to Workfront, but there is no task trigger or notification. Marcus manually checks for new assets twice a day.",
            "Key handoff pain: Campaign Manager → Legal has no SLA tracking. Legal has no visibility into campaign volume until briefs arrive. Peaks in campaign launches cause Legal queue backup with no advance warning.",
            "Overall pattern: Every stage of this process has at least one informal workaround built by someone trying to make it work. The workarounds are invisible to leadership and create single points of failure when the person who built them is unavailable.",
        ],
        "needs_statements": [
            "Campaign Managers need a way to see real-time trafficking status for all active campaigns so that they can proactively manage client expectations and stop spending 30% of their day on status chasing.",
            "Campaign Managers need a way to submit briefs with all required fields validated before submission so that downstream teams can begin work immediately without follow-up.",
            "Trafficking Specialists need a way to receive complete, validated briefs on first submission so that they can eliminate the intake triage step and reduce rework from 40% to under 10% of weekly hours.",
            "Trafficking Specialists need a way to see Legal approval status in the platform so that they can begin setup immediately upon approval without waiting for a manual notification from the Campaign Manager.",
            "Analytics team needs a way to receive automated campaign launch notifications with all parameters pre-populated so that they can configure tracking same-day and enable mid-flight optimization from day one.",
            "Analytics team needs a way to access a live connection between the trafficking platform and the analytics stack so that reporting data is available within 4 hours of launch without manual data entry.",
            "Creative team needs a way to receive consolidated, final feedback in a single submission so that they can deliver approved assets on the first revision and eliminate the back-and-forth loop.",
            "Creative team needs a way to know which campaign requests are highest priority so that they can allocate capacity intentionally instead of treating every request as equally urgent.",
            "Brand Directors need a way to access a live campaign health dashboard with auto-updating data so that they can conduct Monday reviews without manually reconciling numbers from multiple sources.",
            "Brand Directors need a way to receive proactive alerts when campaign performance deviates from expected benchmarks so that they can intervene before issues compound.",
            "Legal needs a way to see incoming campaign volume 5–7 business days in advance so that they can resource review queues before peaks create SLA misses.",
            "Marketing Ops needs a way to be included in process design decisions so that platform capabilities are aligned with process requirements before new workflows are implemented.",
            "All roles need a way to access a single authoritative version of the current brief so that every team is working from the same source of truth regardless of when they joined the process.",
            "All roles need a way to understand how their handoff affects the next step in the process so that brief quality and timing decisions are made with full awareness of downstream impact.",
        ],
        "assumptions_grid": [
            "HIGH RISK / UNCERTAIN — Must validate: Trafficking is the primary bottleneck in time-to-launch. If the actual bottleneck is Legal review or brief quality, optimizing trafficking will have minimal impact on lead time.",
            "HIGH RISK / UNCERTAIN — Must validate: Teams will adopt a standardized brief template without executive mandate. Voluntary adoption of new standards has failed twice in the past 3 years at Lilly.",
            "HIGH RISK / UNCERTAIN — Must validate: A live status dashboard will reduce the need for status check meetings and Slack messages. If teams don't trust the dashboard data, they'll continue to use manual channels.",
            "HIGH RISK / CERTAIN — Monitor closely: Campaign volume will increase 20–30% next year due to new product launches. Any process improvement must scale — solutions designed for current volume may break at higher load.",
            "HIGH RISK / CERTAIN — Monitor closely: IT prioritization is required for any platform integration fix. IT queue is typically 6–9 months for non-critical projects. Analytics fix may be delayed regardless of business priority.",
            "LOW RISK / UNCERTAIN — Worth exploring: A weekly cross-functional sync could replace most status check Slack messages. Worth piloting with 2 campaign managers for 30 days to validate the hypothesis.",
            "LOW RISK / UNCERTAIN — Worth exploring: Creative output quality would improve if the revision process were more structured. Hypothesis: one formal review round with consolidated feedback produces better outcomes than 3 informal rounds.",
            "LOW RISK / CERTAIN: Teams want better tooling. Every persona in today's workshop named a tool gap. Appetite for change exists — the question is readiness and resourcing.",
            "KEY QUESTION — Must answer before proceeding: Who owns enforcement of the brief standard? Without a named owner and consequence for non-compliance, the brief template will be ignored within 60 days.",
            "KEY QUESTION — Must answer before proceeding: Does leadership understand the cost of the current state? James expressed concern that the upstream cost (Creative rework hours, Analytics delay cost) has never been quantified for leadership.",
            "KEY QUESTION — Must answer before proceeding: Can the existing trafficking platform (Workfront) support automated analytics notifications, or does this require a new integration? If new integration, IT estimate is needed before committing to the roadmap.",
            "Sarah: I'd add one more critical assumption — we're assuming the brief template will be adopted by Campaign Managers, but Campaign Managers haven't been in the room today to agree to it. We need their buy-in before we finalize anything.",
            "Marcus: The biggest unknown for me is whether Legal can actually do 3-day reviews at current volume. We've never tested whether the 4.7-day average is a resource problem or a process problem. Those require different solutions.",
        ],
        "big_idea_vignettes": [
            "Theme: Unified Campaign Workspace. Idea: A single platform where the brief, creative, trafficking specs, Legal approval, and analytics setup all live in one place — no email handoffs, no version conflicts, no duplicate entry.",
            "Theme: Unified Campaign Workspace. Idea: Auto-populate trafficking fields from the approved brief. When a Campaign Manager marks a brief complete, Workfront pre-fills the trafficking setup with no manual re-entry. Eliminates 30% of trafficking rework.",
            "Theme: Unified Campaign Workspace. Idea: Auto-notify Analytics at launch with all campaign parameters pre-populated. A system trigger when trafficking marks a campaign live sends a structured data packet to the analytics stack.",
            "Theme: Brief Quality at Intake. Idea: An AI-powered brief validator that checks for required fields, conflicting specs, and missing legal language before the brief can be submitted — like a spell-checker for campaign briefs.",
            "Theme: Brief Quality at Intake. Idea: A 'brief readiness score' visible to Campaign Managers before submission. Color-coded (green / yellow / red) against a checklist of what trafficking and legal need. Forces self-review before handoff.",
            "Theme: Brief Quality at Intake. Idea: Mandatory brief template with hard-required fields. No submission without completion. Sounds obvious — but we've never enforced it. This alone would eliminate 40% of trafficking follow-up.",
            "Theme: Visibility and Status. Idea: A live campaign status dashboard for Brand Directors — auto-updated from the trafficking platform, showing trafficking status, Legal approval, and launch date with no manual input required.",
            "Theme: Visibility and Status. Idea: A 'campaign heartbeat' Slack notification — automated daily digest showing status of all in-flight campaigns, replacing the 15+ individual status messages Campaign Managers currently send.",
            "Theme: Visibility and Status. Idea: Traffic light indicators in Workfront visible to all stakeholders — green (on track), yellow (at risk), red (delayed) — with a mandatory comment required to move from green to yellow.",
            "Theme: Process Coordination. Idea: A weekly 20-minute cross-functional launch review — Campaign Managers, Trafficking, Legal, Analytics, and Creative — replacing 40+ ad-hoc status messages. Fixed agenda. Rotating facilitator.",
            "Theme: Process Coordination. Idea: Legal volume forecasting — Campaign Managers submit a 2-week forward look of planned briefs. Legal uses this to staff reviews proactively instead of reacting to peaks.",
            "Theme: Process Coordination. Idea: A formal Creative-to-Trafficking handoff trigger in Workfront. When Creative marks assets as final, an automated task is created in Marcus's queue — no more manual checking twice a day.",
            "Theme: Analytics and Measurement. Idea: Real-time mid-flight optimization dashboard — live performance data available from hour 1 of launch, enabling intra-campaign budget reallocation based on actual performance.",
            "Theme: Analytics and Measurement. Idea: A campaign performance retrospective template — structured 30-minute post-campaign review that captures what worked, what didn't, and feeds back into brief quality for the next campaign.",
        ],
        "prioritization_grid": [
            "HIGH IMPORTANCE / HIGH FEASIBILITY — No-Brainer: Standardized brief template with required field enforcement. Zero cost, immediate impact, within our control. Brief quality is root cause of 40%+ of downstream rework. Do this first.",
            "HIGH IMPORTANCE / HIGH FEASIBILITY — No-Brainer: Weekly cross-functional launch review (20 min). Replaces fragmented Slack status checking. Can be started next week with no tools or budget. James to own calendar invite.",
            "HIGH IMPORTANCE / HIGH FEASIBILITY — No-Brainer: Formal Creative-to-Trafficking handoff trigger in Workfront. IT estimates 2 days of configuration. Eliminates Marcus's twice-daily manual asset check immediately.",
            "HIGH IMPORTANCE / HIGH FEASIBILITY — No-Brainer: Legal volume forecasting — Campaign Managers submit 2-week forward brief outlook. No tools required, just a process change. Legal gets to staff proactively instead of reactively.",
            "HIGH IMPORTANCE / LOWER FEASIBILITY — Big Bet: Live campaign status dashboard for Brand Directors. High value — eliminates James's Sunday manual reconciliation. Requires IT prioritization and 3–6 month build. Needs exec sponsorship.",
            "HIGH IMPORTANCE / LOWER FEASIBILITY — Big Bet: Auto-populated trafficking fields from approved brief. Requires Workfront configuration and data mapping. Eliminates 30% of trafficking rework. 2–3 month project with IT involvement.",
            "HIGH IMPORTANCE / LOWER FEASIBILITY — Big Bet: Automated analytics notification at campaign launch. Requires integration between Workfront and the analytics stack. IT-dependent. 4–6 month timeline. Highest analytical ROI.",
            "LOWER IMPORTANCE / HIGH FEASIBILITY — Utility: Campaign heartbeat Slack digest — automated daily status summary. Easy to build, nice to have. But if the brief template and weekly sync are working, this may become unnecessary.",
            "LOWER IMPORTANCE / HIGH FEASIBILITY — Utility: Brief readiness score visible before submission. Good UX improvement. Doesn't solve the root problem if Campaign Managers can still submit incomplete briefs — needs hard enforcement to matter.",
            "LOWER IMPORTANCE / LOWER FEASIBILITY — Revisit: AI-powered brief validator. Technically interesting but complex and unproven at Lilly. Not a Q1 or Q2 priority. Revisit when foundational fixes are in place.",
            "Team consensus emerged: Start with the 4 no-brainers — brief template, weekly sync, Workfront handoff trigger, Legal forecasting. These can all be running within 30 days and cost essentially nothing.",
            "Big bets priority order: (1) auto-populated trafficking fields, (2) analytics auto-notification, (3) Brand Director dashboard. Sequence matters — brief fix must come before trafficking automation or we automate bad data.",
            "James: I'm committing to sponsor the dashboard project in Q2 if the no-brainer improvements show measurable reduction in rework in Q1. We need to demonstrate ROI before asking IT for 6 months of engineering time.",
        ],
        "tobe_process_map": [
            "Phase 1 — Brief Creation (To-Be): Campaign Manager completes a standardized digital brief form with required field validation. The system blocks submission if mandatory fields — including legal disclaimers, asset specs, and trafficking parameters — are empty.",
            "Phase 1 continued: Upon submission, the brief is automatically versioned and stored as the single source of truth. All downstream teams see the same current version. No more email attachments or version confusion.",
            "Phase 2 — Legal Review (To-Be): Legal receives an automated notification with the complete brief and a pre-calculated review deadline based on the campaign launch date. Legal queue visibility allows proactive resourcing.",
            "Phase 2 continued: Legal approval is captured in the platform as a status flag. When Legal approves, Trafficking receives an automatic task with all brief fields pre-populated. No manual notification step.",
            "Phase 3 — Creative (To-Be): Creative receives a structured brief with prioritization score (High / Standard / Low). All feedback is collected in a single consolidated review round. Creative marks assets final in the platform — triggering the Trafficking task.",
            "Phase 3 continued: Only one round of revision is formally supported. If additional changes are required after the first revision, a new mini-brief is required — making scope creep visible and accountable.",
            "Phase 4 — Trafficking (To-Be): Trafficking setup begins from pre-populated brief data. Spec re-entry is eliminated for standard fields. Trafficker validates and confirms — this step now takes 15–20 minutes instead of 45–90.",
            "Phase 4 continued: QA reviews against the live, versioned brief — not an email. Version mismatch errors are eliminated. QA pass rate is expected to improve from 78% to 95%+ in the first quarter.",
            "Phase 5 — Launch (To-Be): When Trafficking marks the campaign live, an automated trigger sends campaign parameters to Analytics. Tracking is configured within 4 hours of launch. Mid-flight data is available same day.",
            "Phase 5 continued: Brand Directors receive a real-time dashboard view of all active campaigns — trafficking status, launch date, early performance indicators. Monday review prep time drops from 2–3 hours to 15 minutes.",
            "Pain resolved: Brief quality gaps — eliminated by required field validation and submission block. Rework caused by incomplete briefs drops from 40% to under 5%.",
            "Pain resolved: Analytics delay — automated launch notification and pre-populated parameter handoff enables same-day tracking. The 2.3-day average delay is eliminated.",
            "Pain resolved: Status visibility — live platform status visible to all stakeholders eliminates reactive status chasing and reduces Campaign Manager Slack messages by an estimated 70%.",
            "New capability required: Workfront configuration for required field validation, automated triggers, and versioned brief storage. IT estimate: 6–8 weeks of configuration work.",
            "Quick wins (can start immediately): Brief template standard with enforcement, weekly launch review meeting, Creative-to-Trafficking handoff trigger, Legal volume forecasting process.",
        ],
        "experience_roadmap": [
            "NEAR-TERM (0–3 months): Our users can submit campaign briefs using a standardized template with required field validation — eliminating incomplete brief submissions from day one of rollout.",
            "NEAR-TERM (0–3 months): Our users can see cross-functional campaign status in a single weekly 20-minute review meeting, replacing fragmented Slack status checking.",
            "NEAR-TERM (0–3 months): Our trafficking specialists can begin work immediately upon Creative sign-off via an automated Workfront handoff trigger — eliminating the manual twice-daily asset check.",
            "MID-TERM (3–6 months): Our Campaign Managers can track trafficking status in real time without sending a single Slack message — through a live status view in the campaign platform.",
            "MID-TERM (3–6 months): Our trafficking specialists can set up campaigns in 15–20 minutes instead of 45–90 minutes, because spec fields are auto-populated from the approved brief.",
            "MID-TERM (3–6 months): Our analytics team receives automated launch notifications with pre-populated parameters, enabling tracking configuration within 4 hours of launch — down from 2.3 days.",
            "LONG-TERM (6–12 months): Our Brand Directors can access a live, auto-updating campaign health dashboard that replaces the manual PowerPoint prep before Monday reviews.",
            "LONG-TERM (6–12 months): Our analytics team can run mid-flight campaign optimization from day one of launch, because tracking data is live same-day and connected to budget controls.",
            "LONG-TERM (6–12 months): Our Campaign Managers can predict Legal review timelines and plan launches accordingly — because Legal has a forward visibility tool and manages reviews proactively.",
            "LONG-TERM (12–18 months): Our teams can work within a fully integrated campaign platform where brief, creative, trafficking, legal approval, and analytics are connected with no manual handoffs.",
            "Learning at each stage — Near-term: Does brief quality actually improve with enforcement, or do Campaign Managers find workarounds? Measure: % of briefs requiring follow-up within 30 days.",
            "Learning at mid-term: Does trafficking auto-population actually reduce rework, or do edge cases still require manual correction? Measure: trafficking re-entry rate before vs. after.",
            "Learning at long-term: Does the dashboard actually change how Brand Directors make decisions, or do they revert to manual review? Measure: Sunday prep time and dashboard active usage rate.",
        ],
        "hills_objectives": [
            "Hill 1: A Campaign Manager with 15 active campaigns can see the real-time trafficking status of every campaign without sending a single Slack message or email — from any device, at any time.",
            "Hill 1 — WHO: Campaign Managers across all brand teams. WHAT: Real-time status visibility without manual status requests. WOW: Zero campaign status Slack messages per week, measured at 30 days post-launch.",
            "Hill 1 — Pain addressed: Chronic reactive posture, daily status chasing, anxiety about missing delays. Deliverables: live status dashboard, automated status triggers in Workfront. Timeline: Q2.",
            "Hill 2: A Trafficking Specialist receives zero incomplete briefs by end of Q3 2026 — measured as briefs that require at least one follow-up before trafficking can begin.",
            "Hill 2 — WHO: Trafficking Specialists. WHAT: Complete, validated briefs on first submission. WOW: Brief rework rate drops from 40% to under 5% of weekly hours.",
            "Hill 2 — Pain addressed: Manual intake triage, chronic rework, 45–90 min setup time per campaign. Deliverables: mandatory brief template with field validation, brief readiness score, enforcement process. Timeline: Q1.",
            "Hill 3: The Analytics team can configure campaign tracking and have live data available within 4 hours of campaign launch — down from the current 2.3-day average.",
            "Hill 3 — WHO: Analytics team. WHAT: Same-day tracking configuration via automated parameter handoff. WOW: Mid-flight optimization active on 100% of campaigns from launch day.",
            "Hill 3 — Pain addressed: Blind first 72 hours of every campaign, no mid-flight optimization, manual UTM workaround. Deliverables: Workfront-to-Analytics integration, automated launch trigger. Timeline: Q3 with IT dependency.",
            "Quantified value across all 3 Hills: 30% reduction in trafficking rework hours (15 hrs/week saved), 2.3 days recovered in time-to-reporting per campaign, 2–3 hours/week saved for Brand Directors on manual prep.",
            "James: I want to be explicit that Hill 2 is the prerequisite for Hills 1 and 3. If briefs are still incomplete, automating downstream steps just automates bad data. Brief quality must be fixed first.",
        ],
        "gantt_roadmap": [
            "Month 1: Design and socialize standardized brief template with all campaign managers. Get explicit sign-off from top 5 brief-writers. Marcus and Sarah co-own. Target: template live in Workfront by end of Month 1.",
            "Month 1–2: Establish weekly cross-functional launch review cadence. James facilitates first 4 sessions to establish norms. Hand off facilitation rotation to Campaign Ops after 30 days.",
            "Month 2: Configure Creative-to-Trafficking handoff trigger in Workfront. IT estimate: 2 days of configuration. Marcus defines trigger logic. Go live end of Month 2.",
            "Month 2–3: Implement Legal volume forecasting process. Campaign Managers submit 2-week brief outlook every Friday. Legal uses to staff proactively. Priya tracks on-time Legal review rate as leading indicator.",
            "Month 3: Measure Q1 no-brainer impact. KPIs: brief rework rate (target: <10%), Legal SLA miss rate (target: <5%), Campaign Manager Slack status messages (target: -50%). Decision gate for Q2 investment.",
            "Month 3–5: Configure Workfront auto-population of trafficking fields from approved brief. IT owns with Marketing Ops as product owner. Marcus defines data mapping requirements. Target: 10 campaign pilot in Month 5.",
            "Month 4–6: Build analytics auto-notification integration. Priya defines parameter schema. IT builds Workfront-to-Analytics data trigger. Target: pilot on 3 campaigns in Month 6 before full rollout.",
            "Month 5–7: Brand Director dashboard design and build. James defines KPIs and layout requirements. IT and Marketing Ops co-own build. Target: live for Monday review use in Month 7.",
            "Month 6: Mid-point review of all Q2 big-bet initiatives. Measure trafficking setup time (target: <20 min), analytics configuration time (target: <4 hours), dashboard active usage rate.",
            "Month 9: Full rollout of integrated workflow — brief → trafficking auto-population → analytics auto-trigger → live dashboard. All campaign managers required to use new workflow. Marketing Ops owns ongoing support.",
            "Month 10–12: Evaluate unified platform options if integrated workflow shows gaps. RFP only if KPIs are not met by Month 9. James owns the go/no-go decision based on Q3 data.",
            "Dependencies: IT availability is the critical path for Months 3–7. Workfront auto-population, analytics integration, and dashboard all require IT. Marcus to submit IT requests in Month 1 to secure queue position.",
        ],
        "resource_plan": [
            "Initiative: Brief Template + Enforcement. Owner: Marcus Rodriguez (Campaign Ops). Support: Sarah Chen (Campaign Managers). Skills: process design, Workfront configuration. Effort: Low (2 weeks). Cost: $0. No gaps.",
            "Initiative: Weekly Cross-Functional Launch Review. Owner: James Okafor (Brand Strategy). Facilitation rotation after Month 1. Skills: meeting design, facilitation. Effort: Low (ongoing 20 min/week). Cost: $0. No gaps.",
            "Initiative: Creative-to-Trafficking Workfront Trigger. Owner: Marketing Ops + IT. Skills: Workfront configuration. Effort: Low (2 days IT). Cost: $0 if internal. Risk: IT queue dependency — submit request Month 1.",
            "Initiative: Legal Volume Forecasting Process. Owner: Marcus Rodriguez + Legal lead (TBD). Skills: process design, calendar coordination. Effort: Low. Cost: $0. Gap: Legal lead has not been identified as owner yet.",
            "Initiative: Trafficking Auto-Population from Brief. Owner: Marcus Rodriguez (business), IT (technical). Skills: Workfront data mapping, API configuration. Effort: Medium (6–8 weeks IT). Cost: $15–25K if IT capacity purchased.",
            "Initiative: Analytics Launch Auto-Notification. Owner: Priya Nair (business), IT (technical). Skills: Workfront-to-analytics integration, API development. Effort: High (12–16 weeks IT). Cost: $40–60K. Highest ROI initiative.",
            "Initiative: Brand Director Live Dashboard. Owner: James Okafor (business), IT + Marketing Ops (technical). Skills: data visualization, Workfront API, BI platform. Effort: High (16–20 weeks). Cost: $60–100K. Requires exec funding approval.",
            "Exec Sponsor: James Okafor owns budget approval and removes blockers for all Q2+ initiatives. Must formally commit to Q2 funding in writing by end of Month 1.",
            "External resource need: May need 1 FTE contractor for platform integration work in Months 3–9 if IT internal capacity is insufficient. Estimated cost: $80–120K for 6-month engagement.",
            "Skills gap: No internal campaign analytics engineer. The Workfront-to-analytics integration requires skills Priya's team doesn't have. Options: (1) IT internal, (2) contractor, (3) analytics platform vendor professional services.",
            "Budget summary: No-brainer initiatives: $0. Q2 big bets: $55–85K (trafficking auto-pop + analytics trigger). Q3 dashboard: $60–100K. Total investment range: $115–185K for full roadmap. Excludes ongoing maintenance.",
            "Marcus: The biggest resource gap is IT prioritization, not budget. Every technical initiative on this roadmap requires IT queue position. If we don't submit requests in Month 1, we'll be waiting until Q3 to start.",
        ],
        "feedback_grid": [
            "WORKED WELL: The empathy map exercise. For the first time, Marcus understood why Lily sends assets to three different places — and Lily understood why Marcus can't just 'start on the brief' when fields are missing. Real empathy was built.",
            "WORKED WELL: The process mapping exercise. Seeing the whole as-is flow on one page made it obvious where the friction is. The visual was more powerful than any summary document.",
            "WORKED WELL: James setting a clear scope at the start — 'I want 3 things we're actually going to do.' That framing kept us from going too broad and helped us prioritize ruthlessly.",
            "WORKED WELL: The assumptions grid surfaced a critical assumption nobody had named — that brief quality enforcement requires a named owner, not just a template. That insight is worth the whole day.",
            "WORKED WELL: Priya quantifying the analytics delay cost — 2.3 days of lost optimization per campaign. Putting a number on it changed how seriously James is taking the analytics integration.",
            "CHANGE: The stakeholder map felt rushed. We needed more time to map the Legal and Agency relationships — those turned out to be more important than we initially thought.",
            "CHANGE: We should have had a Campaign Manager in the room today. We made assumptions about brief quality and adoption that we can't validate without their perspective.",
            "CHANGE: The big ideas session needed more time. We generated 14 ideas in 20 minutes — the quality was high but the discussion was shallow. Consider 40 minutes next time.",
            "CHANGE: Start with the process map, not introductions. Understanding the as-is flow is the foundation for everything else. Front-loading it would make the empathy exercise more concrete.",
            "NEW IDEA: Record a 3-minute video summary of the day's outputs — Hills, top ideas, immediate next steps — and share with stakeholders who weren't in the room. Builds alignment without another meeting.",
            "NEW IDEA: Run a follow-up 'brief writer' workshop specifically with Campaign Managers. They need to understand what happens downstream of their brief — that visibility alone would improve quality.",
            "QUESTION: Who owns the brief enforcement process? We named the template but didn't name an enforcer. Without enforcement, adoption will decay within 60 days — we've seen this before.",
            "QUESTION: How do we keep momentum? The risk is that this session becomes a great memory but not a catalyst. We need a 30-day check-in with the same group to review Hill 2 progress.",
            "OVERALL: This was the best DTW we've run. High energy, honest conversation, specific outputs. The Hills are real and the no-brainers are already assigned. The question now is execution.",
        ],
        "executive_summary": [
            "Core problem: Campaign trafficking delays at Lilly Digital Marketing stem from three interdependent root causes — incomplete brief submissions, fragmented cross-functional visibility, and disconnected analytics setup at launch.",
            "Workshop context: 5 participants representing Campaign Management, Campaign Operations, Analytics, Brand Strategy, and Creative. Full-day Design Thinking Workshop using IBM Enterprise Design Thinking methodology.",
            "Key insight 1: All 5 personas experience the same broken process from entirely different vantage points. The dysfunction is structural, not behavioral. Teams aren't failing — the system is failing them.",
            "Key insight 2: Every role has built a personal workaround to compensate for process gaps. These workarounds are invisible to leadership and create single points of failure when the person who built them is unavailable.",
            "Key insight 3: Brief quality is the root cause upstream of every downstream problem. Fixing trafficking, analytics, and visibility without fixing brief quality would automate bad data through a faster process.",
            "Top pain theme 1 — Brief Quality: 40% of briefs require at least one follow-up before trafficking can begin. Trafficking rework accounts for 35–40% of Marcus's weekly hours. Root cause: no required field validation, no submission enforcement.",
            "Top pain theme 2 — Visibility: Campaign Managers send 3–5 status Slack messages per campaign per day. Brand Directors spend 2–3 hours every Sunday manually reconciling data for Monday reviews. No live status view exists.",
            "Top pain theme 3 — Analytics Delay: Tracking setup averages 2.3 days post-launch. Mid-flight optimization is impossible in the first 72 hours. Root cause: no automated launch notification, no parameter handoff from trafficking.",
            "Top 5 ideas prioritized: (1) Mandatory brief template with required field enforcement, (2) Weekly cross-functional launch review, (3) Auto-populated trafficking fields from approved brief, (4) Automated analytics launch notification, (5) Live Brand Director campaign dashboard.",
            "Agreed Hills: Hill 1 — Zero Slack status messages from Campaign Managers (real-time visibility). Hill 2 — Zero incomplete briefs by Q3 (brief quality enforcement). Hill 3 — Analytics live within 4 hours of launch.",
            "Immediate next steps: James approves brief template standard by end of week. Marcus configures Workfront handoff trigger in Month 1. Both submit IT requests for Q2 platform work before next Friday.",
            "Investment required: No-brainer fixes cost $0 and can be live within 30 days. Q2 platform improvements estimated at $55–85K. Q3 dashboard at $60–100K. Total roadmap: $115–185K over 12 months.",
        ],
        "playback_deck": [
            "We started this workshop with one question: why does it take so long to get a campaign live, and what will it take to fix it?",
            "We mapped the current state — an as-is process with 6 phases, 12+ documented pain points, and at least 8 informal workarounds built by individuals to make a broken process function.",
            "We built empathy maps for all 5 personas in the room. The most important thing we learned: everyone is suffering from the same systemic failures, but experiencing them from completely different vantage points.",
            "We identified the root causes: incomplete briefs upstream, no cross-functional visibility mid-process, and a disconnected analytics setup at launch. These three causes account for roughly 80% of our time-to-launch delays.",
            "We generated 14 big ideas and prioritized them using a feasibility-vs-impact grid. 4 ideas are no-brainers — high impact, low cost, within our control. We can start all 4 next week.",
            "We defined 3 Hills that represent what success looks like: Campaign Managers with real-time status visibility, zero incomplete briefs, and same-day analytics at launch.",
            "Our roadmap: Month 1 — brief template, weekly sync, Workfront trigger, Legal forecasting. Months 3–6 — trafficking auto-population, analytics integration. Months 6–12 — Brand Director dashboard, full platform integration.",
            "The investment: $0 for the no-brainers. $115–185K for the full roadmap over 12 months. The cost of the current state — in rework hours, optimization windows lost, and manual prep time — far exceeds that investment.",
            "What we need from leadership: Approve the brief template standard this week. Commit Q2 budget for trafficking and analytics improvements. Protect IT capacity for the integrations on the roadmap.",
            "Next step — everyone in this room: Brief template reviewed and approved by Friday. IT requests submitted by Marcus and Priya before next Friday. 30-day check-in scheduled before we leave today.",
        ],
    }

    # Build all exercises with synthetic messages
    all_ex_ids = [
        "intros_outcomes", "persona_cards", "executive_summary",
        "hopes_fears", "stakeholder_map", "empathy_maps",
        "asis_process_map", "needs_statements", "assumptions_grid",
        "big_idea_vignettes", "prioritization_grid", "tobe_process_map",
        "experience_roadmap", "hills_objectives", "gantt_roadmap",
        "resource_plan", "feedback_grid", "playback_deck",
    ]

    session_id = str(uuid.uuid4())
    room_code  = room_code_from_id(session_id)
    while room_code in _room_code_map:
        session_id = str(uuid.uuid4())
        room_code  = room_code_from_id(session_id)

    exercises_state: dict = {}
    for ex_id in all_ex_ids:
        lines = SYNTHETIC_CHATS.get(ex_id, [])
        messages = []
        for i, text in enumerate(lines):
            participant = FAKE_PARTICIPANTS[i % len(FAKE_PARTICIPANTS)]
            messages.append({
                "sender":    participant["name"],
                "text":      text,
                "timestamp": f"2026-03-18T{10 + i // 60:02d}:{i % 60:02d}:00Z",
                "role":      participant["role"],
            })
        exercises_state[ex_id] = {
            "messages": messages,
            "summary": None,
            "summary_generated_at": None,
            "status": "pending",
        }
    exercises_state[all_ex_ids[0]]["status"] = "active"

    session = {
        "session_id": session_id,
        "session_name": "⚗ Test Session — Synthetic Data",
        "room_code": room_code,
        "created_at": utcnow(),
        "exercise_order": all_ex_ids,
        "current_exercise_index": 0,
        "status": "active",
        "participants": [{"name": p["name"], "role": p["role"], "joined_at": utcnow()} for p in FAKE_PARTICIPANTS],
        "exercises": exercises_state,
    }

    save_session(session)
    _room_code_map[room_code] = session_id
    log.info("Created test session %s (room %s)", session_id, room_code)
    return jsonify({"session_id": session_id, "room_code": room_code}), 201


@app.route("/api/sessions", methods=["GET"])
def list_sessions():
    sessions = []
    if _s3_enabled():
        try:
            paginator = _s3().get_paginator("list_objects_v2")
            prefix = f"{S3_PREFIX}/sessions/"
            objs = []
            for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=prefix):
                objs.extend(page.get("Contents", []))
            # Sort by LastModified descending
            objs.sort(key=lambda o: o["LastModified"], reverse=True)
            for obj in objs:
                try:
                    raw = _s3().get_object(Bucket=S3_BUCKET, Key=obj["Key"])
                    s = json.loads(raw["Body"].read().decode("utf-8"))
                    sessions.append({
                        "session_id": s["session_id"],
                        "session_name": s["session_name"],
                        "room_code": s["room_code"],
                        "created_at": s["created_at"],
                        "status": s["status"],
                        "participant_count": len(s.get("participants", [])),
                        "exercise_count": len(s.get("exercise_order", [])),
                        "current_exercise_index": s.get("current_exercise_index", 0),
                    })
                except Exception:
                    pass
        except Exception as exc:
            log.error("list_sessions S3 error: %s", exc)
    else:
        for p in sorted(SESSIONS_DIR.glob("session_*.json"), key=lambda x: x.stat().st_mtime, reverse=True):
            try:
                with open(p, encoding="utf-8") as f:
                    s = json.load(f)
                sessions.append({
                    "session_id": s["session_id"],
                    "session_name": s["session_name"],
                    "room_code": s["room_code"],
                    "created_at": s["created_at"],
                    "status": s["status"],
                    "participant_count": len(s.get("participants", [])),
                    "exercise_count": len(s.get("exercise_order", [])),
                    "current_exercise_index": s.get("current_exercise_index", 0),
                })
            except Exception:
                pass
    return jsonify(sessions)


@app.route("/api/sessions/<session_id>", methods=["GET"])
def get_session(session_id):
    # Allow lookup by room code too
    if len(session_id) == 6 and session_id.upper() in _room_code_map:
        session_id = _room_code_map[session_id.upper()]
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404
    # Strip passcode from response
    safe = {k: v for k, v in s.items() if k != "passcode"}
    return jsonify(safe)


@app.route("/api/sessions/<session_id>", methods=["DELETE"])
def delete_session(session_id):
    """Delete a session and all its artifacts."""
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    room_code = s.get("room_code", "")

    if _s3_enabled():
        # Delete session JSON
        try:
            _s3().delete_object(Bucket=S3_BUCKET, Key=_s3_session_key(session_id))
        except Exception as exc:
            log.error("S3 delete session error: %s", exc)
        # Delete all artifacts for this session
        try:
            paginator = _s3().get_paginator("list_objects_v2")
            for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=f"{S3_PREFIX}/artifacts/{session_id}_"):
                keys = [{"Key": o["Key"]} for o in page.get("Contents", [])]
                if keys:
                    _s3().delete_objects(Bucket=S3_BUCKET, Delete={"Objects": keys})
        except Exception as exc:
            log.error("S3 delete artifacts error: %s", exc)
    else:
        # Local filesystem
        p = session_path(session_id)
        if p.exists():
            p.unlink()
        for f in list(ARTIFACTS_DIR.glob(f"{session_id}_*.pptx")) + \
                  list(ARTIFACTS_DIR.glob(f"{session_id}_*.docx")):
            f.unlink(missing_ok=True)

    # Remove from in-memory room code map
    if room_code and room_code in _room_code_map:
        del _room_code_map[room_code]

    log.info("Deleted session %s and its artifacts", session_id)
    return jsonify({"deleted": session_id})


@app.route("/api/sessions/<session_id>/advance", methods=["POST"])
def advance_exercise(session_id):
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    idx = s["current_exercise_index"]
    max_idx = len(s["exercise_order"]) - 1
    if idx >= max_idx:
        return jsonify({"error": "Already at last exercise"}), 400

    s["current_exercise_index"] = idx + 1
    # Mark previous exercise complete
    prev_ex_id = s["exercise_order"][idx]
    s["exercises"][prev_ex_id]["status"] = "completed"
    # Mark new exercise active
    new_ex_id = s["exercise_order"][idx + 1]
    s["exercises"][new_ex_id]["status"] = "active"

    save_session(s)

    ex_meta = get_exercise_meta(new_ex_id) or {}
    socketio.emit("exercise_changed", {
        "exercise_id": new_ex_id,
        "index": idx + 1,
        "exercise_meta": ex_meta,
    }, room=session_id)

    log.info("Session %s advanced to exercise %d (%s)", session_id, idx + 1, new_ex_id)
    return jsonify({"current_exercise_index": idx + 1, "exercise_id": new_ex_id})


@app.route("/api/sessions/<session_id>/back", methods=["POST"])
def back_exercise(session_id):
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    idx = s["current_exercise_index"]
    if idx <= 0:
        return jsonify({"error": "Already at first exercise"}), 400

    s["current_exercise_index"] = idx - 1
    prev_ex_id = s["exercise_order"][idx]
    s["exercises"][prev_ex_id]["status"] = "pending"
    new_ex_id  = s["exercise_order"][idx - 1]
    s["exercises"][new_ex_id]["status"] = "active"

    save_session(s)

    ex_meta = get_exercise_meta(new_ex_id) or {}
    socketio.emit("exercise_changed", {
        "exercise_id": new_ex_id,
        "index": idx - 1,
        "exercise_meta": ex_meta,
    }, room=session_id)

    return jsonify({"current_exercise_index": idx - 1, "exercise_id": new_ex_id})


@app.route("/api/sessions/<session_id>/exercises", methods=["PUT"])
def update_exercises(session_id):
    """Replace the exercise order for a session, preserving existing messages/summaries."""
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    body = request.json or {}
    new_order = body.get("exercise_order") or []
    if not new_order:
        return jsonify({"error": "exercise_order must not be empty"}), 400

    # Preserve existing exercise state; initialise any new exercises
    existing = s.get("exercises", {})
    new_exercises: dict = {}
    for ex_id in new_order:
        new_exercises[ex_id] = existing.get(ex_id, {
            "messages": [], "summary": None, "summary_generated_at": None, "status": "pending",
        })

    # Clamp current_exercise_index to valid range
    new_idx = min(s.get("current_exercise_index", 0), len(new_order) - 1)

    s["exercise_order"] = new_order
    s["exercises"] = new_exercises
    s["current_exercise_index"] = new_idx
    save_session(s)

    # Notify all connected clients
    safe = {k: v for k, v in s.items() if k != "passcode"}
    socketio.emit("session_state", safe, room=session_id)

    log.info("Updated exercise list for session %s (%d exercises)", session_id, len(new_order))
    return jsonify({"exercise_order": new_order, "current_exercise_index": new_idx})


# ── Artifact generation ─────────────────────────────────────────────────────────

def _parse_markdown_sections(md: str) -> list[dict]:
    """Split markdown into sections: [{title, bullets}].
    Treats ## headings as section titles and collects body lines as bullets."""
    sections: list[dict] = []
    current: dict | None = None
    for line in md.splitlines():
        line = line.strip()
        if line.startswith("## "):
            if current:
                sections.append(current)
            current = {"title": line[3:].strip(), "bullets": []}
        elif line.startswith("# "):
            if current:
                sections.append(current)
            current = {"title": line[2:].strip(), "bullets": []}
        elif current is not None and line:
            # strip leading markdown list markers
            text = line.lstrip("-*• ").strip()
            if text:
                current["bullets"].append(text)
    if current:
        sections.append(current)
    return sections


def _safe_filename(text: str) -> str:
    """Convert text to a filesystem-safe slug."""
    import re
    return re.sub(r"[^a-zA-Z0-9_-]", "_", text)[:40]


def _generate_docx(session: dict, ex_id: str, ex_meta: dict, summary_text: str) -> Path:
    """Create a .docx summary and return its path."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    sections = _parse_markdown_sections(summary_text)
    if sections:
        for sec in sections:
            doc.add_heading(sec["title"], 2)
            for bullet in sec["bullets"]:
                doc.add_paragraph(bullet, style="List Bullet")
    else:
        # fallback: raw text
        for line in summary_text.splitlines():
            if line.strip():
                doc.add_paragraph(line.strip())

    slug = _safe_filename(ex_meta.get("name", ex_id))
    filename = f"{session['session_id']}_{slug}.docx"
    if _s3_enabled():
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        key = f"{S3_PREFIX}/artifacts/{filename}"
        _s3().put_object(
            Bucket=S3_BUCKET, Key=key, Body=buf.read(),
            ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        log.info("Uploaded DOCX to S3: %s", key)
        return Path(filename)
    out_path = ARTIFACTS_DIR / filename
    doc.save(str(out_path))
    return out_path


def _generate_pptx(session: dict, ex_id: str, ex_meta: dict, summary_text: str) -> Path:
    """Create a .pptx summary styled after Cloud_DTW_Updated.pptx.

    Layout strategy by exercise type:
    - empathy_maps:        left context panel + 4-quadrant sticky grid (SAYS/THINKS/FEELS/DOES)
    - asis/tobe_process_map: left sidebar + horizontal stage columns with sticky notes
    - big_idea_vignettes:  left sidebar with tag legend + card grid
    - everything else:     white bg, dark-red header bar, card columns (1–3 per slide)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Slide dimensions: 13.3333" × 7.5" (matches Cloud_DTW_Updated.pptx) ──────
    prs = Presentation()
    prs.slide_width  = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    blank = prs.slide_layouts[6]

    # ── Color palette ─────────────────────────────────────────────────────────────
    DARK_RED   = RGBColor(0x8B, 0x00, 0x00)  # header bars
    LILLY_RED  = RGBColor(0xD5, 0x2B, 0x1E)  # top/bottom master bar
    DARK_BURG  = RGBColor(0x51, 0x12, 0x07)  # sidebar fill
    ACCENT_RED = RGBColor(0xE1, 0x25, 0x1B)  # sidebar edge, tags
    NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)  # primary text
    DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)  # body text
    MID_GRAY   = RGBColor(0x55, 0x55, 0x55)  # secondary text
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    CARD_BODY  = RGBColor(0xFD, 0xF5, 0xF5)  # card body (very light pink)
    PALE_PINK  = RGBColor(0xFB, 0xCF, 0xC8)  # sidebar secondary text
    LT_GRAY    = RGBColor(0xDD, 0xDD, 0xDD)  # dividers on dark
    SHADOW     = RGBColor(0xCC, 0xBF, 0xB0)  # card shadow

    # Empathy map quadrant colors (from slides 8–15)
    EM_YELLOW = RGBColor(0xFF, 0xE5, 0x7A)   # SAYS / THINKS
    EM_PINK   = RGBColor(0xFF, 0xB3, 0xB3)   # FEELS
    EM_GREEN  = RGBColor(0xB7, 0xE1, 0xA1)   # DOES

    # Big-ideas tag palette: (tag_fill, card_header_bg)
    TAG_PALETTE = [
        (RGBColor(0xE1, 0x25, 0x1B), RGBColor(0xFF, 0xF8, 0xF8)),
        (RGBColor(0x0F, 0x3A, 0x85), RGBColor(0xEB, 0xF4, 0xFF)),
        (RGBColor(0x14, 0x4B, 0x2D), RGBColor(0xE8, 0xF0, 0xEB)),
        (RGBColor(0xFF, 0xC7, 0x09), RGBColor(0xFF, 0xF8, 0xE1)),
        (RGBColor(0xF5, 0x8E, 0x7D), RGBColor(0xFC, 0xEE, 0xE8)),
        (RGBColor(0x99, 0xBF, 0xE5), RGBColor(0xEB, 0xF4, 0xFF)),
    ]

    # ── Low-level helpers ─────────────────────────────────────────────────────────
    def _rect(slide, x, y, w, h, fill_rgb):
        shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb
        shp.line.fill.background()
        return shp

    def _tbx(slide, x, y, w, h):
        return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

    def _set_text(tf, text, size_pt, bold=False, italic=False, color=None,
                  align=PP_ALIGN.LEFT, name="Arial", first=False, space_after=0):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = text
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if p.runs:
            r = p.runs[0]
            r.font.size  = Pt(size_pt)
            r.font.bold  = bold
            r.font.italic = italic
            r.font.name  = name
            if color:
                r.font.color.rgb = color
        return p

    # ── Master decorations: thin red top line + red bottom bar ────────────────────
    def _decorate(slide):
        _rect(slide, 0, 0,        13.3333, 0.05, LILLY_RED)   # top line
        _rect(slide, 0, 7.1944,   13.3333, 0.30, LILLY_RED)   # bottom bar

    # ── Dark-red full-width header bar ────────────────────────────────────────────
    def _header(slide, title_text: str):
        bar = _rect(slide, 0, 0.05, 13.3333, 0.75, DARK_RED)
        tf = bar.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.4)
        tf.margin_top  = Inches(0.0)
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.space_before = Pt(11)
        if p.runs:
            r = p.runs[0]
            r.font.size  = Pt(16)
            r.font.bold  = True
            r.font.name  = "Arial Black"
            r.font.color.rgb = WHITE

    # ── Left sidebar (dark-burgundy panel) ───────────────────────────────────────
    def _sidebar(slide, section_label: str, title: str, body_lines: list[str],
                 width: float = 2.82):
        _rect(slide, 0,     0, width, 7.5, DARK_BURG)
        _rect(slide, 0,     0, 0.09,  7.5, ACCENT_RED)
        lbl = _tbx(slide, 0.18, 0.20, width - 0.28, 0.38)
        lbl.text_frame.word_wrap = True
        _set_text(lbl.text_frame, section_label.upper(), 8, bold=True,
                  color=ACCENT_RED, name="Arial", first=True)
        ttl = _tbx(slide, 0.18, 0.62, width - 0.28, 0.85)
        ttl.text_frame.word_wrap = True
        _set_text(ttl.text_frame, title, 15, bold=True,
                  color=WHITE, name="Arial Black", first=True)
        _rect(slide, 0.26, 1.55, width - 0.44, 0.025, ACCENT_RED)
        if body_lines:
            bx = _tbx(slide, 0.18, 1.62, width - 0.28, 5.4)
            bx.text_frame.word_wrap = True
            for i, ln in enumerate(body_lines[:14]):
                _set_text(bx.text_frame, ln, 8, italic=True,
                          color=PALE_PINK, name="Calibri", first=(i == 0),
                          space_after=2)

    # ── Sticky note card ──────────────────────────────────────────────────────────
    def _sticky(slide, text: str, x: float, y: float, color: RGBColor,
                w: float = 3.5, h: float = 0.634):
        _rect(slide, x + 0.025, y + 0.025, w, h, SHADOW)  # shadow
        card = _rect(slide, x, y, w, h, color)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_right  = Inches(0.08)
        tf.margin_top  = Inches(0.04); tf.margin_bottom = Inches(0.04)
        n = len(text)
        pt = 9 if n <= 80 else (8 if n <= 130 else 7)
        if n > 160:
            text = text[:157] + "…"
        p = tf.paragraphs[0]; p.text = text
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(pt); r.font.name = "Calibri"
            r.font.color.rgb = NEAR_BLACK

    # ─────────────────────────────────────────────────────────────────────────────
    # EMPATHY MAP  (slides 8–15 style from template)
    # ─────────────────────────────────────────────────────────────────────────────
    def _empathy_slide(persona_section: dict):
        slide = prs.slides.add_slide(blank)
        _decorate(slide)

        QUADRANT_KEYS = {"SAYS", "THINKS", "FEELS", "DOES"}
        LEFT_KEYS     = {"CONTEXT", "PAINS", "GAINS"}
        left_buckets:     dict[str, list[str]] = {k: [] for k in LEFT_KEYS}
        quadrant_bullets: dict[str, list[str]] = {q: [] for q in QUADRANT_KEYS}

        for b in persona_section["bullets"]:
            placed = False
            # Check all known keys (left panel + quadrant)
            for key in list(QUADRANT_KEYS) + list(LEFT_KEYS):
                if (b.upper().startswith(key + ":")
                        or b.upper().startswith(f"**{key}:**")
                        or b.upper().startswith(f"**{key}:**".upper())):
                    clean = b.split(":", 1)[-1].strip().lstrip("* ").strip()
                    if key in QUADRANT_KEYS:
                        quadrant_bullets[key].append(clean)
                    else:
                        left_buckets[key].append(clean)
                    placed = True; break
            if not placed:
                left_buckets["CONTEXT"].append(b)

        # ── Left panel ────────────────────────────────────────────────────────────
        # "Empathy Map" label
        lbl = _tbx(slide, 0.28, 0.10, 5.2, 0.38)
        _set_text(lbl.text_frame, "EMPATHY MAP", 10, bold=True,
                  color=DARK_RED, name="Arial", first=True)

        # Persona name
        name_box = _tbx(slide, 0.28, 0.52, 5.2, 0.55)
        name_box.text_frame.word_wrap = True
        _set_text(name_box.text_frame, persona_section["title"], 20, bold=True,
                  color=NEAR_BLACK, name="Arial Black", first=True)

        # Thin red divider under name
        _rect(slide, 0.28, 1.12, 4.9, 0.025, DARK_RED)

        # Context description
        ctx_y = 1.18
        if left_buckets["CONTEXT"]:
            ctx = _tbx(slide, 0.28, ctx_y, 5.0, 1.10)
            ctx.text_frame.word_wrap = True
            for i, cb in enumerate(left_buckets["CONTEXT"][:2]):
                _set_text(ctx.text_frame, cb, 10, italic=True,
                          color=MID_GRAY, name="Calibri",
                          first=(i == 0), space_after=3)
            ctx_y += 1.18

        # PAINS section
        if left_buckets["PAINS"]:
            _rect(slide, 0.28, ctx_y, 1.0, 0.24, RGBColor(0x8B, 0x00, 0x00))
            pl = _tbx(slide, 0.28, ctx_y, 1.02, 0.24)
            p = pl.text_frame.paragraphs[0]; p.text = "PAINS"
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                r = p.runs[0]; r.font.size = Pt(8); r.font.bold = True
                r.font.name = "Arial"; r.font.color.rgb = WHITE
            pain_box = _tbx(slide, 0.28, ctx_y + 0.28, 5.0, 1.80)
            pain_box.text_frame.word_wrap = True
            for i, pain in enumerate(left_buckets["PAINS"][:4]):
                _set_text(pain_box.text_frame, f"• {pain}", 10,
                          color=DARK_GRAY, name="Calibri",
                          first=(i == 0), space_after=4)
            ctx_y += 0.28 + min(len(left_buckets["PAINS"]), 4) * 0.42 + 0.15

        # GAINS section
        if left_buckets["GAINS"] and ctx_y < 6.5:
            GAINS_COLOR = RGBColor(0x14, 0x4B, 0x2D)
            _rect(slide, 0.28, ctx_y, 1.05, 0.24, GAINS_COLOR)
            gl = _tbx(slide, 0.28, ctx_y, 1.07, 0.24)
            p = gl.text_frame.paragraphs[0]; p.text = "GAINS"
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                r = p.runs[0]; r.font.size = Pt(8); r.font.bold = True
                r.font.name = "Arial"; r.font.color.rgb = WHITE
            gain_box = _tbx(slide, 0.28, ctx_y + 0.28, 5.0, 7.10 - ctx_y - 0.30)
            gain_box.text_frame.word_wrap = True
            for i, gain in enumerate(left_buckets["GAINS"][:4]):
                _set_text(gain_box.text_frame, f"• {gain}", 10,
                          color=DARK_GRAY, name="Calibri",
                          first=(i == 0), space_after=4)

        # ── Right: 4-quadrant grid ────────────────────────────────────────────────
        _rect(slide, 5.58, 0.70, 7.40, 6.45, WHITE)
        _rect(slide, 9.27, 0.70, 0.02, 6.45, NEAR_BLACK)   # vertical divider
        _rect(slide, 5.58, 3.92, 7.40, 0.02, NEAR_BLACK)   # horizontal divider

        # Quadrant label bars
        QLABELS = [
            ("SAYS",   5.58, 0.70, RGBColor(0xFF, 0xE5, 0x7A)),
            ("THINKS", 9.29, 0.70, RGBColor(0xFF, 0xE5, 0x7A)),
            ("FEELS",  5.58, 3.92, EM_PINK),
            ("DOES",   9.29, 3.92, EM_GREEN),
        ]
        for label, qx, qy, lcolor in QLABELS:
            bar = _rect(slide, qx, qy, 3.67, 0.30, lcolor)
            tf = bar.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.0)
            p = tf.paragraphs[0]; p.text = label; p.space_before = Pt(4)
            if p.runs:
                r = p.runs[0]; r.font.size = Pt(11)
                r.font.bold = True; r.font.name = "Arial"
                r.font.color.rgb = NEAR_BLACK

        # Sticky cards (3 per quadrant)
        CARD_W = 3.50
        CARD_H = 0.87
        ROW_GAP = 0.07
        QLAYOUT = [
            ("SAYS",   5.62, [1.04, 1.98, 2.92], EM_YELLOW),
            ("THINKS", 9.31, [1.04, 1.98, 2.92], EM_YELLOW),
            ("FEELS",  5.62, [4.24, 5.18, 6.12], EM_PINK),
            ("DOES",   9.31, [4.24, 5.18, 6.12], EM_GREEN),
        ]
        for quad, qx, y_rows, color in QLAYOUT:
            items = quadrant_bullets[quad]
            for i, ypos in enumerate(y_rows):
                if i < len(items):
                    _sticky(slide, items[i], qx, ypos, color,
                            w=CARD_W, h=CARD_H)

    # ─────────────────────────────────────────────────────────────────────────────
    # PROCESS MAP  (left sidebar + horizontal stage columns)
    # ─────────────────────────────────────────────────────────────────────────────
    def _process_slide(sections: list[dict]):
        slide = prs.slides.add_slide(blank)
        _decorate(slide)

        SB_W       = 1.80
        CONT_X     = SB_W + 0.08
        CONT_W     = 13.3333 - CONT_X - 0.10
        BAND_Y     = 0.05
        BAND_H     = 0.60
        STAGE_H    = 0.48
        CONTENT_Y  = BAND_Y + BAND_H + STAGE_H + 0.08
        AVAIL_H    = 7.1944 - CONTENT_Y - 0.10

        # Sidebar
        _rect(slide, 0, 0, SB_W, 7.1944, DARK_BURG)
        _rect(slide, 0, 0, 0.09, 7.1944, ACCENT_RED)
        lbl = _tbx(slide, 0.15, 0.22, SB_W - 0.22, 0.32)
        _set_text(lbl.text_frame, ex_meta.get("phase_name", "PROCESS MAP").upper(),
                  7, bold=True, color=ACCENT_RED, name="Arial", first=True)
        ttl = _tbx(slide, 0.15, 0.58, SB_W - 0.22, 0.75)
        ttl.text_frame.word_wrap = True
        _set_text(ttl.text_frame, ex_meta.get("name", ex_id), 12, bold=True,
                  color=WHITE, name="Arial Black", first=True)

        # Stage header band
        _rect(slide, CONT_X, BAND_Y + BAND_H, CONT_W, STAGE_H,
              RGBColor(0x22, 0x22, 0x22))

        n_stages = max(len(sections), 1)
        stage_w  = CONT_W / n_stages
        CIRCLE_D = 0.24
        NOTES_PER_STAGE = max(1, int(AVAIL_H / 0.85))

        for i, sec in enumerate(sections):
            sx = CONT_X + i * stage_w
            if i > 0:
                _rect(slide, sx, BAND_Y + BAND_H, 0.015,
                      7.1944 - BAND_Y - BAND_H,
                      RGBColor(0x44, 0x44, 0x44))
            # Stage circle
            cx = sx + stage_w / 2 - CIRCLE_D / 2
            cy = BAND_Y + BAND_H + (STAGE_H - CIRCLE_D) / 2
            circ = slide.shapes.add_shape(
                9, Inches(cx), Inches(cy), Inches(CIRCLE_D), Inches(CIRCLE_D))
            circ.fill.solid(); circ.fill.fore_color.rgb = DARK_RED
            circ.line.fill.background()
            nb = _tbx(slide, cx + 0.015, cy + 0.02, CIRCLE_D - 0.03, CIRCLE_D - 0.04)
            _set_text(nb.text_frame, str(i + 1), 7, bold=True,
                      color=WHITE, name="Arial", align=PP_ALIGN.CENTER, first=True)
            # Stage label
            lbl2 = _tbx(slide, sx + 0.04, BAND_Y + BAND_H + STAGE_H - 0.04,
                        stage_w - 0.08, 0.42)
            lbl2.text_frame.word_wrap = True
            _set_text(lbl2.text_frame, sec["title"], 7,
                      color=LT_GRAY, name="Arial",
                      align=PP_ALIGN.CENTER, first=True)
            # Sticky notes
            STICKY_W  = stage_w - 0.14
            STICKY_H  = min(0.82, AVAIL_H / max(NOTES_PER_STAGE, 1) - 0.08)
            NOTE_COLORS = [EM_YELLOW, RGBColor(0xFF, 0xCC, 0xA8),
                           RGBColor(0xA8, 0xD8, 0xFF), EM_GREEN]
            for j, bullet in enumerate(sec["bullets"][:NOTES_PER_STAGE]):
                bx = sx + 0.07
                by = CONTENT_Y + j * (STICKY_H + 0.08)
                if by + STICKY_H > 7.1944:
                    break
                _sticky(slide, bullet, bx, by, NOTE_COLORS[j % 4],
                        w=STICKY_W, h=STICKY_H)

    # ─────────────────────────────────────────────────────────────────────────────
    # BIG IDEAS  (left sidebar + card grid with colored tags)
    # ─────────────────────────────────────────────────────────────────────────────
    def _big_ideas_slide(sections: list[dict]):
        slide = prs.slides.add_slide(blank)
        _decorate(slide)

        SB_W      = 3.00
        TAG_W     = 0.90; TAG_H  = 0.185
        HDR_H     = 0.72; BODY_H = 1.10
        CW        = 2.25; GAP    = 0.12
        COLS      = 4
        START_X   = SB_W + 0.22
        START_Y   = 0.42
        ROW_H     = HDR_H + BODY_H + 0.22

        # Sidebar
        _rect(slide, 0, 0, SB_W, 7.5, DARK_BURG)
        _rect(slide, 0, 0, 0.09, 7.5, ACCENT_RED)
        lbl_s = _tbx(slide, 0.20, 0.22, SB_W - 0.30, 0.32)
        _set_text(lbl_s.text_frame, "BIG IDEAS", 9, bold=True,
                  color=ACCENT_RED, name="Calibri", first=True)
        ttl_s = _tbx(slide, 0.20, 0.57, SB_W - 0.30, 0.80)
        ttl_s.text_frame.word_wrap = True
        _set_text(ttl_s.text_frame, ex_meta.get("name", ex_id), 17, bold=True,
                  color=WHITE, name="Arial Black", first=True)
        _rect(slide, 0.28, 1.45, SB_W - 0.52, 0.025, ACCENT_RED)

        # Theme list in sidebar
        sy = 1.53
        for i, sec in enumerate(sections):
            if sy > 6.6: break
            tc, _ = TAG_PALETTE[i % len(TAG_PALETTE)]
            _rect(slide, 0.22, sy + 0.04, 0.13, 0.13, tc)
            tb = _tbx(slide, 0.42, sy, SB_W - 0.55, 0.22)
            _set_text(tb.text_frame, sec["title"][:48], 8, bold=True,
                      color=WHITE, name="Calibri", first=True)
            sy += 0.27

        # Card grid
        all_cards: list[tuple] = []
        for si, sec in enumerate(sections):
            tc, hbg = TAG_PALETTE[si % len(TAG_PALETTE)]
            for b in sec["bullets"][:3]:
                all_cards.append((sec["title"], b, tc, hbg))

        for ci, (theme, bullet, tag_color, hdr_bg) in enumerate(all_cards[:12]):
            col = ci % COLS
            row = ci // COLS
            cx_ = START_X + col * (CW + GAP)
            cy_ = START_Y + row * ROW_H
            if cy_ + HDR_H + BODY_H > 7.5:
                break
            _rect(slide, cx_ + 0.025, cy_ + 0.025, CW, HDR_H + BODY_H, SHADOW)
            _rect(slide, cx_, cy_, CW, HDR_H, hdr_bg)
            tag = _rect(slide, cx_ + 0.07, cy_ + 0.07, TAG_W, TAG_H, tag_color)
            p = tag.text_frame.paragraphs[0]
            p.text = theme[:22].upper(); p.alignment = PP_ALIGN.CENTER
            if p.runs:
                r = p.runs[0]; r.font.size = Pt(6)
                r.font.bold = True; r.font.name = "Arial"
                r.font.color.rgb = WHITE
            title_box = _tbx(slide, cx_ + 0.07, cy_ + 0.31, CW - 0.14, 0.40)
            title_box.text_frame.word_wrap = True
            t_text = bullet[:65] if len(bullet) > 65 else bullet
            _set_text(title_box.text_frame, t_text, 9, bold=True,
                      color=NEAR_BLACK, name="Calibri", first=True)
            _rect(slide, cx_, cy_ + HDR_H, CW, BODY_H, WHITE)
            _rect(slide, cx_, cy_ + HDR_H, CW, 0.015, RGBColor(0xCC, 0xCC, 0xCC))
            rest = bullet[65:]
            if rest:
                body_box = _tbx(slide, cx_ + 0.07, cy_ + HDR_H + 0.04,
                                CW - 0.14, BODY_H - 0.08)
                body_box.text_frame.word_wrap = True
                _set_text(body_box.text_frame, rest[:120], 8,
                          color=DARK_GRAY, name="Calibri", first=True)

    # ─────────────────────────────────────────────────────────────────────────────
    # COLUMN-CARD slide  (white bg, dark-red header, N card columns)
    # ─────────────────────────────────────────────────────────────────────────────
    def _column_slide(slide_title: str, sections: list[dict], subtitle: str = ""):
        slide = prs.slides.add_slide(blank)
        _decorate(slide)
        _header(slide, slide_title)

        if subtitle:
            st = _tbx(slide, 0.45, 0.84, 12.4, 0.28)
            _set_text(st.text_frame, subtitle, 10, italic=True,
                      color=MID_GRAY, name="Arial", first=True)

        TOP_Y   = 1.14
        AVAIL_H = 7.1944 - TOP_Y - 0.12
        AVAIL_W = 12.43
        LEFT_X  = 0.45
        n       = max(len(sections), 1)
        GAP     = 0.18
        COL_W   = (AVAIL_W - GAP * (n - 1)) / n
        font_pt = 9 if n >= 3 else (10 if n == 2 else 12)

        for i, sec in enumerate(sections):
            cx_ = LEFT_X + i * (COL_W + GAP)
            hdr = _rect(slide, cx_, TOP_Y, COL_W, 0.36, DARK_RED)
            hdr_tf = hdr.text_frame
            hdr_tf.margin_left = Inches(0.12); hdr_tf.margin_top = Inches(0.0)
            p = hdr_tf.paragraphs[0]
            p.text = sec["title"].upper(); p.space_before = Pt(8)
            if p.runs:
                r = p.runs[0]; r.font.size = Pt(9)
                r.font.bold = True; r.font.name = "Arial"
                r.font.color.rgb = WHITE
            body_h = AVAIL_H - 0.36
            _rect(slide, cx_, TOP_Y + 0.36, COL_W, body_h, CARD_BODY)
            txt = _tbx(slide, cx_ + 0.12, TOP_Y + 0.44, COL_W - 0.20, body_h - 0.15)
            txt.text_frame.word_wrap = True
            for bi, bullet in enumerate(sec["bullets"][:16]):
                _set_text(txt.text_frame, f"• {bullet}", font_pt,
                          color=DARK_GRAY, name="Arial",
                          first=(bi == 0), space_after=3)

    # ─────────────────────────────────────────────────────────────────────────────
    # BUILD SLIDES
    # ─────────────────────────────────────────────────────────────────────────────
    EMPATHY_EXERCISES   = {"empathy_maps"}
    PROCESS_EXERCISES   = {"asis_process_map", "tobe_process_map"}
    BIG_IDEAS_EXERCISES = {"big_idea_vignettes"}

    sections = _parse_markdown_sections(summary_text)
    if not sections:
        sections = [{"title": ex_meta.get("name", ex_id),
                     "bullets": [l.strip() for l in summary_text.splitlines() if l.strip()]}]

    if ex_id in EMPATHY_EXERCISES:
        for sec in sections:
            if sec["title"].lower().startswith("common thread"):
                _column_slide(sec["title"], [sec])
            else:
                _empathy_slide(sec)

    elif ex_id in PROCESS_EXERCISES:
        STAGES_PER_SLIDE = 6
        for chunk_start in range(0, len(sections), STAGES_PER_SLIDE):
            _process_slide(sections[chunk_start:chunk_start + STAGES_PER_SLIDE])

    elif ex_id in BIG_IDEAS_EXERCISES:
        THEMES_PER_SLIDE = 6
        for chunk_start in range(0, len(sections), THEMES_PER_SLIDE):
            _big_ideas_slide(sections[chunk_start:chunk_start + THEMES_PER_SLIDE])

    else:
        MAX_COLS = 3
        if len(sections) <= MAX_COLS:
            _column_slide(ex_meta.get("name", ex_id), sections)
        else:
            for chunk_start in range(0, len(sections), MAX_COLS):
                chunk = sections[chunk_start:chunk_start + MAX_COLS]
                slide_title = (ex_meta.get("name", ex_id)
                               if chunk_start == 0
                               else f"{ex_meta.get('name', ex_id)} (cont.)")
                _column_slide(slide_title, chunk)

    slug = _safe_filename(ex_meta.get("name", ex_id))
    filename = f"{session['session_id']}_{slug}.pptx"
    if _s3_enabled():
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        key = f"{S3_PREFIX}/artifacts/{filename}"
        _s3().put_object(
            Bucket=S3_BUCKET, Key=key, Body=buf.read(),
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        log.info("Uploaded PPTX to S3: %s", key)
        return Path(filename)
    out_path = ARTIFACTS_DIR / filename
    prs.save(str(out_path))
    return out_path


def _generate_artifacts(session: dict, ex_id: str, ex_meta: dict, summary_text: str) -> list[str]:
    """Generate PPTX + DOCX for a summary. Returns list of filenames created."""
    created = []
    try:
        p = _generate_pptx(session, ex_id, ex_meta, summary_text)
        created.append(p.name)
        log.info("Generated PPTX: %s", p.name)
    except Exception as e:
        log.error("PPTX generation failed: %s", e)
    try:
        d = _generate_docx(session, ex_id, ex_meta, summary_text)
        created.append(d.name)
        log.info("Generated DOCX: %s", d.name)
    except Exception as e:
        log.error("DOCX generation failed: %s", e)
    return created


@app.route("/api/sessions/<session_id>/summary", methods=["POST"])
def generate_summary(session_id):
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    idx    = s["current_exercise_index"]
    ex_id  = s["exercise_order"][idx]
    ex_state = s["exercises"].get(ex_id, {})
    messages = ex_state.get("messages", [])

    if not messages:
        return jsonify({"error": "No messages in current exercise to summarize"}), 400

    ex_meta = get_exercise_meta(ex_id)
    if not ex_meta:
        return jsonify({"error": f"Exercise metadata not found for '{ex_id}'"}), 500

    try:
        summary_text = generate_exercise_summary(ex_meta, messages)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        log.error("Summary generation failed: %s", e)
        return jsonify({"error": f"AI summary failed: {e}"}), 500

    s["exercises"][ex_id]["summary"] = summary_text
    s["exercises"][ex_id]["summary_generated_at"] = utcnow()
    save_session(s)

    # Generate PPTX + DOCX artifacts in background
    ex_meta_full = get_exercise_meta(ex_id) or {}
    _generate_artifacts(s, ex_id, ex_meta_full, summary_text)

    socketio.emit("summary_ready", {
        "exercise_id": ex_id,
        "summary_text": summary_text,
    }, room=session_id)

    log.info("Summary generated for session %s exercise %s", session_id, ex_id)
    return jsonify({"exercise_id": ex_id, "summary_text": summary_text})


@app.route("/api/sessions/<session_id>/artifacts")
def list_artifacts(session_id):
    """List all PPTX/DOCX artifacts for a session."""
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    result = []
    prefix_str = session_id  # used for slug parsing

    if _s3_enabled():
        s3_prefix = f"{S3_PREFIX}/artifacts/{session_id}_"
        try:
            paginator = _s3().get_paginator("list_objects_v2")
            for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=s3_prefix):
                for obj in page.get("Contents", []):
                    fname = obj["Key"].split("/")[-1]
                    ext   = fname.rsplit(".", 1)[-1].upper()
                    stem  = fname.rsplit(".", 1)[0]
                    slug  = stem[len(prefix_str) + 1:]
                    label = slug.replace("_", " ").title()
                    result.append({
                        "filename": fname,
                        "label":    label,
                        "type":     ext,
                        "size":     obj["Size"],
                        "url":      f"/api/sessions/{session_id}/artifacts/{fname}",
                    })
        except Exception as exc:
            log.error("list_artifacts S3 error: %s", exc)
    else:
        files = sorted(ARTIFACTS_DIR.glob(f"{session_id}_*.pptx")) + \
                sorted(ARTIFACTS_DIR.glob(f"{session_id}_*.docx"))
        for f in files:
            ext   = f.suffix.upper().lstrip(".")
            slug  = f.stem[len(prefix_str) + 1:]
            label = slug.replace("_", " ").title()
            result.append({
                "filename": f.name,
                "label":    label,
                "type":     ext,
                "size":     f.stat().st_size,
                "url":      f"/api/sessions/{session_id}/artifacts/{f.name}",
            })

    # Sort: PPTX before DOCX, then alphabetically by label
    result.sort(key=lambda a: (a["label"], a["type"]))
    return jsonify(result)


@app.route("/api/sessions/<session_id>/artifacts/<filename>")
def download_artifact(session_id, filename):
    """Download a specific artifact — proxies from S3 or serves from disk."""
    if not filename.startswith(session_id) or ".." in filename:
        return jsonify({"error": "Not found"}), 404
    if _s3_enabled():
        key = f"{S3_PREFIX}/artifacts/{filename}"
        try:
            obj  = _s3().get_object(Bucket=S3_BUCKET, Key=key)
            data = obj["Body"].read()
            ext  = filename.rsplit(".", 1)[-1].lower()
            ct   = {
                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            }.get(ext, "application/octet-stream")
            return Response(
                data,
                headers={
                    "Content-Disposition": f'attachment; filename="{filename}"',
                    "Content-Type": ct,
                },
            )
        except Exception as exc:
            log.error("S3 download error: %s", exc)
            return jsonify({"error": "File not found"}), 404
    return send_from_directory(str(ARTIFACTS_DIR), filename, as_attachment=True)


@app.route("/api/sessions/<session_id>/artifacts/<filename>", methods=["DELETE"])
def delete_artifact(session_id, filename):
    """Delete a specific artifact file."""
    if not filename.startswith(session_id) or ".." in filename:
        return jsonify({"error": "Not found"}), 404
    if _s3_enabled():
        key = f"{S3_PREFIX}/artifacts/{filename}"
        try:
            _s3().delete_object(Bucket=S3_BUCKET, Key=key)
            log.info("Deleted S3 artifact: %s", key)
            return jsonify({"deleted": filename})
        except Exception as exc:
            log.error("S3 delete error: %s", exc)
            return jsonify({"error": "Delete failed"}), 500
    path = ARTIFACTS_DIR / filename
    if not path.exists():
        return jsonify({"error": "File not found"}), 404
    path.unlink()
    log.info("Deleted artifact: %s", filename)
    return jsonify({"deleted": filename})


@app.route("/api/sessions/<session_id>/goto", methods=["POST"])
def goto_exercise(session_id):
    """Jump the session to any exercise by index."""
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    body = request.json or {}
    try:
        target = int(body["index"])
    except (KeyError, TypeError, ValueError):
        return jsonify({"error": "index required"}), 400

    max_idx = len(s["exercise_order"]) - 1
    if not (0 <= target <= max_idx):
        return jsonify({"error": f"index must be 0–{max_idx}"}), 400

    old_idx = s["current_exercise_index"]
    if old_idx != target:
        s["exercises"][s["exercise_order"][old_idx]]["status"] = "completed" if target > old_idx else "pending"
        s["current_exercise_index"] = target
        s["exercises"][s["exercise_order"][target]]["status"] = "active"
        save_session(s)

    new_ex_id = s["exercise_order"][target]
    ex_meta = get_exercise_meta(new_ex_id) or {}
    socketio.emit("exercise_changed", {
        "exercise_id": new_ex_id,
        "index": target,
        "exercise_meta": ex_meta,
    }, room=session_id)

    log.info("Session %s jumped to exercise %d (%s)", session_id, target, new_ex_id)
    return jsonify({"current_exercise_index": target, "exercise_id": new_ex_id})


@app.route("/api/config")
def get_config():
    """Return public Azure AD config for frontend MSAL.js auth."""
    return jsonify({
        "azure_client_id": os.environ.get("AZURE_CLIENT_ID", ""),
        "azure_tenant_id": os.environ.get("AZURE_TENANT_ID", ""),
    })


@app.route("/api/me")
def get_me():
    """Return current user identity decoded from the keychain JWT."""
    identity = _identity_from_keychain()
    return jsonify(identity)


@app.route("/api/facilitator-auth", methods=["POST"])
def facilitator_auth():
    """Authenticate a facilitator by name (fallback when keychain unavailable)."""
    body = request.json or {}
    name = (body.get("name") or "").strip()

    if name.lower() not in [f.lower() for f in FACILITATORS]:
        return jsonify({"error": "Name not on facilitator list"}), 403

    return jsonify({"name": name, "is_facilitator": True})


@app.route("/api/sessions/<session_id>/export", methods=["GET"])
def export_session(session_id):
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404
    safe_name = s["session_name"].encode("ascii", "ignore").decode().replace(" ", "_")[:40].strip("_") or session_id[:8]
    filename  = f"DTW_{safe_name}_{session_id[:8]}.json"
    resp = app.response_class(
        response=json.dumps(s, indent=2, ensure_ascii=False),
        status=200,
        mimetype="application/json",
    )
    resp.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return resp


@app.route("/api/sessions/<session_id>/export-chat", methods=["GET"])
def export_chat_log(session_id):
    """Export all chat messages grouped by exercise as a .docx file."""
    s = load_session(session_id)
    if not s:
        return jsonify({"error": "Session not found"}), 404

    all_ex_meta = load_exercises()

    # python-docx uses lxml which segfaults under eventlet's monkey-patch.
    # subprocess.run() also deadlocks inside eventlet green threads.
    # Fix: run the blocking subprocess call in a real OS thread via tpool.execute()
    # so eventlet's patched I/O primitives are never involved.
    # All subprocess I/O uses temp files (no pipes at all).
    helper = BASE_DIR / "build_chat_docx.py"
    s_json   = json.dumps(s, ensure_ascii=False)
    ex_json  = json.dumps(all_ex_meta)

    def _build_docx():
        with tempfile.NamedTemporaryFile(suffix=".json", delete=False, mode="w", encoding="utf-8") as sess_f:
            sess_f.write(s_json)
            sess_path = sess_f.name
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as out_f:
            out_path = out_f.name
        with tempfile.NamedTemporaryFile(suffix=".err", delete=False) as err_f:
            err_path = err_f.name
        try:
            with open(out_path, "wb") as out_fh, open(err_path, "wb") as err_fh:
                proc = subprocess.run(
                    [sys.executable, str(helper), sess_path, ex_json],
                    stdout=out_fh,
                    stderr=err_fh,
                    timeout=30,
                )
            err_text = open(err_path).read() if proc.returncode != 0 else ""
            return proc.returncode, err_text, open(out_path, "rb").read()
        finally:
            for p in (sess_path, out_path, err_path):
                try:
                    os.unlink(p)
                except OSError:
                    pass

    returncode, err_text, docx_bytes = eventlet.tpool.execute(_build_docx)
    if returncode != 0:
        log.error("build_chat_docx failed: %s", err_text)
        return jsonify({"error": "Failed to generate document"}), 500

    safe_name = s["session_name"].encode("ascii", "ignore").decode().replace(" ", "_")[:40].strip("_") or session_id[:8]
    filename  = f"DTW_{safe_name}_{session_id[:8]}_ChatLog.docx"

    return Response(
        docx_bytes,
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Content-Type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        },
    )


# ── SocketIO Events ────────────────────────────────────────────────────────────

@socketio.on("join_session")
def handle_join(data):
    session_id   = data.get("session_id", "").strip()
    display_name = (data.get("display_name") or "").strip()
    role         = (data.get("role") or "").strip()

    # Allow room code lookup
    if len(session_id) == 6 and session_id.upper() in _room_code_map:
        session_id = _room_code_map[session_id.upper()]

    s = load_session(session_id)
    if not s:
        emit("error", {"message": "Session not found. Check your room code."})
        return

    if not display_name:
        emit("error", {"message": "Display name is required."})
        return

    join_room(session_id)

    # Add participant if new
    existing = [p["name"] for p in s.get("participants", [])]
    if display_name not in existing:
        s["participants"].append({"name": display_name, "role": role, "joined_at": utcnow()})
        save_session(s)
        socketio.emit("participant_joined", {"name": display_name, "role": role}, room=session_id)
        log.info("Participant '%s' joined session %s", display_name, session_id)

    # Send full session state to this client
    safe = {k: v for k, v in s.items() if k != "passcode"}
    emit("session_state", safe)


@socketio.on("facilitator_join")
def handle_facilitator_join(data):
    session_id = data.get("session_id", "").strip()

    s = load_session(session_id)
    if not s:
        emit("error", {"message": "Session not found."})
        return

    join_room(session_id)
    safe = {k: v for k, v in s.items() if k != "passcode"}
    emit("session_state", safe)
    log.info("Facilitator joined session %s", session_id)


@socketio.on("send_message")
def handle_message(data):
    session_id = data.get("session_id", "").strip()
    sender     = (data.get("sender") or "").strip()
    text       = (data.get("text") or "").strip()
    exercise_id = (data.get("exercise_id") or "").strip()

    if not session_id or not sender or not text:
        emit("error", {"message": "Missing required fields."})
        return

    s = load_session(session_id)
    if not s:
        emit("error", {"message": "Session not found."})
        return

    msg = {
        "sender":      sender,
        "timestamp":   utcnow(),
        "text":        text,
        "exercise_id": exercise_id,
    }

    # Store under the correct exercise
    if exercise_id and exercise_id in s["exercises"]:
        s["exercises"][exercise_id]["messages"].append(msg)
    else:
        # Fallback: store under current exercise
        idx    = s["current_exercise_index"]
        ex_id  = s["exercise_order"][idx] if s["exercise_order"] else None
        if ex_id and ex_id in s["exercises"]:
            msg["exercise_id"] = ex_id
            s["exercises"][ex_id]["messages"].append(msg)

    save_session(s)
    socketio.emit("new_message", msg, room=session_id)


# ── Startup ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    log.info("DTW Workshop App starting on http://localhost:%d", PORT)
    if _s3_enabled():
        log.info("S3 storage enabled: s3://%s/%s/ ✓", S3_BUCKET, S3_PREFIX)
    else:
        log.info("S3 not configured — using local filesystem (set S3_BUCKET in .env to enable)")
        log.info("Sessions directory: %s", SESSIONS_DIR)
    log.info("LLM Gateway configured: %s (model: %s) — keychain read deferred to first request", LLM_GATEWAY_URL, LLM_MODEL)
    socketio.run(app, host="0.0.0.0", port=PORT, debug=False)
