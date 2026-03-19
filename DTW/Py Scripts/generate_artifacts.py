#!/usr/bin/env python3
"""
Design Thinking Workshop Artifacts Generator — Parts 1–3
Hopes & Fears Board, Personas, Stakeholder Map

Design system follows Cloud_DTW_Updated.pptx template:
  - Slide: 13.33" × 7.50" (standard widescreen)
  - Fonts: Arial / Calibri (NO Comic Sans)
  - Lilly brand colors: Red #D52B1E, Blue #0063BE, Green #007A33,
                        Purple #6D2077, Orange #E87722, Teal #00A3E0
  - Layout: full-width red header bar (0.75") + thin accent divider
  - Content cards: white/light bg with 0.06"-wide colored left accent bar
  - NO rotation, NO sticky-note aesthetic
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Lilly Brand Colors ──────────────────────────────────────────────────────
RED     = RGBColor(0xD5, 0x2B, 0x1E)   # primary
BLUE    = RGBColor(0x00, 0x63, 0xBE)
GREEN   = RGBColor(0x00, 0x7A, 0x33)
PURPLE  = RGBColor(0x6D, 0x20, 0x77)
ORANGE  = RGBColor(0xE8, 0x77, 0x22)
TEAL    = RGBColor(0x00, 0xA3, 0xE0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x58, 0x58, 0x58)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xC8, 0xC8, 0xC8)

ACCENT_COLORS = [RED, BLUE, GREEN, PURPLE, ORANGE, TEAL]

# ── Slide dimensions (standard widescreen, matches Cloud_DTW_Updated.pptx) ─
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── Layout constants ────────────────────────────────────────────────────────
HEADER_H     = Inches(0.75)
DIVIDER_H    = Inches(0.04)
LEFT_MARGIN  = Inches(0.5)
ACCENT_W     = Inches(0.06)   # left accent bar width on cards
CARD_PADDING = Inches(0.12)   # text indent from accent bar

# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, prs, title, subtitle=""):
    """Full-width red header bar + thin accent divider + title text."""
    # Red bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    # Title in bar
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.12),
                                  prs.slide_width - Inches(1), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = Pt(28)
    p.font.bold  = True
    p.font.color.rgb = WHITE

    # Thin accent divider below header
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, HEADER_H, prs.slide_width, DIVIDER_H)
    div.fill.solid()
    div.fill.fore_color.rgb = MID_GRAY
    div.line.fill.background()

    # Optional subtitle
    if subtitle:
        stb = slide.shapes.add_textbox(
            LEFT_MARGIN, HEADER_H + DIVIDER_H + Inches(0.08),
            prs.slide_width - Inches(1), Inches(0.35))
        stf = stb.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.name  = "Calibri"
        sp.font.size  = Pt(12)
        sp.font.italic = True
        sp.font.color.rgb = DARK_GRAY


def add_card(slide, left, top, width, height, text, accent_color,
             title=None, font_size=Pt(10), bg_color=None):
    """
    Card with colored left accent bar (matching Cloud_DTW template pattern).
    Optional bold title line above body text.
    """
    bg = bg_color or WHITE

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(0.5)

    # Left accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, ACCENT_W, height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Text box indented past accent bar
    text_left   = left + ACCENT_W + CARD_PADDING
    text_width  = width - ACCENT_W - CARD_PADDING - Inches(0.08)
    text_top    = top + Inches(0.07)
    text_height = height - Inches(0.14)

    tb = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = tb.text_frame
    tf.word_wrap = True

    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name  = "Arial"
        p0.font.size  = Pt(10)
        p0.font.bold  = True
        p0.font.color.rgb = NEAR_BLACK
        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.name  = "Calibri"
        p1.font.size  = font_size
        p1.font.color.rgb = DARK_GRAY
    else:
        p0 = tf.paragraphs[0]
        p0.text = text
        p0.font.name  = "Calibri"
        p0.font.size  = font_size
        p0.font.color.rgb = NEAR_BLACK

    return card


def add_section_label(slide, left, top, width, text, color=None):
    """Bold uppercase section label, optionally colored."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = Pt(11)
    p.font.bold  = True
    p.font.color.rgb = color or NEAR_BLACK
    return tb


def add_numbered_circle(slide, cx, cy, number, color):
    """0.45" diameter circle with bold number — for process flows."""
    d = Inches(0.45)
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - d // 2, cy - d // 2, d, d)
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.text = str(number)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name  = "Arial"
    p.font.size  = Pt(14)
    p.font.bold  = True
    p.font.color.rgb = WHITE


def content_top():
    """Y position where content starts (below header + divider + small gap)."""
    return HEADER_H + DIVIDER_H + Inches(0.18)


# ── Artifact 1: Hopes & Fears Board ────────────────────────────────────────

def create_hopes_fears_board():
    """OUTPUT 1: Hopes & Fears Board"""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Hopes & Fears")

    ct = content_top()
    usable_h = SLIDE_H - ct - Inches(0.15)
    half_h   = usable_h / 2 - Inches(0.08)

    col_w  = (SLIDE_W - Inches(1.0)) / 2 - Inches(0.1)
    left_x = LEFT_MARGIN
    right_x = LEFT_MARGIN + col_w + Inches(0.2)

    # Section headers
    add_section_label(slide, left_x,  ct, Inches(2.5), "HOPES", GREEN)
    add_section_label(slide, right_x, ct, Inches(2.5), "FEARS", RED)

    label_h = Inches(0.32)

    hopes = [
        "Visibility on the hand offs so ADCO/MCO can jump in and help manage speed",
        "Quicker handoff to Media partners means faster time to market",
        "Standardization across lines of business and business units",
        "Role clarity among everyone involved in the process",
        "Better systems in place that talk to each other",
        "Clear ownership of Creative Grid process within Lilly",
        "Revamp QA process before assets sent to agency",
        "Clear, end-to-end DTC production timeframes established",
        "Consistency in approval process timing"
    ]

    fears = [
        "So many roles constantly changing — impossible to keep up",
        "A lot of black boxes — unclear who is supposed to be doing what",
        "Ongoing improvement needed in trafficking process to cut revisions and swirl",
        "Too many cooks in the kitchen for the trafficking process",
        "Roles are not clearly owned across teams",
        "Digital team may traffic assets too quickly before marketer sends email",
        "Too many handoffs can feel heavy and frustrating",
        "Loss of control when it gets to the creative grid process",
        "No clear standard on who fills grids, who traffics, who fixes failures",
        "Number of handoffs in this process can feel heavy and sometimes frustrating"
    ]

    card_h  = Inches(0.54)
    gap     = Inches(0.07)
    start_y = ct + label_h + Inches(0.04)

    for i, hope in enumerate(hopes):
        y = start_y + i * (card_h + gap)
        if y + card_h > SLIDE_H - Inches(0.1):
            break
        add_card(slide, left_x, y, col_w, card_h, hope, GREEN,
                 font_size=Pt(9))

    for i, fear in enumerate(fears):
        y = start_y + i * (card_h + gap)
        if y + card_h > SLIDE_H - Inches(0.1):
            break
        add_card(slide, right_x, y, col_w, card_h, fear, RED,
                 font_size=Pt(9))

    out = '/Users/V5X8512/Downloads/01_hopes_fears.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 2: Personas ─────────────────────────────────────────────────────

def create_personas():
    """OUTPUT 2: Persona Cards — one slide per persona."""
    prs = new_prs()

    personas = [
        {
            "name": "Mia — HCP Marketer",
            "role": "HCP Marketing Lead",
            "quote": "I feel like there are a lot of handoffs that can make me feel frustrated or wish it could be a click of a button.",
            "goals":     ["Kick off banners with brand strategy",
                          "Project manage through MLR and creative grid",
                          "Be proactive 1–2 months out",
                          "Think through HCP experience"],
            "pains":     ["Too many handoffs in the process",
                          "5–10 min to enter inputs — hard to find time",
                          "Loss of control in creative grid process",
                          "Communication gaps leave banners hanging"],
            "behaviors": ["Presents idea to LMS or E2E colleagues",
                          "Takes through MLR for approval",
                          "Sends CMI team updated grid with banner rotation",
                          "Checks in via email, Teams, or NPP weekly status"],
            "tools":     ["MLR", "Excel / Creative Grid", "Email", "Teams"],
            "color": RED,
        },
        {
            "name": "Julie — DTC Media Lead",
            "role": "DTC Media Strategy",
            "quote": "We need to work more strongly with the Lilly organization to understand impact when assets are not available.",
            "goals":     ["Ensure placements are strategically placed",
                          "Coordinate timing for market impact",
                          "Drive scripts through targeted messaging"],
            "pains":     ["Constant moving pieces to coordinate",
                          "Assets delayed or unavailable",
                          "Dilutes marketplace impact to drive scripts"],
            "behaviors": ["Strategically determine placements for ad unit types",
                          "Provide specs to TA & GCO CC&A team",
                          "Confirm creative grid inputs are accurate",
                          "Approve for agency handoff"],
            "tools":     ["Media planning tools", "Creative Grid", "Google Studio"],
            "color": BLUE,
        },
        {
            "name": "Jenn — HCP Agency (CMI)",
            "role": "Agency Media Operations",
            "quote": "CMI often has to help facilitate the task to see it through to completion in a timely manner.",
            "goals":     ["Receive and QA assets properly",
                          "Build and launch media accurately",
                          "Meet launch timelines"],
            "pains":     ["Too many cooks in the kitchen",
                          "Roles not clearly owned",
                          "Click tags often missing; expiration dates not updated"],
            "behaviors": ["Receive assets via Google Studio or SharePoint",
                          "QA for web specs & creative grids",
                          "Assign creatives in internal systems",
                          "Build tags and send to suppliers"],
            "tools":     ["Google Studio", "SharePoint", "Ad serving platforms"],
            "color": GREEN,
        },
        {
            "name": "Matt — Digital Brand Lead",
            "role": "Digital Development Strategy",
            "quote": "There is no clear standard on who should be filling out the grids, who is trafficking, and who is fixing failed creatives.",
            "goals":     ["Ensure assets delivered on time",
                          "Validate creatives prior to trafficking",
                          "Remove blocks and escalate issues"],
            "pains":     ["No clear standard on responsibilities",
                          "Called in when information is missing",
                          "Unclear who fixes failed creatives"],
            "behaviors": ["Partners with brands on digital best practices",
                          "Consults on digital tactic design for compliance",
                          "Oversees QA to validate creatives"],
            "tools":     ["Adobe Workfront (proposed)", "Creative trafficking forms"],
            "color": PURPLE,
        },
        {
            "name": "Kate — Campaign Execution",
            "role": "Campaign Activation Strategy",
            "quote": "There are many hands in the pot. Order of operations is not clear. Roles and responsibilities are fuzzy.",
            "goals":     ["Work with brand & omnichannel on banner needs",
                          "Program banners in campaign software when approved"],
            "pains":     ["Many hands in the pot",
                          "Order of operations unclear",
                          "Combination creates delays"],
            "behaviors": ["Receive/assess media plan",
                          "Determine if banners need to be created",
                          "Alert marketing and brand DCO to include in backlog"],
            "tools":     ["CAP (Campaign Activation Platform)", "Media plans"],
            "color": ORANGE,
        },
        {
            "name": "Molly — Media Operations",
            "role": "Media Ops Process Owner",
            "quote": "There are so many roles in this organization constantly changing — it's impossible to keep up.",
            "goals":     ["Ensure trafficking handoff processes work well",
                          "Provide necessary tools like Smartsheet",
                          "Enable teams to work together confidently"],
            "pains":     ["Roles constantly changing",
                          "Black boxes — unclear who does what",
                          "When problems arise, hard to know who to triage"],
            "behaviors": ["Provides process documentation",
                          "Manages Smartsheet tools",
                          "Facilitates communication between groups"],
            "tools":     ["Smartsheet", "Process documentation"],
            "color": TEAL,
        },
    ]

    for idx, p in enumerate(personas):
        slide = blank_slide(prs)
        add_header(slide, prs,
                   p["name"],
                   subtitle=p["role"])

        ct = content_top() + Inches(0.38)   # extra room for subtitle

        # Quote bar — full-width light card
        quote_h = Inches(0.50)
        qcard = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            LEFT_MARGIN, ct,
            SLIDE_W - Inches(1.0), quote_h)
        qcard.fill.solid()
        qcard.fill.fore_color.rgb = LIGHT_GRAY
        qcard.line.fill.background()

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            LEFT_MARGIN, ct, ACCENT_W, quote_h)
        accent.fill.solid()
        accent.fill.fore_color.rgb = p["color"]
        accent.line.fill.background()

        qtb = slide.shapes.add_textbox(
            LEFT_MARGIN + ACCENT_W + CARD_PADDING, ct + Inches(0.08),
            SLIDE_W - Inches(1.3), Inches(0.38))
        qtf = qtb.text_frame
        qtf.text = f'"{p["quote"]}"'
        qpar = qtf.paragraphs[0]
        qpar.font.name   = "Calibri"
        qpar.font.size   = Pt(10)
        qpar.font.italic = True
        qpar.font.color.rgb = DARK_GRAY

        # 4 section columns: Goals | Pains | Behaviors | Tools
        sections = [
            ("GOALS & MOTIVATIONS", p["goals"],     GREEN),
            ("PAIN POINTS",         p["pains"],      RED),
            ("KEY BEHAVIORS",       p["behaviors"],  BLUE),
            ("TOOLS & SYSTEMS",     p["tools"],      TEAL),
        ]

        col_area_top = ct + quote_h + Inches(0.14)
        col_area_h   = SLIDE_H - col_area_top - Inches(0.12)
        n_cols       = len(sections)
        col_w        = (SLIDE_W - Inches(1.0)) / n_cols - Inches(0.08)

        for ci, (label, items, color) in enumerate(sections):
            cx = LEFT_MARGIN + ci * (col_w + Inches(0.08))
            # Column header
            add_section_label(slide, cx, col_area_top, col_w, label, color)
            # Items as cards
            card_h  = Inches(0.56)
            gap     = Inches(0.06)
            item_y  = col_area_top + Inches(0.32)
            for item in items[:5]:
                if item_y + card_h > SLIDE_H - Inches(0.08):
                    break
                add_card(slide, cx, item_y, col_w, card_h, item, color,
                         font_size=Pt(9))
                item_y += card_h + gap

    out = '/Users/V5X8512/Downloads/02_personas.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 3: Stakeholder Map ──────────────────────────────────────────────

def create_stakeholder_map():
    """OUTPUT 3: Stakeholder Map — concentric ring layout."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Stakeholder Map",
               subtitle="Banner Display Trafficking & Handoff Process")

    ct = content_top()

    # Center oval — the initiative
    cx  = SLIDE_W / 2
    cy  = ct + (SLIDE_H - ct) / 2 + Inches(0.1)
    ow  = Inches(2.2)
    oh  = Inches(1.1)
    center_oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - ow / 2, cy - oh / 2, ow, oh)
    center_oval.fill.solid()
    center_oval.fill.fore_color.rgb = RED
    center_oval.line.fill.background()
    ctf = center_oval.text_frame
    ctf.text = "Banner\nTrafficking\nProcess"
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cp.font.name  = "Arial"
    cp.font.size  = Pt(12)
    cp.font.bold  = True
    cp.font.color.rgb = WHITE
    for pp in ctf.paragraphs[1:]:
        pp.alignment = PP_ALIGN.CENTER
        pp.font.name  = "Arial"
        pp.font.size  = Pt(12)
        pp.font.bold  = True
        pp.font.color.rgb = WHITE

    # Stakeholder groups: (group, role_detail, quote, x_offset, y_offset, color)
    stakeholders = [
        ("ADCO / MCO",      "Carrie Sieglitz / Jordyn King",
         "We tend to be asked: where is the banner? Why doesn't CMI have the files?",
         -4.5, -1.8, RED),
        ("HCP Marketers",   "Mia McCrumb",
         "I have to be proactive (1–2 months out) to ensure the right message gets to market.",
         -4.5,  0.0, BLUE),
        ("DTC Marketers",   "Jamie Doyle",
         "The DTC marketer should own the media grid so the marketer stays in the loop.",
         -4.5,  1.8, BLUE),
        ("Media Leads",     "Julie Gogan / Alex Lund",
         "We need to understand impact when assets are not available.",
         -1.5, -2.4, GREEN),
        ("Media Operations","Molly Hudlow",
         "So many roles constantly changing — impossible to keep up.",
         1.5,  -2.4, GREEN),
        ("HCP Agency (CMI)","Jenn Margiloff",
         "CMI often has to help facilitate to see it through completion in a timely manner.",
         4.5,  -1.2, PURPLE),
        ("Consumer Agency", "Alex Fox",
         "Better learning of processes between teams to aid in collaboration.",
         4.5,   0.8, PURPLE),
        ("Digital Brand Lead","Matt Casolaro",
         "No clear standard on who fills the grids, who traffics, who fixes failures.",
         2.0,   2.2, ORANGE),
        ("Campaign Execution","Kate Kilgore",
         "Many hands in the pot. Order of operations is not clear.",
         -2.0,  2.2, ORANGE),
    ]

    card_w = Inches(2.5)
    card_h = Inches(1.1)

    for group, role, quote, dx, dy, color in stakeholders:
        x = cx + Inches(dx) - card_w / 2
        y = cy + Inches(dy) - card_h / 2
        # Clamp to slide
        x = max(Inches(0.1), min(x, SLIDE_W - card_w - Inches(0.1)))
        y = max(ct + Inches(0.05), min(y, SLIDE_H - card_h - Inches(0.05)))
        add_card(slide, x, y, card_w, card_h,
                 quote, color,
                 title=f"{group}  ·  {role}",
                 font_size=Pt(8))

    out = '/Users/V5X8512/Downloads/03_stakeholder_map.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Entry point ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating Design Thinking Workshop Artifacts 1–3...")
    print("=" * 60)
    create_hopes_fears_board()
    create_personas()
    create_stakeholder_map()
    print("=" * 60)
    print("Completed artifacts 1–3.")
