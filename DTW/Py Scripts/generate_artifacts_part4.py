#!/usr/bin/env python3
"""
Design Thinking Workshop Artifacts Generator — Parts 11–14
Experience Roadmap, Hills, Gantt Roadmap, Resource Plan
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

RED     = RGBColor(0xD5, 0x2B, 0x1E)
BLUE    = RGBColor(0x00, 0x63, 0xBE)
GREEN   = RGBColor(0x00, 0x7A, 0x33)
PURPLE  = RGBColor(0x6D, 0x20, 0x77)
ORANGE  = RGBColor(0xE8, 0x77, 0x22)
TEAL    = RGBColor(0x00, 0xA3, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x58, 0x58, 0x58)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xC8, 0xC8, 0xC8)

SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.50)
HEADER_H     = Inches(0.75)
DIVIDER_H    = Inches(0.04)
LEFT_MARGIN  = Inches(0.5)
ACCENT_W     = Inches(0.06)
CARD_PADDING = Inches(0.12)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, prs, title, subtitle=""):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.12),
                                  prs.slide_width - Inches(1), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, HEADER_H, prs.slide_width, DIVIDER_H)
    div.fill.solid()
    div.fill.fore_color.rgb = MID_GRAY
    div.line.fill.background()
    if subtitle:
        stb = slide.shapes.add_textbox(
            LEFT_MARGIN, HEADER_H + DIVIDER_H + Inches(0.06),
            prs.slide_width - Inches(1), Inches(0.32))
        stf = stb.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.name   = "Calibri"
        sp.font.size   = Pt(11)
        sp.font.italic = True
        sp.font.color.rgb = DARK_GRAY


def add_card(slide, left, top, width, height, text, accent_color,
             title=None, font_size=Pt(10), bg_color=None):
    bg = bg_color or WHITE
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(0.5)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, ACCENT_W, height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    text_left  = left + ACCENT_W + CARD_PADDING
    text_width = width - ACCENT_W - CARD_PADDING - Inches(0.08)
    tb = slide.shapes.add_textbox(
        text_left, top + Inches(0.06), text_width, height - Inches(0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name  = "Arial"
        p0.font.size  = Pt(9)
        p0.font.bold  = True
        p0.font.color.rgb = NEAR_BLACK
        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.name  = "Calibri"
        p1.font.size  = font_size
        p1.font.color.rgb = DARK_GRAY
    else:
        p0 = tf.paragraphs[0]
        p0.text = text
        p0.font.name  = "Calibri"
        p0.font.size  = font_size
        p0.font.color.rgb = NEAR_BLACK


def add_section_label(slide, left, top, width, text, color=None, size=Pt(10)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = size
    p.font.bold  = True
    p.font.color.rgb = color or NEAR_BLACK


def content_top(has_subtitle=False):
    extra = Inches(0.38) if has_subtitle else Inches(0.14)
    return HEADER_H + DIVIDER_H + extra


# ── Artifact 11: Experience-Based Roadmap ────────────────────────────────────

def create_experience_roadmap():
    """OUTPUT 11: Experience-Based Roadmap — 3 columns (Near / Mid / Long term)."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs,
               "Experience-Based Roadmap",
               subtitle="Our user will be able to...")

    ct = content_top(has_subtitle=True)

    columns = [
        ("NEAR-TERM  (Months 1–3)",  RED,    [
            "See clear roles and responsibilities for each step of the process",
            "Access a marketer cheat sheet with required steps before handoff",
            "Receive Teams channel notifications when handoffs are ready",
            "View a quarterly rotation matrix aligned with the creative grid",
            "Track banner status: Not Started → In Progress → Ready → Launched",
            "Complete grid inputs in under 2 minutes with pre-filled fields",
        ]),
        ("MID-TERM  (Months 4–6)",   BLUE,   [
            "Benefit from automated grid notifications when steps complete",
            "Access a high-level creative map showing what launches when and why",
            "Leverage auto-populated metadata in the creative grid",
            "See agencies receive assets with complete information in one handoff",
            "Utilize standardized banner templates across HCP and DTC lines",
            "Experience one-day turnaround from AFD to CMI QA",
        ]),
        ("LONG-TERM  (Months 7–12)", GREEN,  [
            "Benefit from fully automated metadata collection and grid population",
            "Work in a unified workflow system with Adobe Workfront integration",
            "Leverage a fast-track path for high-urgency messages",
            "Experience seamless system integration where tools talk to each other",
            "Feel confident banners reach market without communication gaps",
            "Achieve cross-BU standardization across all business units and lines",
        ]),
    ]

    col_gap = Inches(0.12)
    col_w   = (SLIDE_W - Inches(1.0) - col_gap * 2) / 3
    card_h  = Inches(0.70)
    card_gap = Inches(0.06)
    label_h = Inches(0.36)

    for ci, (label, color, items) in enumerate(columns):
        cx = LEFT_MARGIN + ci * (col_w + col_gap)

        # Column header band
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, ct, col_w, Inches(0.38))
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.fill.background()
        btb = slide.shapes.add_textbox(cx + Inches(0.1), ct + Inches(0.06),
                                       col_w - Inches(0.2), Inches(0.30))
        btf = btb.text_frame
        btf.text = label
        bp = btf.paragraphs[0]
        bp.font.name  = "Arial"
        bp.font.size  = Pt(10)
        bp.font.bold  = True
        bp.font.color.rgb = WHITE

        iy = ct + Inches(0.44)
        for item in items:
            if iy + card_h > SLIDE_H - Inches(0.1):
                break
            add_card(slide, cx, iy, col_w, card_h, item, color,
                     font_size=Pt(9))
            iy += card_h + card_gap

    # Vertical dividers between columns
    for ci in range(1, 3):
        dx = LEFT_MARGIN + ci * (col_w + col_gap) - col_gap / 2
        dv = slide.shapes.add_connector(2, dx, ct, dx, SLIDE_H - Inches(0.1))
        dv.line.width = Pt(0.5)
        dv.line.color.rgb = MID_GRAY

    out = '/Users/V5X8512/Downloads/11_experience_roadmap.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 12: Hills / Objectives ─────────────────────────────────────────

def create_hills():
    """OUTPUT 12: Hills — numbered cards with WHO / WHAT / WOW structure."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Hills",
               subtitle="Meaningful outcomes we commit to delivering")

    ct = content_top(has_subtitle=True)

    hills = [
        {
            "num":  1,
            "who":  "Marketers and Media Teams",
            "what": "Complete banner trafficking from AFD to live in market within 1 business day with zero communication gaps",
            "wow":  "Reducing time-to-market by 80% and eliminating all handoff-related delays",
            "color": RED,
        },
        {
            "num":  2,
            "who":  "All trafficking process stakeholders",
            "what": "Know exactly who owns each step, what actions are required, and receive automated notifications when it's their turn",
            "wow":  "Achieving 100% role clarity and eliminating all black box confusion across teams",
            "color": BLUE,
        },
        {
            "num":  3,
            "who":  "HCP and Consumer Agencies",
            "what": "Receive complete, accurate, error-free banner assets with all required metadata in a single automated handoff",
            "wow":  "Reducing agency QA rework by 90% and consistently meeting SLA requirements",
            "color": GREEN,
        },
    ]

    card_h  = (SLIDE_H - ct - Inches(0.2)) / len(hills) - Inches(0.10)
    card_w  = SLIDE_W - Inches(1.0)

    for i, hill in enumerate(hills):
        cy = ct + i * (card_h + Inches(0.10))
        color = hill["color"]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MID_GRAY
        card.line.width = Pt(0.5)

        # Left accent bar
        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, cy, ACCENT_W, card_h)
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        # Numbered circle
        circ_d = Inches(0.48)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            LEFT_MARGIN + ACCENT_W + Inches(0.12), cy + (card_h - circ_d) / 2,
            circ_d, circ_d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.text = str(hill["num"])
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.font.name  = "Arial"
        cp.font.size  = Pt(18)
        cp.font.bold  = True
        cp.font.color.rgb = WHITE

        # Text content
        text_x = LEFT_MARGIN + ACCENT_W + Inches(0.75)
        text_w = card_w - ACCENT_W - Inches(0.82)
        tb = slide.shapes.add_textbox(text_x, cy + Inches(0.10), text_w, card_h - Inches(0.20))
        tf = tb.text_frame
        tf.word_wrap = True

        p_who = tf.paragraphs[0]
        p_who.text = f"WHO:   {hill['who']}"
        p_who.font.name  = "Arial"
        p_who.font.size  = Pt(10)
        p_who.font.color.rgb = DARK_GRAY

        p_what = tf.add_paragraph()
        p_what.text = f"WHAT:  {hill['what']}"
        p_what.font.name  = "Arial"
        p_what.font.size  = Pt(10)
        p_what.font.color.rgb = NEAR_BLACK

        p_wow = tf.add_paragraph()
        p_wow.text = f"WOW:   {hill['wow']}"
        p_wow.font.name  = "Arial"
        p_wow.font.size  = Pt(10)
        p_wow.font.bold  = True
        p_wow.font.color.rgb = color

    out = '/Users/V5X8512/Downloads/12_hills.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 13: Gantt Roadmap ───────────────────────────────────────────────

def create_gantt_roadmap():
    """OUTPUT 13: Gantt-style 12-month roadmap."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "12-Month Roadmap",
               subtitle="Initiatives grouped by theme — short / mid / long term")

    ct = content_top(has_subtitle=True)

    # Timeline dimensions
    name_col_w   = Inches(3.2)
    timeline_x   = LEFT_MARGIN + name_col_w + Inches(0.1)
    timeline_w   = SLIDE_W - timeline_x - Inches(0.4)
    month_w      = timeline_w / 12

    # Zone header
    zones = [
        ("SHORT-TERM",  0,  3, RGBColor(0xFF, 0xEE, 0xEE)),
        ("MID-TERM",    3,  3, RGBColor(0xE8, 0xF0, 0xFB)),
        ("LONG-TERM",   6,  6, RGBColor(0xED, 0xF4, 0xED)),
    ]
    zone_h = Inches(0.30)
    for name, start, dur, fill in zones:
        zx = timeline_x + start * month_w
        zb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, zx, ct, dur * month_w - Inches(0.02), zone_h)
        zb.fill.solid()
        zb.fill.fore_color.rgb = fill
        zb.line.fill.background()
        ztb = slide.shapes.add_textbox(zx + Inches(0.06), ct + Inches(0.04),
                                       dur * month_w - Inches(0.14), Inches(0.24))
        ztf = ztb.text_frame
        ztf.text = name
        zp = ztf.paragraphs[0]
        zp.font.name  = "Arial"
        zp.font.size  = Pt(8)
        zp.font.bold  = True
        zp.font.color.rgb = DARK_GRAY

    # Month numbers
    month_row_y = ct + zone_h
    month_row_h = Inches(0.22)
    for m in range(12):
        mx = timeline_x + m * month_w
        mtb = slide.shapes.add_textbox(mx, month_row_y, month_w, month_row_h)
        mtf = mtb.text_frame
        mtf.text = str(m + 1)
        mp = mtf.paragraphs[0]
        mp.alignment = PP_ALIGN.CENTER
        mp.font.name  = "Calibri"
        mp.font.size  = Pt(8)
        mp.font.color.rgb = DARK_GRAY

    # Vertical month grid lines
    for m in range(1, 12):
        gx = timeline_x + m * month_w
        gv = slide.shapes.add_connector(2, gx, ct, gx, SLIDE_H - Inches(0.1))
        gv.line.width = Pt(0.25)
        gv.line.color.rgb = MID_GRAY

    # Initiative rows
    initiatives = [
        {
            "theme": "QUICK WINS",
            "color": RED,
            "items": [
                ("Teams Channel Notifications",        0, 1),
                ("Role Clarity Documentation (RACI)",  0, 2),
                ("Marketer Cheat Sheet",               1, 1),
                ("Automated Grid Notifications",       0, 2),
            ],
        },
        {
            "theme": "PROCESS & COMMUNICATION",
            "color": BLUE,
            "items": [
                ("Status Tracker Implementation",      1, 2),
                ("Quarterly Rotation Matrix Template", 1, 2),
                ("High-Level Creative Map",            3, 2),
                ("Enhanced QA Process",                4, 2),
            ],
        },
        {
            "theme": "GRID & SYSTEMS",
            "color": GREEN,
            "items": [
                ("Creative Grid Simplification",       2, 3),
                ("Metadata Auto-Population",           4, 4),
                ("Workfront Integration",              6, 4),
                ("Full Grid Automation",               8, 4),
            ],
        },
        {
            "theme": "STANDARDIZATION & SCALE",
            "color": ORANGE,
            "items": [
                ("Fast-Track Templates",               5, 3),
                ("Cross-BU Standardization",           7, 5),
                ("Unified Workflow System",            9, 3),
            ],
        },
    ]

    row_start_y = month_row_y + month_row_h + Inches(0.06)
    avail_row_h = SLIDE_H - row_start_y - Inches(0.12)
    n_items_total = sum(len(g["items"]) for g in initiatives) + len(initiatives)
    row_h = avail_row_h / n_items_total
    item_h = max(Inches(0.20), row_h * 0.85)

    current_y = row_start_y

    for group in initiatives:
        # Theme label
        tlb = slide.shapes.add_textbox(LEFT_MARGIN, current_y,
                                       name_col_w, item_h)
        tlf = tlb.text_frame
        tlf.text = group["theme"]
        tp = tlf.paragraphs[0]
        tp.font.name  = "Arial"
        tp.font.size  = Pt(8.5)
        tp.font.bold  = True
        tp.font.color.rgb = group["color"]
        current_y += item_h

        for name, start_m, dur_m in group["items"]:
            # Initiative name
            nlb = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.1), current_y,
                                           name_col_w - Inches(0.1), item_h)
            nlf = nlb.text_frame
            nlf.text = name
            np_ = nlf.paragraphs[0]
            np_.font.name  = "Calibri"
            np_.font.size  = Pt(8)
            np_.font.color.rgb = NEAR_BLACK

            # Gantt bar
            bx = timeline_x + start_m * month_w
            bw = dur_m * month_w - Inches(0.04)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, bx, current_y + Inches(0.03),
                bw, item_h - Inches(0.06))
            bar.fill.solid()
            bar.fill.fore_color.rgb = group["color"]
            bar.line.fill.background()

            current_y += item_h

    out = '/Users/V5X8512/Downloads/13_roadmap.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 14: Resource Plan ───────────────────────────────────────────────

def create_resource_plan():
    """OUTPUT 14: Resource Plan — table with effort color coding."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Resource Plan",
               subtitle="Roles, owners, effort, and dependencies by initiative")

    ct = content_top(has_subtitle=True)

    headers = ["Initiative", "Required Role / Skill", "Owner / Team", "Effort", "Dependencies", "Gap?"]
    col_widths = [Inches(2.8), Inches(2.4), Inches(2.0), Inches(0.90), Inches(2.8), Inches(0.75)]
    header_h = Inches(0.36)
    row_h    = Inches(0.38)

    effort_colors = {
        "Low":    RGBColor(0xE8, 0xF5, 0xE9),
        "Medium": RGBColor(0xFF, 0xF3, 0xE0),
        "High":   RGBColor(0xFF, 0xEB, 0xEB),
    }
    effort_text_colors = {
        "Low":    GREEN,
        "Medium": ORANGE,
        "High":   RED,
    }

    resources = [
        ("Teams Channel Notifications",     "IT / Platform Admin",           "Media Ops",           "Low",    "Teams workspace setup",             ""),
        ("Automated Grid Notifications",    "Developer",                     "TBD",                 "Medium", "Smartsheet API access",             "⚠"),
        ("RACI Documentation",             "Process Owner",                 "Media Ops — Molly",   "Low",    "Stakeholder input",                 ""),
        ("Marketer Cheat Sheet",           "Process Designer",              "ADCO — Carrie",       "Low",    "RACI complete",                     ""),
        ("Status Tracker",                 "PM / Business Analyst",         "TBD",                 "Medium", "Tool selection",                    "⚠"),
        ("Quarterly Rotation Matrix",      "Marketing Ops",                 "Mia / Jamie",         "Low",    "Template design",                   ""),
        ("Creative Grid Simplification",   "UX Designer + Developer",      "TBD",                 "High",   "User research, stakeholder sign-off","⚠"),
        ("Metadata Auto-Population",       "Developer + Data Analyst",     "TBD",                 "High",   "Source system integration",         "⚠"),
        ("Workfront Integration",          "Workfront Admin + Developer",  "Digital PM Team",     "High",   "Workfront license",                 ""),
        ("Cross-BU Standardization",       "Change Management",            "TBD",                 "High",   "Executive sponsorship",             "⚠"),
    ]

    # Draw header row
    x0 = LEFT_MARGIN
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        hbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x0, ct, cw - Inches(0.02), header_h)
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = DARK_GRAY
        hbox.line.fill.background()
        htb = slide.shapes.add_textbox(
            x0 + Inches(0.07), ct + Inches(0.06), cw - Inches(0.12), header_h - Inches(0.10))
        htf = htb.text_frame
        htf.text = hdr
        hp = htf.paragraphs[0]
        hp.font.name  = "Arial"
        hp.font.size  = Pt(9)
        hp.font.bold  = True
        hp.font.color.rgb = WHITE
        x0 += cw

    # Draw data rows
    for ri, (init, role, owner, effort, deps, gap) in enumerate(resources):
        ry = ct + header_h + ri * row_h
        if ry + row_h > SLIDE_H - Inches(0.06):
            break
        row_bg = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else LIGHT_GRAY
        x0 = LEFT_MARGIN
        cells = [init, role, owner, effort, deps, gap]
        for ci, (cell_text, cw) in enumerate(zip(cells, col_widths)):
            is_effort = (ci == 3)
            bg = effort_colors.get(cell_text, row_bg) if is_effort else row_bg
            cbox = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0, ry, cw - Inches(0.02), row_h - Inches(0.02))
            cbox.fill.solid()
            cbox.fill.fore_color.rgb = bg
            cbox.line.color.rgb = MID_GRAY
            cbox.line.width = Pt(0.5)
            ctb = slide.shapes.add_textbox(
                x0 + Inches(0.07), ry + Inches(0.06),
                cw - Inches(0.12), row_h - Inches(0.10))
            ctf = ctb.text_frame
            ctf.word_wrap = True
            ctf.text = cell_text
            cp = ctf.paragraphs[0]
            cp.font.name  = "Calibri"
            cp.font.size  = Pt(8.5)
            if is_effort:
                cp.font.bold  = True
                cp.font.color.rgb = effort_text_colors.get(cell_text, NEAR_BLACK)
            elif ci == 5 and cell_text:  # gap warning
                cp.font.bold  = True
                cp.font.color.rgb = RED
                cp.alignment = PP_ALIGN.CENTER
            else:
                cp.font.color.rgb = NEAR_BLACK
            x0 += cw

    # Legend
    leg = slide.shapes.add_textbox(
        LEFT_MARGIN, SLIDE_H - Inches(0.30),
        Inches(8), Inches(0.24))
    legtf = leg.text_frame
    legtf.text = "⚠ = Unassigned / capability gap     Effort: Green = Low  |  Amber = Medium  |  Red = High"
    legp = legtf.paragraphs[0]
    legp.font.name   = "Calibri"
    legp.font.size   = Pt(8)
    legp.font.italic = True
    legp.font.color.rgb = DARK_GRAY

    out = '/Users/V5X8512/Downloads/14_resource_plan.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Entry point ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating Design Thinking Workshop Artifacts 11–14...")
    print("=" * 60)
    create_experience_roadmap()
    create_hills()
    create_gantt_roadmap()
    create_resource_plan()
    print("=" * 60)
    print("Completed artifacts 11–14.")
