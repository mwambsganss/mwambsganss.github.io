#!/usr/bin/env python3
"""
Design Thinking Workshop Artifacts Generator — Parts 7–10
Assumptions Grid, Big Ideas, Prioritization Grid, To-Be Process Map
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

RED     = RGBColor(0xD5, 0x2B, 0x1E)
BLUE    = RGBColor(0x00, 0x63, 0xBE)
GREEN   = RGBColor(0x00, 0x7A, 0x33)
PURPLE  = RGBColor(0x6D, 0x20, 0x77)
ORANGE  = RGBColor(0xE8, 0x77, 0x22)
TEAL    = RGBColor(0x00, 0xA3, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x58, 0x58, 0x58)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xC8, 0xC8, 0xC8)

SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.50)
HEADER_H     = Inches(0.75)
DIVIDER_H    = Inches(0.04)
LEFT_MARGIN  = Inches(0.5)
ACCENT_W     = Inches(0.06)
CARD_PADDING = Inches(0.12)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, prs, title, subtitle=""):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.12),
                                  prs.slide_width - Inches(1), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, HEADER_H, prs.slide_width, DIVIDER_H)
    div.fill.solid()
    div.fill.fore_color.rgb = MID_GRAY
    div.line.fill.background()
    if subtitle:
        stb = slide.shapes.add_textbox(
            LEFT_MARGIN, HEADER_H + DIVIDER_H + Inches(0.06),
            prs.slide_width - Inches(1), Inches(0.32))
        stf = stb.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.name   = "Calibri"
        sp.font.size   = Pt(11)
        sp.font.italic = True
        sp.font.color.rgb = DARK_GRAY


def add_card(slide, left, top, width, height, text, accent_color,
             title=None, font_size=Pt(10), bg_color=None):
    bg = bg_color or WHITE
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(0.5)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, ACCENT_W, height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    text_left  = left + ACCENT_W + CARD_PADDING
    text_width = width - ACCENT_W - CARD_PADDING - Inches(0.08)
    tb = slide.shapes.add_textbox(
        text_left, top + Inches(0.06), text_width, height - Inches(0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name  = "Arial"
        p0.font.size  = Pt(9)
        p0.font.bold  = True
        p0.font.color.rgb = NEAR_BLACK
        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.name  = "Calibri"
        p1.font.size  = font_size
        p1.font.color.rgb = DARK_GRAY
    else:
        p0 = tf.paragraphs[0]
        p0.text = text
        p0.font.name  = "Calibri"
        p0.font.size  = font_size
        p0.font.color.rgb = NEAR_BLACK


def add_section_label(slide, left, top, width, text, color=None, size=Pt(10)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = size
    p.font.bold  = True
    p.font.color.rgb = color or NEAR_BLACK


def content_top(has_subtitle=False):
    extra = Inches(0.38) if has_subtitle else Inches(0.14)
    return HEADER_H + DIVIDER_H + extra


def add_axis_label(slide, left, top, width, height, text, color=NEAR_BLACK, size=Pt(9)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name  = "Arial"
    p.font.size  = size
    p.font.bold  = True
    p.font.color.rgb = color


# ── Artifact 7: Assumptions & Questions Grid ────────────────────────────────

def create_assumptions_grid():
    """OUTPUT 7: Assumptions & Questions 2×2 Grid (Risk × Certainty)."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Assumptions & Questions",
               subtitle="Certainty (x-axis)  ·  Risk (y-axis)")

    ct = content_top(has_subtitle=True)

    # Grid dimensions
    axis_label_w = Inches(0.80)
    axis_label_h = Inches(0.30)
    grid_left    = LEFT_MARGIN + axis_label_w + Inches(0.1)
    grid_top     = ct
    grid_w       = SLIDE_W - grid_left - Inches(0.4)
    grid_h       = SLIDE_H - ct - Inches(0.5)
    half_w       = grid_w / 2
    half_h       = grid_h / 2

    # Quadrant backgrounds
    quad_styles = [
        # (col, row, fill_color, label, label_color)
        (0, 0, RGBColor(0xFF, 0xED, 0xED), "HIGH RISK / CERTAIN\nTest & Validate",    RED),
        (1, 0, RGBColor(0xFF, 0xE8, 0xE0), "HIGH RISK / UNCERTAIN\nPriority to Resolve", RED),
        (0, 1, RGBColor(0xED, 0xF4, 0xED), "LOW RISK / CERTAIN\nMonitor",              GREEN),
        (1, 1, RGBColor(0xE8, 0xF0, 0xFB), "LOW RISK / UNCERTAIN\nExplore",            BLUE),
    ]
    for col, row, fill, label, lc in quad_styles:
        qx = grid_left + col * half_w
        qy = grid_top  + row * half_h
        qbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, qx, qy, half_w, half_h)
        qbox.fill.solid()
        qbox.fill.fore_color.rgb = fill
        qbox.line.color.rgb = MID_GRAY
        qbox.line.width = Pt(0.5)
        # Quadrant label (top-left of quadrant)
        ltb = slide.shapes.add_textbox(
            qx + Inches(0.1), qy + Inches(0.06), half_w - Inches(0.2), Inches(0.42))
        ltf = ltb.text_frame
        ltf.text = label
        lp = ltf.paragraphs[0]
        lp.font.name  = "Arial"
        lp.font.size  = Pt(8)
        lp.font.bold  = True
        lp.font.color.rgb = lc
        for extra_p in ltf.paragraphs[1:]:
            extra_p.font.name  = "Arial"
            extra_p.font.size  = Pt(8)
            extra_p.font.color.rgb = DARK_GRAY

    # Axis dividers
    v_div = slide.shapes.add_connector(
        2, grid_left + half_w, grid_top, grid_left + half_w, grid_top + grid_h)
    v_div.line.width = Pt(2)
    v_div.line.color.rgb = MID_GRAY

    h_div = slide.shapes.add_connector(
        1, grid_left, grid_top + half_h, grid_left + grid_w, grid_top + half_h)
    h_div.line.width = Pt(2)
    h_div.line.color.rgb = MID_GRAY

    # Axis labels
    # Y-axis: HIGH RISK (top) / LOW RISK (bottom)
    add_axis_label(slide, LEFT_MARGIN, grid_top,
                   axis_label_w, Inches(0.30), "HIGH RISK", RED)
    add_axis_label(slide, LEFT_MARGIN, grid_top + half_h,
                   axis_label_w, Inches(0.30), "LOW RISK", GREEN)
    # X-axis: CERTAIN (left) / UNCERTAIN (right)
    add_axis_label(slide, grid_left, grid_top + grid_h + Inches(0.04),
                   half_w, axis_label_h, "CERTAIN", DARK_GRAY)
    add_axis_label(slide, grid_left + half_w, grid_top + grid_h + Inches(0.04),
                   half_w, axis_label_h, "UNCERTAIN", RED)

    # Items: (text, col, row, color)
    items = [
        ("I assume other roles know what everyone else is doing in the process",                               1, 0, RED),
        ("Do teams know who to ask to accomplish a specific task?",                                            1, 0, RED),
        ("Are there better tools we should be using to communicate?",                                         1, 0, RED),
        ("Everyone is notified when to check off or sign off — high trust assumed",                           1, 0, RED),
        ("High uncertainty when banner rotations are not discussed a month in advance",                        1, 0, RED),
        ("Clear Communication & Role Ownership is muddy — high risk right now",                               1, 0, RED),
        ("Can we add an additional level of Lilly QA before banners go to agency?",                           1, 1, BLUE),
        ("Creative teams have direct line of sight into approval process to anticipate grid updates",          1, 1, BLUE),
    ]

    card_h   = Inches(0.52)
    card_gap = Inches(0.05)
    quad_counts = {}  # track next y per quadrant

    for text, col, row, color in items:
        key = (col, row)
        n   = quad_counts.get(key, 0)
        qx  = grid_left + col * half_w
        qy  = grid_top  + row * half_h
        card_x = qx + Inches(0.1)
        card_y = qy + Inches(0.50) + n * (card_h + card_gap)
        card_w = half_w - Inches(0.2)
        if card_y + card_h > qy + half_h - Inches(0.05):
            continue
        add_card(slide, card_x, card_y, card_w, card_h, text, color,
                 font_size=Pt(8.5))
        quad_counts[key] = n + 1

    out = '/Users/V5X8512/Downloads/07_assumptions.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 8: Big Idea Vignettes ───────────────────────────────────────────

def create_big_ideas():
    """OUTPUT 8: Big Idea Vignettes — grouped by theme cluster."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Big Ideas",
               subtitle="Clustered themes from ideation — ranked by group energy")

    ct = content_top(has_subtitle=True)

    clusters = [
        {
            "theme": "ROLE CLARITY & PROCESS",
            "color": RED,
            "ideas": [
                ("Standardization",    "Standardize process across lines of business and business units"),
                ("Role Clarity",       "Define role clarity for everyone involved and expectations of that role"),
                ("Clear Ownership",    "Assign clear ownership of Creative Grid process within Lilly"),
                ("Defined Process",    "Develop clear end-to-end process with expectations at each lifecycle stage"),
            ],
        },
        {
            "theme": "COMMUNICATION & VISIBILITY",
            "color": BLUE,
            "ideas": [
                ("Visibility",         "ADCO/MCO visibility on handoffs so they can jump in and manage speed"),
                ("More Pings",         "Post to Teams Channel tagging marketer/CMI when handoffs are ready"),
                ("Better Notifications", "Automatic notification when a step completes; visible owner field per stage"),
                ("Rotation Clarity",   "Make currently live banner rotations clear on the creative grid"),
            ],
        },
        {
            "theme": "SYSTEMS & AUTOMATION",
            "color": GREEN,
            "ideas": [
                ("System Alignment",   "Better systems that notify the right people when an asset is ready"),
                ("Automate Grid",      "Auto-populate repeated fields, add drop-downs, enable auto-notifications"),
                ("Workfront Tasks",    "Implement trafficking tasks in Adobe Workfront with assigned POCs"),
                ("Metadata Automation","Collect all metadata for creative campaign and compile into trafficking form"),
            ],
        },
        {
            "theme": "SPEED & QUALITY",
            "color": ORANGE,
            "ideas": [
                ("Fast Track",         "Build a fast-track path for high-urgency messages with pre-approved templates"),
                ("QA Improvements",    "Revamp QA process before assets sent to agency"),
                ("Timeline Standards", "Establish clear end-to-end DTC production timeframes with consistent approval timing"),
                ("Marketer Cheat Sheet", "A marketer cheat sheet: what steps must I complete every time?"),
            ],
        },
    ]

    n_clusters  = len(clusters)
    n_ideas     = max(len(c["ideas"]) for c in clusters)
    col_gap     = Inches(0.10)
    row_gap     = Inches(0.08)
    usable_w    = SLIDE_W - Inches(1.0)
    col_w       = (usable_w - col_gap * (n_ideas - 1)) / n_ideas
    usable_h    = SLIDE_H - ct - Inches(0.1)
    label_h     = Inches(0.28)
    card_h      = (usable_h / n_clusters) - label_h - row_gap

    for ri, cluster in enumerate(clusters):
        row_top = ct + ri * (label_h + card_h + row_gap)
        # Cluster label spanning full width
        add_section_label(slide, LEFT_MARGIN, row_top,
                          usable_w, cluster["theme"], cluster["color"], size=Pt(11))
        for ci, (headline, desc) in enumerate(cluster["ideas"]):
            cx = LEFT_MARGIN + ci * (col_w + col_gap)
            add_card(slide, cx, row_top + label_h, col_w, card_h,
                     desc, cluster["color"],
                     title=headline, font_size=Pt(8.5))

    out = '/Users/V5X8512/Downloads/08_big_ideas.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 9: Prioritization Grid ─────────────────────────────────────────

def create_prioritization_grid():
    """OUTPUT 9: Prioritization Grid 2×2 (Importance × Feasibility)."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Prioritization Grid",
               subtitle="Importance to user (y-axis)  ·  Ease / feasibility (x-axis)")

    ct = content_top(has_subtitle=True)

    axis_label_w = Inches(0.85)
    axis_label_h = Inches(0.30)
    grid_left    = LEFT_MARGIN + axis_label_w + Inches(0.1)
    grid_top     = ct
    grid_w       = SLIDE_W - grid_left - Inches(0.4)
    grid_h       = SLIDE_H - ct - Inches(0.45)
    half_w       = grid_w / 2
    half_h       = grid_h / 2

    quad_styles = [
        (0, 0, RGBColor(0xE8, 0xF0, 0xFB), "HIGH IMPORTANCE / DIFFICULT\nBIG BETS",     BLUE),
        (1, 0, RGBColor(0xED, 0xF4, 0xED), "HIGH IMPORTANCE / EASY\nNO-BRAINERS ★",    GREEN),
        (0, 1, RGBColor(0xF5, 0xF5, 0xF5), "LOW IMPORTANCE / DIFFICULT\nUNWISE",        DARK_GRAY),
        (1, 1, RGBColor(0xF0, 0xF6, 0xFF), "LOW IMPORTANCE / EASY\nUTILITIES",          TEAL),
    ]
    for col, row, fill, label, lc in quad_styles:
        qx = grid_left + col * half_w
        qy = grid_top  + row * half_h
        qbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, qx, qy, half_w, half_h)
        qbox.fill.solid()
        qbox.fill.fore_color.rgb = fill
        qbox.line.color.rgb = MID_GRAY
        qbox.line.width = Pt(0.5)
        ltb = slide.shapes.add_textbox(
            qx + Inches(0.1), qy + Inches(0.06), half_w - Inches(0.2), Inches(0.42))
        ltf = ltb.text_frame
        ltf.text = label
        lp = ltf.paragraphs[0]
        lp.font.name  = "Arial"
        lp.font.size  = Pt(8.5)
        lp.font.bold  = True
        lp.font.color.rgb = lc
        for extra_p in ltf.paragraphs[1:]:
            extra_p.font.name  = "Arial"
            extra_p.font.size  = Pt(8.5)
            extra_p.font.bold  = True
            extra_p.font.color.rgb = lc

    v_div = slide.shapes.add_connector(
        2, grid_left + half_w, grid_top, grid_left + half_w, grid_top + grid_h)
    v_div.line.width = Pt(2)
    v_div.line.color.rgb = MID_GRAY
    h_div = slide.shapes.add_connector(
        1, grid_left, grid_top + half_h, grid_left + grid_w, grid_top + half_h)
    h_div.line.width = Pt(2)
    h_div.line.color.rgb = MID_GRAY

    add_axis_label(slide, LEFT_MARGIN, grid_top,
                   axis_label_w, Inches(0.30), "HIGH", NEAR_BLACK)
    add_axis_label(slide, LEFT_MARGIN, grid_top + half_h,
                   axis_label_w, Inches(0.30), "LOW", DARK_GRAY)
    add_axis_label(slide, grid_left, grid_top + grid_h + Inches(0.04),
                   half_w, axis_label_h, "DIFFICULT / EXPENSIVE", DARK_GRAY)
    add_axis_label(slide, grid_left + half_w, grid_top + grid_h + Inches(0.04),
                   half_w, axis_label_h, "EASY / CHEAP", GREEN)

    # Items: (text, col, row, color)
    ideas_to_plot = [
        ("Automate Grid Notifications",     1, 0, GREEN),
        ("Teams Channel Pings",             1, 0, GREEN),
        ("Role Clarity Documentation (RACI)", 1, 0, GREEN),
        ("Marketer Cheat Sheet",            1, 0, GREEN),
        ("Status Tracker",                  1, 0, GREEN),
        ("Quarterly Rotation Matrix",       1, 0, GREEN),
        ("Metadata Auto-Population",        0, 0, BLUE),
        ("Workfront Implementation",        0, 0, BLUE),
        ("Complete Grid Redesign",          0, 0, BLUE),
        ("Cross-BU Standardization",        0, 0, BLUE),
        ("Additional QA Layer",             1, 1, TEAL),
        ("High-Level Creative Map",         1, 1, TEAL),
    ]

    card_h   = Inches(0.48)
    card_gap = Inches(0.05)
    quad_counts = {}

    for text, col, row, color in ideas_to_plot:
        key = (col, row)
        n   = quad_counts.get(key, 0)
        qx  = grid_left + col * half_w
        qy  = grid_top  + row * half_h
        cx  = qx + Inches(0.1)
        cy  = qy + Inches(0.52) + n * (card_h + card_gap)
        cw  = half_w - Inches(0.2)
        if cy + card_h > qy + half_h - Inches(0.04):
            continue
        add_card(slide, cx, cy, cw, card_h, text, color, font_size=Pt(9))
        quad_counts[key] = n + 1

    out = '/Users/V5X8512/Downloads/09_prioritization_grid.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


def add_axis_label(slide, left, top, width, height, text, color=NEAR_BLACK, size=Pt(9)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name  = "Arial"
    p.font.size  = size
    p.font.bold  = True
    p.font.color.rgb = color


# ── Artifact 10: To-Be Process Map ──────────────────────────────────────────

def create_tobe_process_map():
    """OUTPUT 10: To-Be Scenario Map — same swimlane layout as As-Is."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "To-Be Scenario Map",
               subtitle="Banner Display Trafficking & Handoff Process — Future State")

    ct = content_top(has_subtitle=True)

    phases = [
        "Planning &\nWorkflow",
        "Creative &\nProduction",
        "Review &\nApproval",
        "Trafficking &\nActivation",
        "Confirmation",
    ]
    rows        = ["DOING", "THINKING", "FEELING"]
    row_colors  = [GREEN, TEAL, BLUE]

    label_col_w = Inches(1.1)
    avail_w     = SLIDE_W - LEFT_MARGIN - label_col_w - Inches(0.5)
    col_w       = avail_w / len(phases)
    phase_h     = Inches(0.52)
    avail_row_h = SLIDE_H - ct - phase_h - Inches(0.12)
    row_h       = avail_row_h / len(rows)
    table_left  = LEFT_MARGIN + label_col_w

    # Phase headers (teal — future state)
    for ci, phase in enumerate(phases):
        px = table_left + ci * col_w
        ph_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px, ct, col_w - Inches(0.04), phase_h)
        ph_box.fill.solid()
        ph_box.fill.fore_color.rgb = TEAL
        ph_box.line.fill.background()
        ptf = ph_box.text_frame
        ptf.text = phase
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.word_wrap = True
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pp.font.name  = "Arial"
        pp.font.size  = Pt(9)
        pp.font.bold  = True
        pp.font.color.rgb = WHITE
        for ep in ptf.paragraphs[1:]:
            ep.alignment = PP_ALIGN.CENTER
            ep.font.name  = "Arial"
            ep.font.size  = Pt(9)
            ep.font.bold  = True
            ep.font.color.rgb = WHITE

    # Row labels
    for ri, (row_label, rc) in enumerate(zip(rows, row_colors)):
        ry = ct + phase_h + ri * row_h
        rl_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, ry, label_col_w - Inches(0.04), row_h - Inches(0.04))
        rl_box.fill.solid()
        rl_box.fill.fore_color.rgb = LIGHT_GRAY
        rl_box.line.color.rgb = MID_GRAY
        rl_box.line.width = Pt(0.5)
        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, ry, ACCENT_W, row_h - Inches(0.04))
        acc.fill.solid()
        acc.fill.fore_color.rgb = rc
        acc.line.fill.background()
        rltb = slide.shapes.add_textbox(
            LEFT_MARGIN + ACCENT_W + Inches(0.06), ry + Inches(0.08),
            label_col_w - ACCENT_W - Inches(0.1), row_h - Inches(0.16))
        rltf = rltb.text_frame
        rltf.text = row_label
        rlp = rltf.paragraphs[0]
        rlp.font.name  = "Arial"
        rlp.font.size  = Pt(9)
        rlp.font.bold  = True
        rlp.font.color.rgb = rc

    cell_data = {
        "DOING": [
            "Quarterly matrix created\nRoles defined\nBacklog aligned",
            "Auto-populated templates\nAsset build\nMetadata captured automatically",
            "Streamlined approval\nEditor → Owner → AFD",
            "Auto-notification sent\nGrid pre-filled\nAgency receives complete assets",
            "Confirmation sent to marketer\nBanner live + tracked",
        ],
        "THINKING": [
            "Clear visibility on rotation plan",
            "Templates standardized across all BUs",
            "Faster approval process with defined path",
            "I know exactly who to contact — no guessing",
            "Process complete, confident it's live",
        ],
        "FEELING": [
            "Prepared and aligned from the start",
            "Confident and efficient",
            "Clear expectations throughout",
            "In control — communicated well",
            "Accomplished and confident",
        ],
    }

    for ri, (row_label, rc) in enumerate(zip(rows, row_colors)):
        ry = ct + phase_h + ri * row_h
        for ci, text in enumerate(cell_data[row_label]):
            cx = table_left + ci * col_w
            add_card(slide, cx + Inches(0.02), ry + Inches(0.02),
                     col_w - Inches(0.06), row_h - Inches(0.06),
                     text, rc, font_size=Pt(8))

    # Legend
    leg = slide.shapes.add_textbox(
        SLIDE_W - Inches(3.8), SLIDE_H - Inches(0.32),
        Inches(3.6), Inches(0.25))
    legtf = leg.text_frame
    legtf.text = "Teal headers = future state indicator"
    legp = legtf.paragraphs[0]
    legp.font.name   = "Calibri"
    legp.font.size   = Pt(8)
    legp.font.italic = True
    legp.font.color.rgb = TEAL

    out = '/Users/V5X8512/Downloads/10_tobe_map.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Entry point ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating Design Thinking Workshop Artifacts 7–10...")
    print("=" * 60)
    create_assumptions_grid()
    create_big_ideas()
    create_prioritization_grid()
    create_tobe_process_map()
    print("=" * 60)
    print("Completed artifacts 7–10.")
