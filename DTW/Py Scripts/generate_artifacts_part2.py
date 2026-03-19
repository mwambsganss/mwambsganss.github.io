#!/usr/bin/env python3
"""
Design Thinking Workshop Artifacts Generator — Parts 4–6
Empathy Maps, As-Is Process Map, Needs Statements

Same design system as generate_artifacts.py — see that file for full notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Lilly Brand Colors ──────────────────────────────────────────────────────
RED     = RGBColor(0xD5, 0x2B, 0x1E)
BLUE    = RGBColor(0x00, 0x63, 0xBE)
GREEN   = RGBColor(0x00, 0x7A, 0x33)
PURPLE  = RGBColor(0x6D, 0x20, 0x77)
ORANGE  = RGBColor(0xE8, 0x77, 0x22)
TEAL    = RGBColor(0x00, 0xA3, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x58, 0x58, 0x58)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xC8, 0xC8, 0xC8)

ACCENT_COLORS = [RED, BLUE, GREEN, PURPLE, ORANGE, TEAL]

SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.50)
HEADER_H    = Inches(0.75)
DIVIDER_H   = Inches(0.04)
LEFT_MARGIN = Inches(0.5)
ACCENT_W    = Inches(0.06)
CARD_PADDING = Inches(0.12)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, prs, title, subtitle=""):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.12),
                                  prs.slide_width - Inches(1), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = Pt(28)
    p.font.bold  = True
    p.font.color.rgb = WHITE

    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, HEADER_H, prs.slide_width, DIVIDER_H)
    div.fill.solid()
    div.fill.fore_color.rgb = MID_GRAY
    div.line.fill.background()

    if subtitle:
        stb = slide.shapes.add_textbox(
            LEFT_MARGIN, HEADER_H + DIVIDER_H + Inches(0.06),
            prs.slide_width - Inches(1), Inches(0.32))
        stf = stb.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.name   = "Calibri"
        sp.font.size   = Pt(11)
        sp.font.italic = True
        sp.font.color.rgb = DARK_GRAY


def add_card(slide, left, top, width, height, text, accent_color,
             title=None, font_size=Pt(10), bg_color=None):
    bg = bg_color or WHITE
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(0.5)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, ACCENT_W, height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    text_left  = left + ACCENT_W + CARD_PADDING
    text_width = width - ACCENT_W - CARD_PADDING - Inches(0.08)
    tb = slide.shapes.add_textbox(
        text_left, top + Inches(0.06), text_width, height - Inches(0.12))
    tf = tb.text_frame
    tf.word_wrap = True

    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name  = "Arial"
        p0.font.size  = Pt(9)
        p0.font.bold  = True
        p0.font.color.rgb = NEAR_BLACK
        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.name  = "Calibri"
        p1.font.size  = font_size
        p1.font.color.rgb = DARK_GRAY
    else:
        p0 = tf.paragraphs[0]
        p0.text = text
        p0.font.name  = "Calibri"
        p0.font.size  = font_size
        p0.font.color.rgb = NEAR_BLACK


def add_section_label(slide, left, top, width, text, color=None):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name  = "Arial"
    p.font.size  = Pt(10)
    p.font.bold  = True
    p.font.color.rgb = color or NEAR_BLACK


def content_top(has_subtitle=False):
    extra = Inches(0.38) if has_subtitle else Inches(0.14)
    return HEADER_H + DIVIDER_H + extra


# ── Artifact 4: Empathy Maps ────────────────────────────────────────────────

def create_empathy_maps():
    """OUTPUT 4: Empathy Maps — one slide per persona, 2×2 quadrant layout."""
    prs = new_prs()

    personas = [
        {
            "name": "Mia — HCP Marketer",
            "thinks": ["Have to be proactive 1–2 months out to ensure right message gets to market",
                       "Think through the HCP experience and what I want them to see",
                       "Give CMI partners ample heads-up when we have a plan"],
            "feels":  ["Frustrated by the number of handoffs in this new process",
                       "Letting the process down when I'm not fast enough entering inputs",
                       "Loss of control once it reaches the creative grid process"],
            "says":   ["I kick off banners, project manage through MLR and creative grid",
                       "This can take 5–10 minutes sometimes — even that is hard to find time for",
                       "Banners left hanging if I think we're set but CMI doesn't get notification"],
            "does":   ["Kicks off banners with brand strategy and LMS/E2E colleagues",
                       "Takes through MLR to get approval",
                       "Sends CMI updated grid with banner rotation, start/end dates",
                       "Checks in via email, Teams, or NPP weekly status"],
            "pains":  ["Too many handoffs create friction and delays",
                       "5–10 min grid entry hard to find time for",
                       "Loss of control in creative grid process"],
            "gains":  ["Right message reaches HCPs on time",
                       "Feel confident process is complete",
                       "Know when steps are incomplete so team can fix quickly"],
            "color": RED,
        },
        {
            "name": "Julie — DTC Media Lead",
            "thinks": ["Constant moving pieces to coordinate across teams",
                       "Consider full HCP experience and which placements resonate most",
                       "Try to give CMI as much advance notice as possible"],
            "feels":  ["Number of handoffs can feel heavy and frustrating",
                       "Loss of control — communication gaps can stall banners",
                       "Responsible for protecting marketplace investment"],
            "says":   ["We need to understand impact when assets are not available",
                       "Delayed assets dilutes our marketplace impact to drive scripts",
                       "Placements must be strategically placed for patients/consumers"],
            "does":   ["Strategically determines placements for ad unit types",
                       "Provides specs to TA & GCO CC&A team",
                       "Confirms creative grid inputs are accurate before agency handoff"],
            "pains":  ["Constant moving pieces", "Items delayed negatively impact plans",
                       "Dilutes marketplace impact"],
            "gains":  ["Assets delivered on time", "Strategic integrity of plans maintained",
                       "Patients/consumers see intended messages"],
            "color": BLUE,
        },
        {
            "name": "Molly — Media Operations",
            "thinks": ["So many roles constantly changing — impossible to keep up",
                       "Black boxes — people don't understand who does what",
                       "When problems arise, hard to know who to triage"],
            "feels":  ["Impossible to keep people communicating on what the process is",
                       "Frustrated by role changes making documentation quickly stale"],
            "says":   ["Ensures trafficking handoff processes are working well",
                       "Providing necessary tools like Smartsheet",
                       "Process documentation so groups can confidently work together"],
            "does":   ["Provides Smartsheet tools and process documentation",
                       "Facilitates communication between groups",
                       "Enables teams to work together confidently"],
            "pains":  ["Roles constantly changing", "Black boxes — unclear who does what",
                       "Problems arise and it's unclear who to call"],
            "gains":  ["Understand who is doing what",
                       "Triage who to talk to when problems arise",
                       "Dig into root cause of failures"],
            "color": GREEN,
        },
    ]

    for persona in personas:
        slide = blank_slide(prs)
        add_header(slide, prs,
                   f"Empathy Map — {persona['name']}",
                   subtitle=f"Thinks · Feels · Says · Does · Pains · Gains")

        ct = content_top(has_subtitle=True)
        avail_h = SLIDE_H - ct - Inches(0.12)

        # Layout: 2 columns × 2 rows for quadrants, then pains/gains strip
        # Pains/Gains strip height
        strip_h  = Inches(0.85)
        quad_h   = (avail_h - strip_h - Inches(0.12)) / 2
        half_w   = (SLIDE_W - Inches(1.0)) / 2 - Inches(0.05)

        quadrants = [
            ("THINKS",  persona["thinks"], BLUE,   LEFT_MARGIN,              ct),
            ("FEELS",   persona["feels"],  PURPLE, LEFT_MARGIN + half_w + Inches(0.1), ct),
            ("SAYS",    persona["says"],   ORANGE, LEFT_MARGIN,              ct + quad_h + Inches(0.06)),
            ("DOES",    persona["does"],   TEAL,   LEFT_MARGIN + half_w + Inches(0.1), ct + quad_h + Inches(0.06)),
        ]

        card_h = Inches(0.54)
        gap    = Inches(0.05)

        for label, items, color, qx, qy in quadrants:
            add_section_label(slide, qx, qy, Inches(3), label, color)
            iy = qy + Inches(0.30)
            for item in items[:4]:
                if iy + card_h > qy + quad_h - Inches(0.02):
                    break
                add_card(slide, qx, iy, half_w, card_h, item, color,
                         font_size=Pt(9))
                iy += card_h + gap

        # Pains / Gains strip at bottom
        strip_y    = SLIDE_H - strip_h - Inches(0.08)
        pains_w    = half_w
        gains_w    = half_w
        pains_x    = LEFT_MARGIN
        gains_x    = LEFT_MARGIN + half_w + Inches(0.1)
        item_card_h = Inches(0.55)
        item_card_w = half_w / max(1, len(persona["pains"][:3])) - Inches(0.04)

        add_section_label(slide, pains_x, strip_y, Inches(2), "PAINS", RED)
        add_section_label(slide, gains_x, strip_y, Inches(2), "GAINS", GREEN)

        for i, pain in enumerate(persona["pains"][:3]):
            px = pains_x + i * (item_card_w + Inches(0.04))
            add_card(slide, px, strip_y + Inches(0.28), item_card_w, item_card_h,
                     pain, RED, font_size=Pt(8))

        for i, gain in enumerate(persona["gains"][:3]):
            gx = gains_x + i * (item_card_w + Inches(0.04))
            add_card(slide, gx, strip_y + Inches(0.28), item_card_w, item_card_h,
                     gain, GREEN, font_size=Pt(8))

    out = '/Users/V5X8512/Downloads/04_empathy_maps.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 5: As-Is Process Map ───────────────────────────────────────────

def create_asis_process_map():
    """OUTPUT 5: As-Is Scenario Map — swimlane with phase columns."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs,
               "As-Is Scenario Map",
               subtitle="Banner Display Trafficking & Handoff Process")

    ct = content_top(has_subtitle=True)

    phases = [
        "Planning &\nWorkflow",
        "Creative &\nProduction",
        "Review &\nApproval",
        "Trafficking &\nActivation",
        "Tactic\nDestruction",
    ]
    rows = ["DOING", "THINKING", "FEELING"]
    row_colors = [BLUE, TEAL, ORANGE]

    # Dimensions
    label_col_w = Inches(1.1)
    avail_w     = SLIDE_W - LEFT_MARGIN - label_col_w - Inches(0.5)
    col_w       = avail_w / len(phases)
    phase_h     = Inches(0.52)
    avail_row_h = SLIDE_H - ct - phase_h - Inches(0.12)
    row_h       = avail_row_h / len(rows)
    table_left  = LEFT_MARGIN + label_col_w

    # Phase header row
    for ci, phase in enumerate(phases):
        px = table_left + ci * col_w
        ph_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px, ct, col_w - Inches(0.04), phase_h)
        ph_box.fill.solid()
        ph_box.fill.fore_color.rgb = RED
        ph_box.line.fill.background()
        ptf = ph_box.text_frame
        ptf.text = phase
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.word_wrap = True
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pp.font.name  = "Arial"
        pp.font.size  = Pt(9)
        pp.font.bold  = True
        pp.font.color.rgb = WHITE
        for extra_p in ptf.paragraphs[1:]:
            extra_p.alignment = PP_ALIGN.CENTER
            extra_p.font.name  = "Arial"
            extra_p.font.size  = Pt(9)
            extra_p.font.bold  = True
            extra_p.font.color.rgb = WHITE

    # Row label column
    for ri, (row_label, rc) in enumerate(zip(rows, row_colors)):
        ry = ct + phase_h + ri * row_h
        rl_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, ry, label_col_w - Inches(0.04), row_h - Inches(0.04))
        rl_box.fill.solid()
        rl_box.fill.fore_color.rgb = LIGHT_GRAY
        rl_box.line.color.rgb = MID_GRAY
        rl_box.line.width = Pt(0.5)

        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, LEFT_MARGIN, ry, ACCENT_W, row_h - Inches(0.04))
        acc.fill.solid()
        acc.fill.fore_color.rgb = rc
        acc.line.fill.background()

        rltb = slide.shapes.add_textbox(
            LEFT_MARGIN + ACCENT_W + Inches(0.06), ry + Inches(0.08),
            label_col_w - ACCENT_W - Inches(0.1), row_h - Inches(0.16))
        rltf = rltb.text_frame
        rltf.text = row_label
        rlp = rltf.paragraphs[0]
        rlp.font.name  = "Arial"
        rlp.font.size  = Pt(9)
        rlp.font.bold  = True
        rlp.font.color.rgb = rc

    # Cell content
    cell_data = {
        "DOING": [
            "Quarterly alignment\nBacklog refinement\nSprint planning",
            "Asset build per approved templates\nReadiness checks",
            "Editor → MLRO → Project Owner → AFD",
            "Assets + grid details sent to agencies\nQA assets\nBuild tags",
            "Upon tactic expiration",
        ],
        "THINKING": [
            "Need to align with brand strategy",
            "Must follow templates and specs exactly",
            "Multiple approval layers needed",
            "Who has the files? Are they correct?",
            "Is the banner out of market?",
        ],
        "FEELING": [
            "Proactive planning required",
            "Rushing to meet deadlines",
            "Waiting — anxious about approvals",
            "Frustrated by handoffs\nLoss of control\nCommunication gaps",
            "Relief when complete",
        ],
    }
    pain_indices = {("FEELING", 3)}  # flag pain point

    for ri, (row_label, rc) in enumerate(zip(rows, row_colors)):
        ry = ct + phase_h + ri * row_h
        for ci, text in enumerate(cell_data[row_label]):
            cx = table_left + ci * col_w
            is_pain = (row_label, ci) in pain_indices
            cell_bg = RGBColor(0xFF, 0xF0, 0xED) if is_pain else WHITE
            add_card(slide, cx + Inches(0.02), ry + Inches(0.02),
                     col_w - Inches(0.06), row_h - Inches(0.06),
                     text, rc if not is_pain else RED,
                     font_size=Pt(8), bg_color=cell_bg)

    # Pain point legend
    leg = slide.shapes.add_textbox(
        SLIDE_W - Inches(3.2), SLIDE_H - Inches(0.32),
        Inches(3.0), Inches(0.25))
    legtf = leg.text_frame
    legtf.text = "Shaded = key pain point"
    legp = legtf.paragraphs[0]
    legp.font.name  = "Calibri"
    legp.font.size  = Pt(8)
    legp.font.italic = True
    legp.font.color.rgb = RED

    out = '/Users/V5X8512/Downloads/05_asis_map.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Artifact 6: Needs Statements ────────────────────────────────────────────

def create_needs_statements():
    """OUTPUT 6: Needs Statements — card grid grouped by theme."""
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs,
               "Needs Statements",
               subtitle="[User] needs a way to... so that...")

    ct = content_top(has_subtitle=True)

    needs = [
        # (persona, need, benefit, accent_color)
        ("Mia (HCP Marketer)",
         "easily see rotation and banner status",
         "I can quickly assess HCP experience and understand if changes are needed",
         RED),
        ("Mia (HCP Marketer)",
         "know when banners went to customers",
         "I feel confident the process is complete and HCPs are seeing intended messages",
         RED),
        ("Mia (HCP Marketer)",
         "know when steps are incomplete",
         "myself or team can address it for speedy launches",
         RED),
        ("Julie (DTC Media)",
         "ensure assets are delivered on time",
         "to maintain strategic integrity of plans, impact patients, and protect agency resources",
         BLUE),
        ("Molly (Media Ops)",
         "understand who is doing what",
         "when problems arise, we can triage who to talk to and dig into what went wrong",
         GREEN),
        ("Matt (Digital Brand Lead)",
         "know when assets are to be trafficked with all information needed",
         "to ensure accuracy of the handoff",
         PURPLE),
        ("Matt (Digital Brand Lead)",
         "a process or confirmation when the banner is out of market",
         "to ensure we are compliant with tactic expiration",
         PURPLE),
        ("HCP Agency (CMI)",
         "clear communication, role ownership, and rotation direction from the creative grid",
         "to avoid delays and mitigate risk for errors",
         ORANGE),
        ("HCP Agency (CMI)",
         "banner assets to arrive without errors",
         "to avoid delays in getting assets into market",
         ORANGE),
        ("Consumer Agency (PA)",
         "accurate creative grid updates and asset delivery",
         "to adhere to SLA and allow priority campaigns to break timeline when approvals are delayed",
         TEAL),
        ("ADCO",
         "know who to connect to when files or handoffs are not complete",
         "we can ensure speed of the work",
         RED),
        ("ADCO",
         "confirmation once files are received and QA has started",
         "we can track progress and manage timeline expectations",
         RED),
    ]

    # 3-column grid
    n_cols  = 3
    col_gap = Inches(0.1)
    col_w   = (SLIDE_W - Inches(1.0) - col_gap * (n_cols - 1)) / n_cols
    card_h  = Inches(0.90)
    row_gap = Inches(0.08)

    for idx, (persona, need, benefit, color) in enumerate(needs):
        ci = idx % n_cols
        ri = idx // n_cols
        x  = LEFT_MARGIN + ci * (col_w + col_gap)
        y  = ct + ri * (card_h + row_gap)
        if y + card_h > SLIDE_H - Inches(0.08):
            break
        full_text = f"{persona} needs a way to {need} so that {benefit}."
        add_card(slide, x, y, col_w, card_h,
                 full_text, color,
                 title=persona, font_size=Pt(8.5))

    out = '/Users/V5X8512/Downloads/06_needs_statements.pptx'
    prs.save(out)
    print(f"✓ Generated: {out}")


# ── Entry point ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating Design Thinking Workshop Artifacts 4–6...")
    print("=" * 60)
    create_empathy_maps()
    create_asis_process_map()
    create_needs_statements()
    print("=" * 60)
    print("Completed artifacts 4–6.")
