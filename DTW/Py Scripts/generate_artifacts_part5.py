#!/usr/bin/env python3
"""
Design Thinking Workshop Artifacts Generator - Part 5
Feedback Grid, Executive Summary, Workshop Playback Deck
Follows Cloud DTW template design system (Lilly brand).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Design system constants (matches generate_artifacts parts 1-4)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
HEADER_H = Inches(0.75)
DIVIDER_H = Inches(0.04)
LEFT_MARGIN = Inches(0.5)
RIGHT_MARGIN = Inches(0.5)
ACCENT_W = Inches(0.06)
CARD_PADDING = Inches(0.12)

RED    = RGBColor(0xD5, 0x2B, 0x1E)
BLUE   = RGBColor(0x00, 0x63, 0xBE)
GREEN  = RGBColor(0x00, 0x7A, 0x33)
PURPLE = RGBColor(0x6D, 0x20, 0x77)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
TEAL   = RGBColor(0x00, 0xA3, 0xE0)
GRAY   = RGBColor(0x58, 0x59, 0x5B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# Tinted/pastel versions for backgrounds
RED_LIGHT    = RGBColor(0xFA, 0xD4, 0xD2)
BLUE_LIGHT   = RGBColor(0xCC, 0xDF, 0xF3)
GREEN_LIGHT  = RGBColor(0xCC, 0xE8, 0xD7)
TEAL_LIGHT   = RGBColor(0xCC, 0xED, 0xF7)
ORANGE_LIGHT = RGBColor(0xFA, 0xE5, 0xCC)
PURPLE_LIGHT = RGBColor(0xE1, 0xCC, 0xE3)


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _solid_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _tf_para(tf, text, font_name, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             space_before=0, italic=False):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    p.text = ""
    r = run
    r.font.name  = font_name
    r.font.size  = font_size
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if space_before:
        p.space_before = Pt(space_before)
    return p


def _textbox(slide, left, top, width, height, text, font_name, font_size,
             bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_header(slide, prs, title, subtitle=""):
    """Full-width red header bar with white title, optional subtitle, and gray divider."""
    _solid_rect(slide, 0, 0, prs.slide_width, HEADER_H, RED)
    # Title text
    tb = slide.shapes.add_textbox(LEFT_MARGIN, 0, prs.slide_width - LEFT_MARGIN * 2, HEADER_H)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_top  = 0
    from pptx.util import Pt as _Pt
    tf.auto_size   = None
    p = tf.paragraphs[0]
    p.alignment    = PP_ALIGN.LEFT
    run = p.add_run()
    run.text       = title
    run.font.name  = "Arial"
    run.font.size  = Pt(22) if not subtitle else Pt(20)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name  = "Arial"
        r2.font.size  = Pt(11)
        r2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    # Divider
    _solid_rect(slide, 0, HEADER_H, prs.slide_width, DIVIDER_H, GRAY)


def content_top(has_subtitle=False):
    return HEADER_H + DIVIDER_H + Inches(0.18)


def add_card(slide, left, top, width, height, text, accent_color,
             title=None, font_size=Pt(10), bg_color=None):
    """White card with colored left accent bar. Optional bold title."""
    bg = bg_color or WHITE
    _solid_rect(slide, left, top, width, height, bg, RGBColor(0xD0, 0xD0, 0xD0), 0.5)
    _solid_rect(slide, left, top, ACCENT_W, height, accent_color)
    # Text area
    tx_left  = left + ACCENT_W + CARD_PADDING
    tx_width = width - ACCENT_W - CARD_PADDING * 2
    tx_top   = top + CARD_PADDING
    tx_height = height - CARD_PADDING * 2
    tb = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if title:
        pt = tf.paragraphs[0]
        pt.alignment = PP_ALIGN.LEFT
        rt = pt.add_run()
        rt.text = title
        rt.font.name  = "Arial"
        rt.font.size  = font_size + Pt(1)
        rt.font.bold  = True
        rt.font.color.rgb = accent_color
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name  = "Calibri"
    run.font.size  = font_size
    run.font.color.rgb = RGBColor(0x30, 0x30, 0x30)


def add_section_label(slide, left, top, width, text, color=GRAY):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name  = "Arial"
    run.font.size  = Pt(9)
    run.font.bold  = True
    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# OUTPUT 15: Feedback Grid
# ---------------------------------------------------------------------------

def create_feedback_grid():
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Feedback Grid",
               "What worked, what to change, open questions, and new ideas")

    ct = content_top(has_subtitle=True)
    avail_h = SLIDE_H - ct - Inches(0.2)
    avail_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN

    half_w = avail_w / 2 - Inches(0.05)
    half_h = avail_h / 2 - Inches(0.05)

    quadrants = [
        # (col, row, title, color, accent, bg)
        (0, 0, "+ THINGS THAT WORKED",  GREEN,  GREEN,  GREEN_LIGHT),
        (1, 0, "DELTA  THINGS TO CHANGE", ORANGE, ORANGE, ORANGE_LIGHT),
        (0, 1, "?  QUESTIONS WE STILL HAVE", BLUE,   BLUE,   BLUE_LIGHT),
        (1, 1, "IDEA  NEW IDEAS TO TRY",  TEAL,   TEAL,   TEAL_LIGHT),
    ]

    for col, row, title, accent, label_color, bg in quadrants:
        qx = LEFT_MARGIN + col * (half_w + Inches(0.1))
        qy = ct + row * (half_h + Inches(0.1))
        # Quadrant background
        _solid_rect(slide, qx, qy, half_w, half_h, bg,
                    RGBColor(0xC0, 0xC0, 0xC0), 0.5)
        # Top label band
        _solid_rect(slide, qx, qy, half_w, Inches(0.3), accent)
        # Label text
        tb = slide.shapes.add_textbox(
            qx + Inches(0.1), qy + Inches(0.04), half_w - Inches(0.2), Inches(0.26))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_top = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name  = "Arial"
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = WHITE
        # Placeholder bullet lines
        sample = {
            "+ THINGS THAT WORKED":   ["Workshop format brought all stakeholders together",
                                        "Role-by-role empathy mapping surfaced real pain",
                                        "Transcript captured all perspectives accurately",
                                        "Clear identification of cross-functional pain points"],
            "DELTA  THINGS TO CHANGE": ["More time needed for roadmap prioritization",
                                          "Include digital production team members upfront",
                                          "Schedule follow-up to finalize ownership",
                                          "Identify executive sponsors before workshop"],
            "?  QUESTIONS WE STILL HAVE": ["What is timeline for automated notifications?",
                                             "Who will own the RACI documentation?",
                                             "What budget is available for integrations?",
                                             "How do we align HCP and Consumer processes?"],
            "IDEA  NEW IDEAS TO TRY":  ["AI/ML for predicting banner performance",
                                          "Banner production SLA dashboard",
                                          "Natural language grid input tool",
                                          "Chatbot for 'who owns what' queries"],
        }
        lines = sample.get(title, [])
        item_h = (half_h - Inches(0.3) - Inches(0.1)) / max(len(lines), 1)
        for i, line in enumerate(lines):
            iy = qy + Inches(0.3) + Inches(0.05) + i * item_h
            add_card(slide, qx + Inches(0.1), iy,
                     half_w - Inches(0.2), item_h - Inches(0.04),
                     line, accent, font_size=Pt(9))

    prs.save('/Users/V5X8512/Downloads/15_feedback_grid.pptx')
    print("Generated: 15_feedback_grid.pptx")


# ---------------------------------------------------------------------------
# OUTPUT 16: Executive Summary
# ---------------------------------------------------------------------------

def create_executive_summary():
    prs = new_prs()
    slide = blank_slide(prs)
    add_header(slide, prs, "Workshop Executive Summary",
               "Cloud Technology Modernization — Design Thinking Workshop")

    ct = content_top(has_subtitle=True)
    col_gap  = Inches(0.15)
    col_w    = (SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN - col_gap) / 2
    left1    = LEFT_MARGIN
    left2    = LEFT_MARGIN + col_w + col_gap

    # ---------- META ROW ----------
    meta_h = Inches(0.42)
    meta_items = [
        ("WORKSHOP",      "Cloud Technology Modernization — Design Thinking Session"),
        ("DATE",          "January 2026"),
        ("PARTICIPANTS",  "Cloud Platform, App Dev, Infrastructure, Security, Finance, PMO"),
    ]
    mx = left1
    mw = (SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN - col_gap * 2) / 3
    for label, value in meta_items:
        _solid_rect(slide, mx, ct, mw, meta_h, LIGHT_GRAY,
                    RGBColor(0xD0, 0xD0, 0xD0), 0.5)
        _solid_rect(slide, mx, ct, ACCENT_W, meta_h, BLUE)
        lbl_tb = slide.shapes.add_textbox(
            mx + ACCENT_W + CARD_PADDING, ct + Inches(0.04),
            mw - ACCENT_W - CARD_PADDING, Inches(0.16))
        lp = lbl_tb.text_frame.paragraphs[0]
        lp.add_run().text = label
        lp.runs[0].font.name  = "Arial"
        lp.runs[0].font.size  = Pt(7)
        lp.runs[0].font.bold  = True
        lp.runs[0].font.color.rgb = BLUE
        val_tb = slide.shapes.add_textbox(
            mx + ACCENT_W + CARD_PADDING, ct + Inches(0.20),
            mw - ACCENT_W - CARD_PADDING * 2, Inches(0.2))
        vp = val_tb.text_frame.paragraphs[0]
        vp.add_run().text = value
        vp.runs[0].font.name  = "Calibri"
        vp.runs[0].font.size  = Pt(8)
        vp.runs[0].font.color.rgb = GRAY
        mx += mw + col_gap

    row_top = ct + meta_h + Inches(0.14)
    row_h   = Inches(1.32)
    row_gap = Inches(0.1)

    # ---------- LEFT COLUMN ----------
    sections_left = [
        ("PROBLEM STATEMENT", RED,
         "Understand the current-state cloud adoption journey, identify barriers to modernization, "
         "map stakeholder needs across platform and application teams, and define a prioritized "
         "path forward for cloud-native transformation."),
        ("KEY USERS", BLUE,
         "Cloud Platform Engineers, Application Developers, Infrastructure Teams, "
         "Security & Compliance, Finance / FinOps, Product Owners, PMO"),
        ("TOP PAIN POINTS", ORANGE,
         "1.  Inconsistent cloud governance creates shadow IT and compliance risk\n"
         "2.  Lack of self-service tooling slows developer velocity\n"
         "3.  Cost visibility gaps prevent effective FinOps decisions"),
    ]
    y = row_top
    for label, color, body in sections_left:
        add_section_label(slide, left1, y - Inches(0.22), col_w, label, color)
        add_card(slide, left1, y, col_w, row_h, body, color, font_size=Pt(9))
        y += row_h + Inches(0.3)

    # ---------- RIGHT COLUMN ----------
    sections_right = [
        ("TOP INSIGHTS", TEAL,
         "1.  Developer experience is the #1 accelerator — friction kills adoption\n"
         "2.  Security guardrails must be automated, not manual checkpoints\n"
         "3.  Shared accountability models (CCoE) outperform central control"),
        ("TOP BIG IDEAS", GREEN,
         "1.  Self-service cloud landing zones with pre-approved patterns\n"
         "2.  Automated compliance-as-code integrated into CI/CD pipelines\n"
         "3.  FinOps dashboard with real-time cost allocation by product team"),
        ("AGREED HILLS", PURPLE,
         "Cloud teams deploy production workloads in <1 day with zero manual "
         "security reviews, reducing time-to-market by 75% while maintaining "
         "100% compliance posture."),
    ]
    y = row_top
    for label, color, body in sections_right:
        add_section_label(slide, left2, y - Inches(0.22), col_w, label, color)
        add_card(slide, left2, y, col_w, row_h, body, color, font_size=Pt(9))
        y += row_h + Inches(0.3)

    # ---------- NEXT STEPS STRIP ----------
    strip_top = SLIDE_H - Inches(0.75)
    _solid_rect(slide, 0, strip_top, SLIDE_W, Inches(0.75), LIGHT_GRAY,
                RGBColor(0xD0, 0xD0, 0xD0), 0.5)
    _solid_rect(slide, 0, strip_top, Inches(1.5), Inches(0.75), BLUE)
    tb_label = slide.shapes.add_textbox(
        Inches(0.1), strip_top + Inches(0.22), Inches(1.3), Inches(0.32))
    lp = tb_label.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lp.add_run().text = "NEXT STEPS"
    lp.runs[0].font.name  = "Arial"
    lp.runs[0].font.size  = Pt(9)
    lp.runs[0].font.bold  = True
    lp.runs[0].font.color.rgb = WHITE
    steps = ("1. Synthesize workshop data and present playback  |  "
             "2. Schedule follow-up alignment session  |  "
             "3. Launch RACI documentation effort  |  "
             "4. Pilot self-service cloud landing zone (Quick Win)")
    tb_steps = slide.shapes.add_textbox(
        Inches(1.65), strip_top + Inches(0.18),
        SLIDE_W - Inches(1.75) - RIGHT_MARGIN, Inches(0.4))
    sp = tb_steps.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sp.add_run().text = steps
    sp.runs[0].font.name  = "Calibri"
    sp.runs[0].font.size  = Pt(9)
    sp.runs[0].font.color.rgb = GRAY

    prs.save('/Users/V5X8512/Downloads/16_executive_summary.pptx')
    print("Generated: 16_executive_summary.pptx")


# ---------------------------------------------------------------------------
# OUTPUT 17: Workshop Playback Deck (multi-slide)
# ---------------------------------------------------------------------------

def _cover_slide(prs):
    slide = blank_slide(prs)
    # Full red background
    _solid_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RED)
    # White diagonal accent block (bottom-left triangle feel — simple rect)
    _solid_rect(slide, 0, prs.slide_height - Inches(1.5),
                prs.slide_width, Inches(1.5), RGBColor(0xBF, 0x1F, 0x15))
    # Workshop type label
    tb_type = slide.shapes.add_textbox(
        LEFT_MARGIN, Inches(1.8), prs.slide_width - Inches(2), Inches(0.4))
    tp = tb_type.text_frame.paragraphs[0]
    tp.add_run().text = "DESIGN THINKING WORKSHOP  |  PLAYBACK"
    tp.runs[0].font.name  = "Arial"
    tp.runs[0].font.size  = Pt(12)
    tp.runs[0].font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    tp.runs[0].font.bold  = True
    # Main title
    tb = slide.shapes.add_textbox(
        LEFT_MARGIN, Inches(2.3), prs.slide_width - Inches(2), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.add_run().text = "Cloud Technology\nModernization"
    p.runs[0].font.name  = "Arial"
    p.runs[0].font.size  = Pt(44)
    p.runs[0].font.bold  = True
    p.runs[0].font.color.rgb = WHITE
    # Date / participants
    tb2 = slide.shapes.add_textbox(
        LEFT_MARGIN, Inches(4.5), prs.slide_width - Inches(2), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.add_run().text = "January 2026  |  Cloud Platform, App Dev, Infrastructure, Security, Finance, PMO"
    p2.runs[0].font.name  = "Calibri"
    p2.runs[0].font.size  = Pt(13)
    p2.runs[0].font.color.rgb = RGBColor(0xFF, 0xE0, 0xE0)
    # Lilly wordmark area (bottom-right placeholder)
    tb3 = slide.shapes.add_textbox(
        prs.slide_width - Inches(2.2),
        prs.slide_height - Inches(0.5),
        Inches(2.0), Inches(0.38))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    p3.add_run().text = "Eli Lilly and Company"
    p3.runs[0].font.name  = "Arial"
    p3.runs[0].font.size  = Pt(9)
    p3.runs[0].font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)


def _section_divider_slide(prs, section_num, section_title, color=BLUE):
    slide = blank_slide(prs)
    _solid_rect(slide, 0, 0, prs.slide_width, prs.slide_height, color)
    _solid_rect(slide, 0, prs.slide_height - Inches(0.08),
                prs.slide_width, Inches(0.08), RGBColor(0xFF, 0xFF, 0xFF))
    num_tb = slide.shapes.add_textbox(
        LEFT_MARGIN, Inches(2.2), Inches(1.2), Inches(1.0))
    np_ = num_tb.text_frame.paragraphs[0]
    np_.add_run().text = str(section_num).zfill(2)
    np_.runs[0].font.name  = "Arial"
    np_.runs[0].font.size  = Pt(54)
    np_.runs[0].font.bold  = True
    np_.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # slightly muted large number — just use white
    sec_tb = slide.shapes.add_textbox(
        LEFT_MARGIN, Inches(3.4), prs.slide_width - Inches(2), Inches(1.0))
    sp = sec_tb.text_frame.paragraphs[0]
    sp.add_run().text = section_title
    sp.runs[0].font.name  = "Arial"
    sp.runs[0].font.size  = Pt(32)
    sp.runs[0].font.bold  = True
    sp.runs[0].font.color.rgb = WHITE


def _content_slide(prs, title, subtitle, bullets, accent_color=BLUE):
    slide = blank_slide(prs)
    add_header(slide, prs, title, subtitle)
    ct = content_top(has_subtitle=True)
    avail_h = SLIDE_H - ct - Inches(0.2)
    card_h  = min(Inches(0.55), (avail_h - Inches(0.1) * (len(bullets) - 1)) / max(len(bullets), 1))
    y = ct
    for i, (label, body) in enumerate(bullets):
        # Numbered circle
        num_x = LEFT_MARGIN
        num_y = y + (card_h - Inches(0.36)) / 2
        circ  = slide.shapes.add_shape(9, num_x, num_y, Inches(0.36), Inches(0.36))  # oval=9
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent_color
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        cr.font.name  = "Arial"
        cr.font.size  = Pt(9)
        cr.font.bold  = True
        cr.font.color.rgb = WHITE
        # Card
        card_x = LEFT_MARGIN + Inches(0.46)
        card_w = SLIDE_W - card_x - RIGHT_MARGIN
        if label:
            add_card(slide, card_x, y, card_w, card_h, body,
                     accent_color, title=label, font_size=Pt(9))
        else:
            add_card(slide, card_x, y, card_w, card_h, body,
                     accent_color, font_size=Pt(10))
        y += card_h + Inches(0.1)


def _two_col_slide(prs, title, subtitle, left_items, right_items,
                   left_color=GREEN, right_color=ORANGE,
                   left_label="", right_label=""):
    slide = blank_slide(prs)
    add_header(slide, prs, title, subtitle)
    ct = content_top(has_subtitle=True)
    avail_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
    col_w   = (avail_w - Inches(0.2)) / 2
    avail_h = SLIDE_H - ct - Inches(0.2)

    for col_idx, (items, color, col_label) in enumerate(
            [(left_items, left_color, left_label),
             (right_items, right_color, right_label)]):
        cx = LEFT_MARGIN + col_idx * (col_w + Inches(0.2))
        if col_label:
            add_section_label(slide, cx, ct, col_w, col_label, color)
        card_h = min(Inches(0.7), (avail_h - Inches(0.32) - Inches(0.08) * (len(items) - 1)) / max(len(items), 1))
        cy = ct + (Inches(0.28) if col_label else 0)
        for body in items:
            add_card(slide, cx, cy, col_w, card_h, body, color, font_size=Pt(9))
            cy += card_h + Inches(0.08)


def _hills_slide(prs, hills):
    """Hills slide: numbered hero blocks."""
    slide = blank_slide(prs)
    add_header(slide, prs, "Our Hills — Desired Outcomes",
               "WHO does WHAT so WOW")
    ct = content_top(has_subtitle=True)
    avail_h = SLIDE_H - ct - Inches(0.2)
    hill_h  = (avail_h - Inches(0.1) * (len(hills) - 1)) / max(len(hills), 1)
    colors  = [RED, BLUE, TEAL]
    y = ct
    for i, (who, what, wow) in enumerate(hills):
        color   = colors[i % len(colors)]
        block_x = LEFT_MARGIN
        block_w = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
        # Background
        _solid_rect(slide, block_x, y, block_w, hill_h,
                    [RED_LIGHT, BLUE_LIGHT, TEAL_LIGHT][i % 3],
                    RGBColor(0xC0, 0xC0, 0xC0), 0.5)
        # Number circle
        circ = slide.shapes.add_shape(9, block_x + Inches(0.12),
                                       y + (hill_h - Inches(0.48)) / 2,
                                       Inches(0.48), Inches(0.48))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        ctf = circ.text_frame
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        cr.font.name  = "Arial"
        cr.font.size  = Pt(14)
        cr.font.bold  = True
        cr.font.color.rgb = WHITE
        # WHO / WHAT / WOW in 3 sub-columns
        sub_labels = [("WHO", who), ("WHAT", what), ("WOW", wow)]
        sub_x = block_x + Inches(0.72)
        sub_w = (block_w - Inches(0.72) - Inches(0.1)) / 3
        for j, (lbl, val) in enumerate(sub_labels):
            sx = sub_x + j * (sub_w + Inches(0.06))
            # Label
            ltb = slide.shapes.add_textbox(
                sx, y + Inches(0.06), sub_w, Inches(0.22))
            lp = ltb.text_frame.paragraphs[0]
            lp.add_run().text = lbl
            lp.runs[0].font.name  = "Arial"
            lp.runs[0].font.size  = Pt(8)
            lp.runs[0].font.bold  = True
            lp.runs[0].font.color.rgb = color
            # Value
            vtb = slide.shapes.add_textbox(
                sx, y + Inches(0.28), sub_w, hill_h - Inches(0.34))
            vtf = vtb.text_frame
            vtf.word_wrap = True
            vtf.margin_left = vtf.margin_top = 0
            vp = vtf.paragraphs[0]
            vp.add_run().text = val
            vp.runs[0].font.name  = "Calibri"
            vp.runs[0].font.size  = Pt(9)
            vp.runs[0].font.color.rgb = RGBColor(0x30, 0x30, 0x30)
        y += hill_h + Inches(0.1)


def _next_steps_slide(prs, rows):
    """Next steps table slide."""
    slide = blank_slide(prs)
    add_header(slide, prs, "Next Steps & Owners", "Immediate priorities from workshop")
    ct = content_top(has_subtitle=True)
    # Column widths
    col_defs = [
        ("#",      Inches(0.4),  PP_ALIGN.CENTER, BLUE),
        ("ACTION", Inches(5.5),  PP_ALIGN.LEFT,   BLUE),
        ("OWNER",  Inches(3.0),  PP_ALIGN.LEFT,   BLUE),
        ("TARGET", Inches(2.0),  PP_ALIGN.CENTER, BLUE),
        ("STATUS", Inches(1.9),  PP_ALIGN.CENTER, BLUE),
    ]
    header_h = Inches(0.34)
    row_h    = Inches(0.46)
    total_w  = sum(c[1] for c in col_defs)
    tx0      = (SLIDE_W - total_w) / 2  # center table

    # Header row
    x = tx0
    for col_name, col_w, align, _ in col_defs:
        _solid_rect(slide, x, ct, col_w, header_h, BLUE)
        tb = slide.shapes.add_textbox(x + Inches(0.06), ct + Inches(0.06),
                                       col_w - Inches(0.12), header_h - Inches(0.08))
        p  = tb.text_frame.paragraphs[0]
        p.alignment = align
        p.add_run().text = col_name
        p.runs[0].font.name  = "Arial"
        p.runs[0].font.size  = Pt(9)
        p.runs[0].font.bold  = True
        p.runs[0].font.color.rgb = WHITE
        x += col_w

    # Data rows
    status_colors = {
        "To Do":     RGBColor(0xD0, 0xD0, 0xD0),
        "In Progress": ORANGE,
        "Done":      GREEN,
    }
    for ri, (action, owner, target, status) in enumerate(rows):
        ry = ct + header_h + ri * row_h
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        x  = tx0
        row_vals = [str(ri + 1), action, owner, target, status]
        for ci, (_, col_w, align, _) in enumerate(col_defs):
            _solid_rect(slide, x, ry, col_w, row_h, bg,
                        RGBColor(0xD8, 0xD8, 0xD8), 0.4)
            cell_text = row_vals[ci]
            if ci == 4:  # status badge
                sc = status_colors.get(cell_text, GRAY)
                badge_w = col_w - Inches(0.2)
                badge_x = x + (col_w - badge_w) / 2
                badge_y = ry + (row_h - Inches(0.26)) / 2
                _solid_rect(slide, badge_x, badge_y, badge_w, Inches(0.26), sc)
                tb = slide.shapes.add_textbox(
                    badge_x, badge_y, badge_w, Inches(0.26))
                p  = tb.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.add_run().text = cell_text
                p.runs[0].font.name  = "Arial"
                p.runs[0].font.size  = Pt(8)
                p.runs[0].font.bold  = True
                p.runs[0].font.color.rgb = WHITE
            else:
                tb = slide.shapes.add_textbox(
                    x + Inches(0.06), ry + Inches(0.06),
                    col_w - Inches(0.12), row_h - Inches(0.1))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_top = 0
                p  = tf.paragraphs[0]
                p.alignment = align
                p.add_run().text = cell_text
                p.runs[0].font.name  = "Calibri"
                p.runs[0].font.size  = Pt(9)
                p.runs[0].font.color.rgb = RGBColor(0x30, 0x30, 0x30)
            x += col_w


def _closing_slide(prs):
    slide = blank_slide(prs)
    _solid_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BLUE)
    _solid_rect(slide, 0, 0, Inches(0.5), prs.slide_height, RED)
    tb = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.8), prs.slide_width - Inches(2.4), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.add_run().text = "Thank You"
    p.runs[0].font.name  = "Arial"
    p.runs[0].font.size  = Pt(54)
    p.runs[0].font.bold  = True
    p.runs[0].font.color.rgb = WHITE
    tb2 = slide.shapes.add_textbox(
        Inches(1.2), Inches(4.2), prs.slide_width - Inches(2.4), Inches(0.5))
    p2  = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.add_run().text = "Questions and continued conversation welcome"
    p2.runs[0].font.name  = "Calibri"
    p2.runs[0].font.size  = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


def create_playback_deck():
    prs = new_prs()

    # 1. Cover
    _cover_slide(prs)

    # 2. Agenda divider
    _section_divider_slide(prs, 1, "Workshop Objectives & Agenda", BLUE)

    # 3. Objectives
    _content_slide(prs,
        "Workshop Objectives",
        "Why we gathered and what we set out to learn",
        [
            ("", "Outline the current-state cloud adoption journey across teams"),
            ("", "Identify key roles, tools, and systems in use today"),
            ("", "Surface pain points and communication gaps in cloud workflows"),
            ("", "Explore automation and standardization opportunities"),
            ("", "Determine clear next steps and owners to address top pain points"),
        ],
        BLUE,
    )

    # 4. Personas divider
    _section_divider_slide(prs, 2, "Who Was in the Room — Key Personas", TEAL)

    # 5. Personas summary
    _two_col_slide(prs,
        "Key Personas",
        "Roles represented in the workshop",
        [
            "Cloud Platform Engineer — 'I spend most of my time fielding ad-hoc requests instead of building shared services.'",
            "Application Developer — 'I just want a fast, secure way to deploy — I don't care about the underlying infrastructure.'",
            "Security / Compliance Lead — 'Every manual review is a bottleneck I wish we could automate away.'",
        ],
        [
            "FinOps / Finance — 'We have no visibility into who is spending what until the bill arrives.'",
            "Infrastructure Lead — 'We've accumulated years of legacy config — nobody wants to touch it.'",
            "PMO — 'Projects stall waiting for cloud provisioning approvals that take weeks.'",
        ],
        TEAL, BLUE,
        "Platform & Engineering", "Business & Operations",
    )

    # 6. Insights divider
    _section_divider_slide(prs, 3, "What We Heard — Key Insights", RED)

    # 7. Insights
    _content_slide(prs,
        "Key Insights from Empathy Work",
        "Patterns that emerged across all personas",
        [
            ("Developer Experience",  "Friction in provisioning kills cloud adoption — every extra step is a reason to do it the old way."),
            ("Security as a Blocker", "Manual security reviews are the #1 cited bottleneck; teams want guardrails, not gates."),
            ("Cost Opacity",          "No team has real-time cost visibility; FinOps decisions are made on month-old data."),
            ("Role Ambiguity",        "Unclear ownership between platform and app teams creates duplicate work and gaps."),
            ("Legacy Drag",           "Accumulated technical debt in infra config prevents teams from adopting newer patterns."),
        ],
        RED,
    )

    # 8. Ideas & Prioritization divider
    _section_divider_slide(prs, 4, "Big Ideas & Prioritization", GREEN)

    # 9. Big Ideas two-col
    _two_col_slide(prs,
        "Top Big Ideas",
        "High-value concepts surfaced in ideation",
        [
            "Self-service cloud landing zones with pre-approved architecture patterns",
            "Automated compliance-as-code integrated directly into CI/CD pipelines",
            "FinOps dashboard with real-time cost allocation per product team",
            "Developer portal (e.g., Backstage) as single pane of glass for cloud resources",
        ],
        [
            "Infrastructure-as-code template library for common workload patterns",
            "Automated onboarding wizard — new teams provisioned in <1 hour",
            "Tagging enforcement automation — no untagged resources in production",
            "Cloud cost anomaly alerting sent directly to product owners",
        ],
        GREEN, TEAL,
        "Platform Capabilities", "Developer Enablement",
    )

    # 10. Hills
    _hills_slide(prs, [
        ("Cloud platform and application teams",
         "Provision production-ready cloud environments in under 1 hour with zero manual security reviews",
         "Reducing provisioning lead time by 95% while maintaining 100% compliance posture"),
        ("All engineering teams",
         "Deploy applications using self-service patterns with full cost visibility from day one",
         "Eliminating shadow IT entirely and cutting monthly cloud spend variance by 80%"),
        ("FinOps and product owners",
         "Receive real-time cost dashboards with automated anomaly alerts",
         "Making data-driven resource decisions that cut waste by 40% quarter over quarter"),
    ])

    # 11. Next Steps divider
    _section_divider_slide(prs, 5, "Next Steps & Owners", PURPLE)

    # 12. Next steps table
    _next_steps_slide(prs, [
        ("Synthesize workshop outputs and distribute playback deck",  "Facilitator",           "Week 1",   "To Do"),
        ("Schedule 1-hour follow-up alignment session with leadership", "PMO",                 "Week 2",   "To Do"),
        ("Stand up CCoE working group with initial charter",          "Cloud Platform Lead",   "Month 1",  "To Do"),
        ("Pilot self-service landing zone for 2 app teams",           "Cloud Platform Eng.",   "Month 1-2","To Do"),
        ("Implement FinOps cost dashboard v1",                        "FinOps + Platform",     "Month 2",  "To Do"),
        ("Automate compliance-as-code in CI/CD (Phase 1)",            "Security + DevOps",     "Month 2-3","To Do"),
    ])

    # 13. Thank you
    _closing_slide(prs)

    prs.save('/Users/V5X8512/Downloads/17_DTW_Playback_Deck.pptx')
    print("Generated: 17_DTW_Playback_Deck.pptx")


# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("Generating Final Artifacts 15-17...")
    print("=" * 60)
    create_feedback_grid()
    create_executive_summary()
    create_playback_deck()
    print("=" * 60)
    print("ALL 17 ARTIFACTS COMPLETED!")
