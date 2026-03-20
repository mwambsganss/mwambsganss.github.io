"""Build ai.lilly.com summary PPT for financial operations audience."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Lilly brand colors ──────────────────────────────────────────────────────
LILLY_RED    = RGBColor(0xDA, 0x29, 0x28)   # #DA2928
LILLY_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY     = RGBColor(0x55, 0x55, 0x55)
ACCENT_BLUE  = RGBColor(0x00, 0x5A, 0xA7)   # secondary brand blue

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_name="Calibri", font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True, valign=None):
    from pptx.enum.text import MSO_ANCHOR
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=14, color=LILLY_DARK, header=None, header_size=15):
    from pptx.enum.text import MSO_ANCHOR
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header
        run.font.name = "Calibri"
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = LILLY_RED
    for item in items:
        p = tf.paragraphs[0] if (first and not header) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None):
    """Red top banner with title."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_color=LILLY_RED)
    add_text(slide, title, 0.35, 0.12, 12.5, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 12.5, 0.4,
                 font_size=14, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC),
                 align=PP_ALIGN.LEFT)
    # light bottom strip
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=LILLY_RED)
    add_text(slide, "ai.lilly.com  |  Lilly AI Program Overview", 0.3, 7.2, 10, 0.3,
             font_size=9, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "CONFIDENTIAL", 10, 7.2, 3, 0.3,
             font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LILLY_DARK)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill_color=LILLY_RED)

# Large title
add_text(slide, "AI at Lilly", 0.6, 1.0, 12, 1.3,
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "What Financial Operations Needs to Know", 0.6, 2.4, 12, 0.8,
         font_size=26, bold=False, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.LEFT)
add_rect(slide, 0.6, 3.35, 5, 0.06, fill_color=LILLY_RED)
add_text(slide, "Based on ai.lilly.com  |  Crawled February 2026", 0.6, 3.55, 10, 0.4,
         font_size=13, color=MID_GRAY, italic=True, align=PP_ALIGN.LEFT)

add_text(slide, "Lilly AI Program  •  Tech@Lilly", 0.6, 5.7, 10, 0.5,
         font_size=17, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "ai.lilly.com", 0.6, 6.2, 10, 0.4,
         font_size=13, color=RGBColor(0xFF, 0xCC, 0xCC), italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is ai.lilly.com?
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "What is ai.lilly.com?",
             "Lilly's internal hub for all things Artificial Intelligence")

# Two column layout
add_rect(slide, 0.3, 1.3, 5.9, 5.6, fill_color=LIGHT_GRAY)
add_rect(slide, 6.5, 1.3, 6.5, 5.6, fill_color=LIGHT_GRAY)

add_text(slide, "The Hub", 0.5, 1.4, 5.5, 0.5,
         font_size=17, bold=True, color=LILLY_RED, align=PP_ALIGN.LEFT)
add_bullet_box(slide,
    ["Single destination for AI products, guidance, and governance",
     "Over 400 AI products tracked enterprise-wide",
     "Maintained by Tech@Lilly and updated continuously",
     "Includes tools available to ALL employees and tools available only to specific groups",
     "Resources include PDFs, training, registry forms, and an AI chatbot (Lilly Assist)"],
    0.5, 1.9, 5.5, 4.8, font_size=13.5, color=LILLY_DARK)

add_text(slide, "Site Sections", 6.7, 1.4, 6.0, 0.5,
         font_size=17, bold=True, color=LILLY_RED, align=PP_ALIGN.LEFT)
add_bullet_box(slide,
    ["Products  —  Browse and access AI tools",
     "AI Stack  —  Reusable AI building blocks for developers",
     "AI Registry  —  Submit and track AI use cases",
     "Guidance  —  Policies, dos/don'ts, decision trees",
     "Learning  —  Training resources and upskilling paths",
     "FAQ  —  Answers to common questions"],
    6.7, 1.9, 6.0, 4.8, font_size=13.5, color=LILLY_DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Enterprise-wide AI Tools (most relevant to fin ops)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "AI Tools Available to All Employees",
             "Enterprise-wide products relevant to financial operations workflows")

tools = [
    ("Copilot Chat for Work",
     "AI assistant in Word, Excel, PowerPoint, Outlook & Teams.\nBoosts productivity for all M365 users."),
    ("Claude: Secure Chat\n(Chat in a Box)",
     "Secure Claude AI chat cleared for Red / CCI data.\nAvailable to all employees & contractors (ex. China)."),
    ("ARTIE",
     "Automated Business Analytics & Insights platform.\nGenAI-powered for financial/operational analysis."),
    ("Forecasting.AI",
     "Upload time-series data, forecast trends with ML.\nNo coding required; no data stored."),
    ("Lilly Translate",
     "Translates into ~80 languages; speech-to-text.\nIdeal for global communications and reporting."),
    ("Prompt Buddy",
     "Share & discover AI prompts in Microsoft Teams.\nAccelerates team-wide AI adoption."),
]

cols = 3
rows = 2
card_w = 4.0
card_h = 2.5
gap_x  = 0.27
gap_y  = 0.35
start_x = 0.35
start_y = 1.3

for idx, (name, desc) in enumerate(tools):
    col = idx % cols
    row = idx // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(slide, x, y, card_w, card_h, fill_color=LIGHT_GRAY,
             border_color=LILLY_RED, border_pt=1.5)
    add_text(slide, name, x + 0.12, y + 0.1, card_w - 0.25, 0.55,
             font_size=13, bold=True, color=LILLY_RED, align=PP_ALIGN.LEFT)
    add_text(slide, desc, x + 0.12, y + 0.68, card_w - 0.25, card_h - 0.85,
             font_size=11.5, color=LILLY_DARK, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Specialized Tools (limited groups)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "Specialized AI Tools (Limited Groups)",
             "Role-specific tools — request access through your Tech@Lilly partner")

spec_tools = [
    ("Pharm(ai) Studio",       "US Only",                    "Build Claude 3.7 Sonnet agents, websites, quality docs. Relevant for process automation teams."),
    ("Cortex Landing Zone",    "Cortex Developers",          "No-code agentic AI builder. Deploy knowledge workers and automated agents without coding."),
    ("AMIGO",                  "Manufacturing Globally",     "GenAI-driven line monitoring with real-time insights, recommendations, and alerts."),
    ("BD360",                  "Enterprise (all)",           "Competitive intelligence dashboard — 85% accurate insights in <30 seconds."),
    ("GitHub Copilot",         "GitHub License Holders",     "AI coding assistant for developers — suggests code, automates repetitive tasks."),
    ("iQ AI Assistant",        "US Field Teams",             "ChatGPT-style assistant for fast field insights and data access."),
]

for idx, (name, access, desc) in enumerate(spec_tools):
    col = idx % 2
    row = idx // 2
    x = 0.35 + col * 6.5
    y = 1.3 + row * 1.9
    add_rect(slide, x, y, 6.1, 1.75, fill_color=LIGHT_GRAY,
             border_color=ACCENT_BLUE, border_pt=1.2)
    add_text(slide, name, x + 0.12, y + 0.08, 4.0, 0.45,
             font_size=13, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
    add_rect(slide, x + 4.2, y + 0.1, 1.7, 0.32, fill_color=ACCENT_BLUE)
    add_text(slide, access, x + 4.21, y + 0.1, 1.68, 0.32,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 0.12, y + 0.55, 5.85, 1.1,
             font_size=11.5, color=LILLY_DARK, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI Registry: How to Submit an AI Idea
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "The AI Registry: All AI Use Cases Must Be Registered",
             "Governance process effective November 7, 2025 — streamlined for low-risk use cases")

# Process arrows
steps = [
    ("1\nPrioritization", "Weekly review: ROI, strategy fit, duplication check"),
    ("2\nFunctional", "Tech@Lilly aligns use case with local strategy & confirms need"),
    ("3\nRisk Assessment", "Architect evaluates data, design, vendor; assigns Fast / Standard / High Risk path"),
    ("4\nAI Tech Review", "Experts recommend existing systems, define architecture, mitigate tech risks"),
    ("5\nLegal & Quality", "Compliance, data protection, quality standards (GMP if applicable)"),
]

box_w = 2.2
box_h = 2.8
start_x = 0.35
y_top = 1.35
arrow_w = 0.28

for i, (title, desc) in enumerate(steps):
    x = start_x + i * (box_w + arrow_w)
    color = LILLY_RED if i == 0 else (ACCENT_BLUE if i < 3 else RGBColor(0x33, 0x77, 0x44))
    add_rect(slide, x, y_top, box_w, box_h, fill_color=color)
    add_text(slide, title, x + 0.1, y_top + 0.12, box_w - 0.2, 0.85,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 0.1, y_top + 1.0, box_w - 0.2, box_h - 1.1,
             font_size=10.5, color=WHITE, align=PP_ALIGN.LEFT)
    # arrow
    if i < 4:
        ax = x + box_w + 0.03
        add_text(slide, "▶", ax, y_top + 1.0, arrow_w, 0.8,
                 font_size=18, bold=True, color=LILLY_DARK, align=PP_ALIGN.CENTER)

# Fast path note
add_rect(slide, 0.35, 4.4, 12.6, 0.9, fill_color=RGBColor(0xFF, 0xF3, 0xCC),
         border_color=RGBColor(0xCC, 0xAA, 0x00), border_pt=1)
add_text(slide, "Low-Risk Fast Path: Low-risk use cases (common data, standard context) only need reviews 1–3, saving significant time.",
         0.5, 4.48, 12.3, 0.75, font_size=12.5, color=LILLY_DARK, bold=False)

# Key notes
add_bullet_box(slide,
    ["Dev work using green data can begin after AI Tech review (step 4)",
     "Work in QA/Prod or processing Lilly data through a vendor cannot begin until ALL reviews complete",
     "Engage reviewers proactively — unresponsive teams delay approvals",
     "Vendor contracts must include standard AI language"],
    0.35, 5.45, 12.6, 1.55, font_size=12, color=LILLY_DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Responsible Use / Dos & Don'ts
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "Using AI Responsibly at Lilly",
             "Key policies and guardrails every employee must follow")

# Left: DOs
add_rect(slide, 0.3, 1.3, 5.9, 5.6, fill_color=RGBColor(0xE8, 0xF5, 0xE9),
         border_color=RGBColor(0x2E, 0x7D, 0x32), border_pt=1.5)
add_text(slide, "DO", 0.5, 1.38, 5.5, 0.55,
         font_size=20, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
add_bullet_box(slide,
    ["Verify AI outputs for accuracy before acting on them",
     "Use Lilly-approved tools (Chat in a Box, Copilot) for confidential / personal data",
     "Disclose AI use in external-facing content and customer interactions",
     "Know your data sources — understand biases and limitations",
     "Engage AI Change Champions for help with adoption",
     "Get help on AI vendor contracts before signing anything",
     "Read the Using AI Responsibly at Lilly procedure before use"],
    0.5, 1.95, 5.5, 4.8, font_size=13, color=RGBColor(0x1B, 0x5E, 0x20))

# Right: DON'Ts
add_rect(slide, 6.5, 1.3, 6.5, 5.6, fill_color=RGBColor(0xFF, 0xEB, 0xEE),
         border_color=LILLY_RED, border_pt=1.5)
add_text(slide, "DON'T", 6.7, 1.38, 6.0, 0.55,
         font_size=20, bold=True, color=LILLY_RED)
add_bullet_box(slide,
    ["Enter confidential or personal data into unapproved public AI tools",
     "Enter into paid AI license agreements for company use (must go through procurement)",
     "Put intellectual property at risk with unvetted tools",
     "Rely on AI output without independent verification",
     "Use public AI tools outside policy guidelines for business decisions",
     "Assume AI content is accurate — hallucinations are real and common",
     "Skip the AI Registry process for new AI use cases"],
    6.7, 1.95, 6.0, 4.8, font_size=13, color=LILLY_RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Financial Operations Relevance
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "What This Means for Financial Operations",
             "Practical opportunities and obligations for FinOps teams")

# 3-column value props
cols_data = [
    ("Efficiency\nOpportunities",
     LILLY_RED,
     ["Forecasting.AI automates time-series trend analysis without coding",
      "ARTIE delivers automated analytics and business insights on demand",
      "Copilot Chat for Work drafts reports, analyzes Excel data, summarizes meetings",
      "Lilly Translate removes language barriers for global financial reporting"]),
    ("Governance\nObligations",
     ACCENT_BLUE,
     ["ALL new AI use cases must be submitted to the AI Registry",
      "Do not sign vendor AI contracts independently — get help",
      "Verify AI-generated financial figures before using in reports or decisions",
      "Disclose AI use in external or stakeholder-facing financial content"]),
    ("Getting\nStarted",
     RGBColor(0x33, 0x77, 0x44),
     ["Visit ai.lilly.com to browse approved tools",
      "Connect with your AI Change Champion for guided adoption",
      "Complete AI upskilling training (ai.lilly.com/ai-upskilling)",
      "Submit AI ideas via the Registry Intake Form",
      "Ask Lilly Assist (chatbot on ai.lilly.com) for quick guidance"]),
]

col_w = 3.95
for i, (heading, color, bullets) in enumerate(cols_data):
    x = 0.3 + i * (col_w + 0.28)
    add_rect(slide, x, 1.3, col_w, 0.75, fill_color=color)
    add_text(slide, heading, x + 0.1, 1.32, col_w - 0.2, 0.7,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 2.05, col_w, 5.0, fill_color=LIGHT_GRAY,
             border_color=color, border_pt=1)
    add_bullet_box(slide, bullets, x + 0.15, 2.15, col_w - 0.3, 4.8,
                   font_size=12.5, color=LILLY_DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Stats & Quick Reference
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "Key Facts & Quick Reference",
             "At a glance — numbers and resources from ai.lilly.com")

# Big stat boxes
stats = [
    ("400+", "AI Products\nat Lilly"),
    ("80",   "Languages\nvia Lilly Translate"),
    ("5",    "Registry Review\nSteps (standard)"),
    ("3",    "Registry Review\nSteps (fast path)"),
]
for i, (num, label) in enumerate(stats):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 1.3, 2.9, 1.6, fill_color=LILLY_RED)
    add_text(slide, num, x + 0.1, 1.35, 2.7, 0.95,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 2.25, 2.7, 0.6,
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Data classification
add_rect(slide, 0.3, 3.1, 12.7, 1.15, fill_color=RGBColor(0xFD, 0xF6, 0xE3),
         border_color=RGBColor(0xCC, 0xAA, 0x00), border_pt=1)
add_text(slide, "Data Classification Reminders",
         0.5, 3.15, 12, 0.4, font_size=14, bold=True, color=RGBColor(0x7A, 0x5C, 0x00))
add_text(slide,
         "Chat in a Box (Claude) and Copilot Chat for Work  are BOTH approved for Red / CI / PI data.  "
         "Enterprise Claude is approved for CI & PI but NOT for Orange+ SPI.  "
         "When in doubt, use Chat in a Box or ask the Digital Legal Office.",
         0.5, 3.55, 12.5, 0.65, font_size=12, color=LILLY_DARK)

# Resource list
add_text(slide, "Key Resources", 0.3, 4.35, 12, 0.4,
         font_size=15, bold=True, color=LILLY_DARK)

resources = [
    ("ai.lilly.com",                  "Main AI hub — browse products, guidance, and registry"),
    ("ai.lilly.com/products/enterprise", "Full list of enterprise-wide AI tools"),
    ("ai.lilly.com/lilly-ai-guidance",   "Policies, decision trees, dos & don'ts"),
    ("ai.lilly.com/ai-registry/review-process", "AI Registry review steps and timelines"),
    ("ai.lilly.com/ai-upskilling",       "Upskilling program for Tech@Lilly leaders"),
    ("ai@lilly.com",                     "Contact the AI team with questions"),
]
for i, (url, desc) in enumerate(resources):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 4.8 + row * 0.72
    add_text(slide, url, x, y, 2.9, 0.45,
             font_size=11, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
    add_text(slide, desc, x + 3.0, y, 3.3, 0.45,
             font_size=11, color=LILLY_DARK, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Disclosure & Compliance Obligations
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
slide_header(slide, "AI Disclosure & Compliance Obligations",
             "What financial operations staff must do to stay compliant")

# Decision tree style flow
add_rect(slide, 0.3, 1.3, 12.7, 0.65, fill_color=LILLY_DARK)
add_text(slide, "WHEN TO DISCLOSE AI USE", 0.5, 1.35, 12.5, 0.55,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

scenarios = [
    ("External-facing content\n(marketing, publications, social media)",
     "DISCLOSURE REQUIRED — label as AI-generated"),
    ("Internal reports / decisions\n(e.g., financial analysis, forecasts)",
     "DISCLOSE if AI output directly drives decisions"),
    ("Drafting emails, translating docs\n(Copilot, LillyTranslate)",
     "Exempt — treated as productivity tools, no disclosure needed"),
    ("Chatbots / AI tools deployed\nto other Lilly staff or customers",
     "REQUIRED — inform users they are interacting with AI"),
]

for i, (scenario, action) in enumerate(scenarios):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 2.1 + row * 2.1
    add_rect(slide, x, y, 2.9, 1.8, fill_color=LIGHT_GRAY,
             border_color=MID_GRAY, border_pt=0.8)
    add_text(slide, scenario, x + 0.12, y + 0.1, 2.65, 0.9,
             font_size=11.5, color=LILLY_DARK, bold=False)
    add_rect(slide, x + 2.92, y, 3.35, 1.8, fill_color=RGBColor(0xFF, 0xEB, 0xEE)
             if "REQUIRED" in action else RGBColor(0xE8, 0xF5, 0xE9),
             border_color=LILLY_RED if "REQUIRED" in action else RGBColor(0x2E, 0x7D, 0x32),
             border_pt=1)
    add_text(slide, action, x + 3.05, y + 0.4, 3.1, 1.0,
             font_size=12, color=LILLY_RED if "REQUIRED" in action else RGBColor(0x1B, 0x5E, 0x20),
             bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 6.35, 12.7, 0.6, fill_color=RGBColor(0xFF, 0xF3, 0xCC),
         border_color=RGBColor(0xCC, 0xAA, 0x00), border_pt=1)
add_text(slide, "Rule of thumb: When in doubt, disclose.  Questions? Contact the Digital Legal Office or ai@lilly.com",
         0.5, 6.38, 12.5, 0.55, font_size=12.5, color=LILLY_DARK, bold=False,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Next Steps / Call to Action
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LILLY_DARK)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill_color=LILLY_RED)

add_text(slide, "Your Next Steps", 0.6, 0.5, 12, 1.0,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_rect(slide, 0.6, 1.5, 4, 0.06, fill_color=LILLY_RED)

actions = [
    ("1", "Visit ai.lilly.com", "Explore the product catalog and find tools relevant to your work"),
    ("2", "Complete AI Upskilling", "Take the 2025 AI Upskilling course at ai.lilly.com/ai-upskilling"),
    ("3", "Identify a Use Case", "Spot a process that could benefit from AI and register it via the Intake Form"),
    ("4", "Connect with AI Change Champions", "Get personalized help with tool adoption and strategy"),
    ("5", "Review AI Guidance", "Read the Using AI Responsibly procedure before deploying any AI solution"),
]

for i, (num, title, desc) in enumerate(actions):
    x = 0.5 + (i % 3) * 4.1
    y = 1.7 + (i // 3) * 1.85
    add_rect(slide, x, y, 0.55, 0.55, fill_color=LILLY_RED)
    add_text(slide, num, x, y, 0.55, 0.55,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.65, y + 0.02, 3.2, 0.4,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, desc, x + 0.65, y + 0.42, 3.2, 0.8,
             font_size=11.5, color=MID_GRAY, align=PP_ALIGN.LEFT)

add_text(slide, "Questions?  ai@lilly.com  |  ai.lilly.com  |  Lilly Assist chatbot",
         0.6, 5.65, 12, 0.5, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Tech@Lilly AI Program  •  Based on ai.lilly.com content crawled February 2026",
         0.6, 6.2, 12, 0.4, font_size=12, color=RGBColor(0xFF, 0xCC, 0xCC),
         italic=True, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/V5X8512/Library/CloudStorage/OneDrive-EliLillyandCompany/Documents/Repos/MWambsganss/mwambsganss.github.io/Web Crawler/ai.lilly.com_crawl/AI_at_Lilly_FinOps_Summary.pptx"
prs.save(out)
print(f"Saved: {out}")
